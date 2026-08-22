"""PPTX Design engine — pure python-pptx logic, zero MCP imports."""

import copy
import subprocess
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # type: ignore[attr-defined]
from pptx.util import Inches, Pt

from shared.file_utils import embed_content, hint_for_error, resolve_path
from shared.live_edit import notify_reload
from shared.platform_utils import get_pdf_converter, open_file, resolve_output_path
from shared.progress import fail, index_range, ok, warn
from shared.slide_visual import (
    NoRoomOnSlide,
    contrast_ratio,
    contrast_warning,
    place_below_content,
    slide_background_hex,
    unreadable_slides,
)
from shared.version_control import snapshot

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

CHART_TYPE_MAP: dict[str, Any] = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}


def _axis_number_format(values: list[Any]) -> str:
    """Pick a tick format that fits, the way the fleet's other charts already do.

    Left unformatted, a spend axis printed "500000 / 1000000 / 2500000". The
    labels were too wide for the plot, so they came out rotated 45 degrees, and
    on a chart squeezed under a bullet list half of them were dropped
    altogether. Each trailing comma in an Excel format divides by a thousand, so
    the same axis reads "0.5M / 1M / 2.5M" and stays level. The K/M thresholds
    match the ones the Data_Analyst charts use, so a figure looks the same in a
    deck as on the dashboard it was copied from.
    """
    peak = 0.0
    for v in values:
        try:
            peak = max(peak, abs(float(v)))
        except (TypeError, ValueError):
            continue
    if peak >= 1_000_000:
        return '#,##0.#,,"M"'
    if peak >= 10_000:
        return '#,##0.#,"K"'
    return "#,##0.##"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _open_prs(path: Path, progress: list[dict[str, Any]]) -> tuple[Any, dict[str, Any] | None]:
    """Load presentation; return (prs, None) or (None, error_dict)."""
    if not path.exists():
        progress.append(fail("File not found", str(path)))
        return None, {
            "success": False,
            "error": f"File not found: {path}",
            "hint": "Check that file_path is absolute and the file exists.",
            "progress": progress,
            "token_estimate": 20,
        }
    if path.suffix.lower() != ".pptx":
        progress.append(fail(f"Wrong file type: {path.suffix}"))
        return None, {
            "success": False,
            "error": f"Expected .pptx file, got {path.suffix}",
            "hint": "Use the correct server for this file type.",
            "progress": progress,
            "token_estimate": 20,
        }
    prs = Presentation(str(path))
    return prs, None


def _check_slide(
    prs: Any, slide_index: int, progress: list[dict[str, Any]], backup: str | None
) -> tuple[Any, dict[str, Any] | None]:
    """Return (slide, None) or (None, error_dict) if index out of range."""
    count = len(prs.slides)
    if slide_index < 0 or slide_index >= count:
        progress.append(fail(f"Slide index {slide_index} out of range", f"Presentation has {count} slide(s)"))
        return None, {
            "success": False,
            "error": f"slide_index {slide_index} out of range {index_range(count, 'slides')}",
            "hint": "Use read_presentation to get current slide count.",
            "backup": backup,
            "progress": progress,
            "token_estimate": 15,
        }
    return prs.slides[slide_index], None


def _find_shape(
    slide: Any, shape_name: str, progress: list[dict[str, Any]], backup: str | None
) -> tuple[Any, dict[str, Any] | None]:
    """Return (shape, None) or (None, error_dict) if shape not found."""
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape, None
    names = [s.name for s in slide.shapes]
    progress.append(fail(f"Shape '{shape_name}' not found"))
    return None, {
        "success": False,
        "error": f"Shape '{shape_name}' not found on slide",
        "hint": f"Available shapes: {', '.join(names)}",
        "backup": backup,
        "progress": progress,
        "token_estimate": 15,
    }


# ---------------------------------------------------------------------------
# Tool implementations
# ---------------------------------------------------------------------------


def _match_chart_text_to_slide(chart: Any, slide: Any) -> str:
    """Colour every piece of chart text to contrast with the slide behind it.

    Returns the hex applied, or "" when the slide background is unknown and the
    chart is left at python-pptx's defaults.
    """
    background = slide_background_hex(slide)
    if background is None:
        return ""

    readable = "FFFFFF" if contrast_ratio("FFFFFF", background) >= contrast_ratio("000000", background) else "000000"
    rgb = RGBColor.from_string(readable)

    # chart.font cascades to anything that has not been styled individually;
    # the axes are set explicitly because their tick labels usually have been.
    try:
        chart.font.color.rgb = rgb
    except (AttributeError, TypeError, ValueError):
        pass

    for axis_name in ("category_axis", "value_axis"):
        try:
            getattr(chart, axis_name).tick_labels.font.color.rgb = rgb
        except (AttributeError, TypeError, ValueError):
            # Pie charts have no axes; nothing to colour.
            continue

    try:
        if chart.has_title:
            for para in chart.chart_title.text_frame.paragraphs:
                para.font.color.rgb = rgb
                for run in para.runs:
                    run.font.color.rgb = rgb
    except (AttributeError, TypeError, ValueError):
        pass

    return readable


def set_background(
    file_path: str,
    slide_index: int,
    color_hex: str = "",
    image_path: str = "",
    open_after: bool = False,
) -> dict[str, Any]:
    """Set slide background to a solid color or image."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    try:
        if not color_hex and not image_path:
            progress.append(fail("Must provide color_hex or image_path"))
            return {
                "success": False,
                "error": "Must provide either color_hex or image_path",
                "hint": "Provide a 6-char hex color like 'FF0000' or an image file path.",
                "progress": progress,
                "token_estimate": 15,
            }

        path = resolve_path(file_path)
        prs, err = _open_prs(path, progress)
        if err:
            return err

        progress.append(ok(f"Opened {path.name}", f"{len(prs.slides)} slides"))

        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        # slide_index=-1 means every slide. Without it, "make this deck dark"
        # took one call per slide, while set_font_all_slides coloured all of
        # them in one -- which is exactly how a deck ended up with a dark slide
        # 1, an untouched white slide 2 and white text on both.
        if slide_index == -1:
            targets = list(prs.slides)
            scope = f"all {len(targets)} slides"
        else:
            slide, err = _check_slide(prs, slide_index, progress, backup)
            if err:
                return err
            targets = [slide]
            scope = f"slide {slide_index}"

        if color_hex:
            clean = color_hex.lstrip("#")
            for target in targets:
                fill = target.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(
                    int(clean[0:2], 16),
                    int(clean[2:4], 16),
                    int(clean[4:6], 16),
                )
            progress.append(ok(f"Set background color #{clean}", scope))

        elif image_path:
            img_path = Path(image_path).resolve()
            if not img_path.exists():
                progress.append(fail(f"Image not found: {img_path.name}"))
                return {
                    "success": False,
                    "error": f"Image file not found: {image_path}",
                    "hint": "Check that image_path is absolute and the file exists.",
                    "backup": backup,
                    "progress": progress,
                    "token_estimate": 15,
                }
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            for target in targets:
                target.shapes.add_picture(str(img_path), 0, 0, slide_width, slide_height)
                # Move image to back (index 0 in spTree after two required elements)
                sp_tree = target.shapes._spTree
                pic_el = sp_tree[-1]
                sp_tree.remove(pic_el)
                sp_tree.insert(2, pic_el)
            progress.append(ok("Set background image", f"{img_path.name} on {scope}"))

        prs.save(str(path))
        if open_after:
            open_file(path)
        progress.append(notify_reload(str(path), "pptx"))

        result: dict[str, Any] = {
            "success": True,
            "op": "set_background",
            "slide_index": slide_index,
            "color_hex": color_hex,
            "image_path": image_path,
            "backup": backup,
            "progress": progress,
        }
        result["token_estimate"] = len(str(result)) // 4
        return result

    except Exception as e:
        progress.append(fail(str(e)))
        return {
            "success": False,
            "error": str(e),
            "hint": hint_for_error(e, path),
            "backup": backup,
            "progress": progress,
            "token_estimate": 15,
        }


def set_font_style(
    file_path: str,
    slide_index: int,
    shape_name: str,
    font_name: str = "",
    font_size: float = 0,
    bold: bool = False,
    color_hex: str = "",
    open_after: bool = False,
) -> dict[str, Any]:
    """Apply font name, size, bold, and/or color to all runs in a shape."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    try:
        path = resolve_path(file_path)
        prs, err = _open_prs(path, progress)
        if err:
            return err

        progress.append(ok(f"Opened {path.name}"))

        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        slide, err = _check_slide(prs, slide_index, progress, backup)
        if err:
            return err

        shape, err = _find_shape(slide, shape_name, progress, backup)
        if err:
            return err

        if not shape.has_text_frame:
            progress.append(fail(f"Shape '{shape_name}' has no text frame"))
            return {
                "success": False,
                "error": f"Shape '{shape_name}' has no text frame",
                "hint": "Only text shapes support font styling.",
                "backup": backup,
                "progress": progress,
                "token_estimate": 15,
            }

        rgb: RGBColor | None = None
        if color_hex:
            clean = color_hex.lstrip("#")
            rgb = RGBColor(
                int(clean[0:2], 16),
                int(clean[2:4], 16),
                int(clean[4:6], 16),
            )

        runs_updated = 0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if font_name:
                    run.font.name = font_name
                if font_size > 0:
                    run.font.size = Pt(font_size)
                if bold:
                    run.font.bold = True
                if rgb is not None:
                    run.font.color.rgb = rgb
                runs_updated += 1

        prs.save(str(path))
        if open_after:
            open_file(path)
        progress.append(notify_reload(str(path), "pptx"))
        progress.append(ok(f"Updated font style in '{shape_name}'", f"{runs_updated} runs"))

        result: dict[str, Any] = {
            "success": True,
            "op": "set_font_style",
            "slide_index": slide_index,
            "shape_name": shape_name,
            "font_name": font_name,
            "font_size": font_size,
            "bold": bold,
            "color_hex": color_hex,
            "runs_updated": runs_updated,
            "backup": backup,
            "progress": progress,
        }
        result["token_estimate"] = len(str(result)) // 4
        return result

    except Exception as e:
        progress.append(fail(str(e)))
        return {
            "success": False,
            "error": str(e),
            "hint": hint_for_error(e, path),
            "backup": backup,
            "progress": progress,
            "token_estimate": 15,
        }


def add_table(
    file_path: str,
    slide_index: int,
    rows: int,
    cols: int,
    data: list[list[str]],
    left: float = 1.0,
    top: float = 2.0,
    width: float = 8.0,
    height: float = 3.0,
    open_after: bool = False,
) -> dict[str, Any]:
    """Insert a table with data on a slide."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    try:
        if rows <= 0 or cols <= 0:
            progress.append(fail("rows and cols must be positive integers"))
            return {
                "success": False,
                "error": "rows and cols must be positive",
                "hint": "Provide positive integers for rows and cols.",
                "progress": progress,
                "token_estimate": 15,
            }

        path = resolve_path(file_path)
        prs, err = _open_prs(path, progress)
        if err:
            return err

        progress.append(ok(f"Opened {path.name}"))

        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        slide, err = _check_slide(prs, slide_index, progress, backup)
        if err:
            return err

        # Default top is 2.0in, which lands squarely on a layout's body text.
        # A swept deck ended up with bullets, a table and a chart stacked in the
        # same region -- three shapes, none readable, all reported as success.
        try:
            left, top, width, height, place_note = place_below_content(prs, slide, left, top, width, height)
        except NoRoomOnSlide as exc:
            progress.append(fail("No room for the table", str(exc)))
            return {
                "success": False,
                "error": str(exc),
                "hint": "Add a new slide with add_slide and put the table there, "
                "or shorten the text already on this one.",
                "backup": backup,
                "progress": progress,
                "token_estimate": 30,
            }
        if place_note:
            progress.append(warn("Table repositioned", place_note))

        table_shape = slide.shapes.add_table(
            rows,
            cols,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        table = table_shape.table

        for ri, row_data in enumerate(data):
            if ri >= rows:
                break
            for ci, cell_text in enumerate(row_data):
                if ci >= cols:
                    break
                table.cell(ri, ci).text = str(cell_text)

        prs.save(str(path))
        if open_after:
            open_file(path)
        progress.append(notify_reload(str(path), "pptx"))
        progress.append(ok(f"Added {rows}×{cols} table", f"slide {slide_index}"))

        result: dict[str, Any] = {
            "success": True,
            "op": "add_table",
            "slide_index": slide_index,
            "rows": rows,
            "cols": cols,
            "backup": backup,
            "progress": progress,
        }
        result["token_estimate"] = len(str(result)) // 4
        return result

    except Exception as e:
        progress.append(fail(str(e)))
        return {
            "success": False,
            "error": str(e),
            "hint": hint_for_error(e, path),
            "backup": backup,
            "progress": progress,
            "token_estimate": 15,
        }


def add_chart(
    file_path: str,
    slide_index: int,
    chart_type: str,
    data: dict[str, Any],
    title: str = "",
    left: float = 1.0,
    top: float = 2.0,
    width: float = 6.0,
    height: float = 4.5,
    open_after: bool = False,
) -> dict[str, Any]:
    """Add a bar, line, or pie chart to a slide."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    try:
        if chart_type not in CHART_TYPE_MAP:
            progress.append(fail(f"Unsupported chart type: {chart_type}"))
            return {
                "success": False,
                "error": f"Unsupported chart type: {chart_type}",
                "hint": f"Allowed types: {', '.join(sorted(CHART_TYPE_MAP))}",
                "progress": progress,
                "token_estimate": 15,
            }

        path = resolve_path(file_path)
        prs, err = _open_prs(path, progress)
        if err:
            return err

        progress.append(ok(f"Opened {path.name}"))

        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        slide, err = _check_slide(prs, slide_index, progress, backup)
        if err:
            return err

        chart_data = ChartData()
        categories = data.get("categories", [])
        chart_data.categories = categories

        series = data.get("series", [])
        # Accept both {"series": {name: [values], ...}} and the documented
        # {"series": [{"name": ..., "values": [...]}, ...]} shape — the tool's
        # own "data: {categories, series}" description doesn't pin down which,
        # and the dict form crashed with a raw "string indices" TypeError.
        if isinstance(series, dict):
            series = [{"name": name, "values": values} for name, values in series.items()]
        if not isinstance(series, list) or not all(
            isinstance(s, dict) and "name" in s and "values" in s for s in series
        ):
            return {
                "success": False,
                "error": "data['series'] must be a list of {'name': str, 'values': list} objects, "
                "or a dict of {name: [values]}.",
                "hint": 'Example: {"categories": ["A","B"], "series": {"Revenue": [10, 20]}}',
                "progress": progress,
                "token_estimate": 40,
            }
        for series_def in series:
            chart_data.add_series(series_def["name"], series_def["values"])

        # python-pptx places a shape wherever it is told: past the edge of the
        # slide, or on top of whatever is already there. A chart drawn that way
        # loses its lower category labels, or lands on the body text and buries
        # it. Both happened in swept decks, both reported success.
        try:
            left, top, width, height, fit_note = place_below_content(prs, slide, left, top, width, height)
        except NoRoomOnSlide as exc:
            progress.append(fail("No room for the chart", str(exc)))
            return {
                "success": False,
                "error": str(exc),
                "hint": "Add a new slide with add_slide and put the chart there, "
                "or shorten the text already on this one.",
                "backup": backup,
                "progress": progress,
                "token_estimate": 30,
            }
        if fit_note:
            progress.append(warn("Chart repositioned", fit_note))

        xl_chart_type = CHART_TYPE_MAP[chart_type]
        chart_shape = slide.shapes.add_chart(
            xl_chart_type,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            chart_data,
        )

        if title:
            chart_shape.chart.has_title = True
            chart_shape.chart.chart_title.text_frame.text = title

        # A pie has no value axis, and a single-series chart names itself in its
        # title -- a legend there is one more thing to read for no information.
        if chart_type == "pie" or len(series) > 1:
            chart_shape.chart.has_legend = True
            chart_shape.chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart_shape.chart.legend.include_in_layout = False

        if chart_type != "pie":
            try:
                ticks = chart_shape.chart.value_axis.tick_labels
                ticks.number_format = _axis_number_format([v for s in series for v in s["values"]])
                ticks.number_format_is_linked = False
            except (AttributeError, ValueError, NotImplementedError):
                # Not every chart type python-pptx builds exposes a value axis.
                pass

        # Every piece of text python-pptx puts in a chart defaults to near-black:
        # the title, the category names, the value ticks and the legend. Dropped
        # on a dark slide they are all unreadable. Colouring only the title, as
        # the first pass here did, left "Instagram / LinkedIn / Google" and the
        # 0-10 axis black on navy -- visible in the render, and still wrong.
        _match_chart_text_to_slide(chart_shape.chart, slide)

        prs.save(str(path))
        if open_after:
            open_file(path)
        progress.append(notify_reload(str(path), "pptx"))
        progress.append(ok(f"Added {chart_type} chart", f"slide {slide_index}"))

        result: dict[str, Any] = {
            "success": True,
            "op": "add_chart",
            "slide_index": slide_index,
            "chart_type": chart_type,
            "title": title,
            "backup": backup,
            "progress": progress,
        }
        result["token_estimate"] = len(str(result)) // 4
        return result

    except Exception as e:
        progress.append(fail(str(e)))
        return {
            "success": False,
            "error": str(e),
            "hint": hint_for_error(e, path),
            "backup": backup,
            "progress": progress,
            "token_estimate": 15,
        }


def duplicate_slide(
    file_path: str,
    slide_index: int,
    insert_at: int = -1,
    open_after: bool = False,
) -> dict[str, Any]:
    """Copy a slide and insert it at the specified position (-1 = end)."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    try:
        path = resolve_path(file_path)
        prs, err = _open_prs(path, progress)
        if err:
            return err

        progress.append(ok(f"Opened {path.name}", f"{len(prs.slides)} slides"))

        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        slide, err = _check_slide(prs, slide_index, progress, backup)
        if err:
            return err

        # Add a blank slide using the same layout as the source
        src_layout = slide.slide_layout
        new_slide = prs.slides.add_slide(src_layout)

        # add_slide() clones the layout's placeholders into the new slide, and
        # the loop below then copies the source's own placeholders across with
        # their text -- so every duplicate ended up holding both. A copied slide
        # carried two shapes named "Title 1", one populated and one empty, and
        # read_slide reported four shapes where the original had two, with
        # nothing to say which one an edit should address. The copies are the
        # ones that carry the content, so the layout's go.
        for cloned in list(new_slide.shapes):
            cloned.element.getparent().remove(cloned.element)

        # Copy all shapes from the source slide, keeping their original order.
        # Inserting each at a fixed index put every new shape ahead of the last,
        # so a copied slide listed its shapes back to front.
        for shape in slide.shapes:
            new_slide.shapes._spTree.append(copy.deepcopy(shape.element))

        new_idx = len(prs.slides) - 1

        # Reorder if insert_at is specified and not -1
        if insert_at >= 0 and insert_at < new_idx:
            # Move the new slide (currently at new_idx) to insert_at
            xml_slides = prs.slides._sldIdLst
            slides_list = list(xml_slides)
            moved = slides_list.pop(new_idx)
            slides_list.insert(insert_at, moved)
            for item in reversed(slides_list):
                xml_slides.remove(item)
                xml_slides.insert(0, item)
            new_idx = insert_at

        prs.save(str(path))
        if open_after:
            open_file(path)
        progress.append(notify_reload(str(path), "pptx"))
        progress.append(
            ok(
                f"Duplicated slide {slide_index} → position {new_idx}",
                f"{len(prs.slides)} slides total",
            )
        )

        result: dict[str, Any] = {
            "success": True,
            "op": "duplicate_slide",
            "source_index": slide_index,
            "new_index": new_idx,
            "slide_count": len(prs.slides),
            "backup": backup,
            "progress": progress,
        }
        result["token_estimate"] = len(str(result)) // 4
        return result

    except Exception as e:
        progress.append(fail(str(e)))
        return {
            "success": False,
            "error": str(e),
            "hint": hint_for_error(e, path),
            "backup": backup,
            "progress": progress,
            "token_estimate": 15,
        }


def export_pdf(
    file_path: str,
    output_path: str = "",
    open_after: bool = True,
    return_content: bool = False,
) -> dict[str, Any]:
    """Export PPTX to PDF using LibreOffice or Microsoft PowerPoint."""
    progress: list[dict[str, Any]] = []
    try:
        path = resolve_path(file_path)
        if not path.exists():
            progress.append(fail("File not found", str(path)))
            return {
                "success": False,
                "error": f"File not found: {path}",
                "hint": "Check that file_path is absolute and the file exists.",
                "progress": progress,
                "token_estimate": 15,
            }
        if path.suffix.lower() != ".pptx":
            progress.append(fail(f"Wrong file type: {path.suffix}"))
            return {
                "success": False,
                "error": f"Expected .pptx file, got {path.suffix}",
                "hint": "This tool only exports .pptx files.",
                "progress": progress,
                "token_estimate": 15,
            }

        converter = get_pdf_converter()
        if not converter:
            progress.append(fail("No PDF converter available"))
            return {
                "success": False,
                "error": "No PDF converter found on this system",
                "hint": (
                    "Install LibreOffice: "
                    "sudo apt install libreoffice (Ubuntu) or "
                    "sudo dnf install libreoffice (Fedora). "
                    "On Windows/macOS, install Microsoft PowerPoint."
                ),
                "progress": progress,
                "token_estimate": 20,
            }

        out = resolve_output_path(output_path or path.stem + ".pdf", path.stem + ".pdf")

        if converter == "libreoffice":
            result_proc = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(out.parent),
                    str(path),
                ],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result_proc.returncode != 0:
                progress.append(fail("LibreOffice conversion failed", result_proc.stderr[:200]))
                return {
                    "success": False,
                    "error": "LibreOffice PDF conversion failed",
                    "hint": result_proc.stderr[:200] if result_proc.stderr else "Check LibreOffice installation.",
                    "progress": progress,
                    "token_estimate": 15,
                }
            # LibreOffice names the output after the input file
            expected_out = out.parent / path.with_suffix(".pdf").name
            if output_path and expected_out != out:
                expected_out.rename(out)

        elif converter == "word":
            try:
                import docx2pdf  # type: ignore[import-untyped]

                docx2pdf.convert(str(path), str(out))
            except Exception as conv_err:
                progress.append(fail("PowerPoint conversion failed", str(conv_err)[:200]))
                return {
                    "success": False,
                    "error": f"PowerPoint PDF conversion failed: {conv_err}",
                    "hint": "Ensure Microsoft PowerPoint is installed.",
                    "progress": progress,
                    "token_estimate": 15,
                }

        progress.append(ok("Exported to PDF", out.name))

        if open_after:
            open_file(out)
            progress.append(ok("Opened PDF in default viewer"))

        result: dict[str, Any] = {
            "success": True,
            "op": "export_pdf",
            "input": str(path),
            "output": str(out),
            "output_name": out.name,
            "converter": converter,
            "progress": progress,
        }
        embed_content(result, out, return_content)
        result["token_estimate"] = len(str(result)) // 4
        return result

    except Exception as e:
        progress.append(fail(str(e)))
        return {
            "success": False,
            "error": str(e),
            "hint": "Check file path and that a PDF converter is installed.",
            "progress": progress,
            "token_estimate": 15,
        }


def add_image_to_all_slides(
    file_path: str,
    image_path: str,
    left: float = 0.1,
    top: float = 0.1,
    width: float = 1.0,
    height: float = 0.5,
    open_after: bool = False,
) -> dict[str, Any]:
    """Add the same image to every slide at a fixed position."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    try:
        path = resolve_path(file_path)
        prs, err = _open_prs(path, progress)
        if err:
            return err

        # Validate image path and format
        img_path = Path(image_path).resolve()
        supported = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}
        if not img_path.exists():
            progress.append(fail(f"Image not found: {img_path.name}"))
            return {
                "success": False,
                "error": f"Image file not found: {image_path}",
                "hint": "Check that image_path is absolute and the file exists.",
                "progress": progress,
                "token_estimate": 15,
            }
        if img_path.suffix.lower() not in supported:
            progress.append(fail(f"Unsupported image format: {img_path.suffix}"))
            return {
                "success": False,
                "error": f"Unsupported image format: {img_path.suffix}",
                "hint": f"Supported formats: {', '.join(sorted(supported))}",
                "progress": progress,
                "token_estimate": 15,
            }

        progress.append(ok(f"Opened {path.name}", f"{len(prs.slides)} slides"))

        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        slide_count = len(prs.slides)
        for slide in prs.slides:
            slide.shapes.add_picture(
                str(img_path),
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height),
            )

        prs.save(str(path))
        if open_after:
            open_file(path)
        progress.append(notify_reload(str(path), "pptx"))
        progress.append(ok(f"Added image to {slide_count} slides", img_path.name))

        result: dict[str, Any] = {
            "success": True,
            "op": "add_image_to_all_slides",
            "image": img_path.name,
            "slide_count": slide_count,
            "backup": backup,
            "progress": progress,
        }
        result["token_estimate"] = len(str(result)) // 4
        return result

    except Exception as e:
        progress.append(fail(str(e)))
        return {
            "success": False,
            "error": str(e),
            "hint": hint_for_error(e, path),
            "backup": backup,
            "progress": progress,
            "token_estimate": 15,
        }


def set_font_all_slides(
    file_path: str,
    font_name: str = "",
    font_size: float = 0,
    bold: bool = False,
    color_hex: str = "",
    open_after: bool = False,
) -> dict[str, Any]:
    """Apply font settings to every text run across all slides."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    try:
        path = resolve_path(file_path)
        prs, err = _open_prs(path, progress)
        if err:
            return err

        progress.append(ok(f"Opened {path.name}", f"{len(prs.slides)} slides"))

        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        rgb: RGBColor | None = None
        if color_hex:
            clean = color_hex.lstrip("#")
            rgb = RGBColor.from_string(clean)

        slides_modified = 0
        shapes_modified = 0
        for slide in prs.slides:
            slide_touched = False
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if font_name:
                            run.font.name = font_name
                        if font_size > 0:
                            run.font.size = Pt(font_size)
                        if bold:
                            run.font.bold = True
                        if rgb is not None:
                            run.font.color.rgb = rgb
                shapes_modified += 1
                slide_touched = True
            if slide_touched:
                slides_modified += 1

        prs.save(str(path))
        if open_after:
            open_file(path)
        progress.append(notify_reload(str(path), "pptx"))
        progress.append(
            ok(
                f"Updated font on {slides_modified} slides",
                f"{shapes_modified} shapes modified",
            )
        )

        # set_background works one slide at a time, this works on all of them.
        # Set a dark background on slide 1 and white text everywhere and slide 2
        # gets white-on-white: applied successfully, and completely unreadable.
        unreadable = unreadable_slides(prs, color_hex) if color_hex else []
        if unreadable:
            progress.append(warn("Text may be unreadable", contrast_warning(unreadable, color_hex)))

        result: dict[str, Any] = {
            "success": True,
            "op": "set_font_all_slides",
            "slides_modified": slides_modified,
            "shapes_modified": shapes_modified,
            "font_name": font_name,
            "font_size": font_size,
            "bold": bold,
            "color_hex": color_hex,
            "backup": backup,
            "progress": progress,
        }
        if unreadable:
            result["unreadable_slides"] = unreadable
            result["hint"] = contrast_warning(unreadable, color_hex)
        result["token_estimate"] = len(str(result)) // 4
        return result

    except Exception as e:
        progress.append(fail(str(e)))
        return {
            "success": False,
            "error": str(e),
            "hint": hint_for_error(e, path),
            "backup": backup,
            "progress": progress,
            "token_estimate": 15,
        }
