"""PPTX New engine — create PowerPoint presentations from scratch."""
from __future__ import annotations

import logging
import shutil
import sys
from pathlib import Path
from typing import Any

# Allow imports from the repo root (shared/)
_ROOT = Path(__file__).resolve().parents[2]
if str(_ROOT) not in sys.path:
    sys.path.insert(0, str(_ROOT))

from pptx import Presentation
from pptx.util import Inches, Pt  # noqa: F401 — available for callers/future use

from shared.platform_utils import open_file
from shared.progress import fail, info, ok, warn

logging.basicConfig(stream=sys.stderr, level=logging.WARNING)
logger = logging.getLogger(__name__)


# ─── Internal helpers ────────────────────────────────────────────────────────


def _ensure_parent(path: Path) -> None:
    """Create parent directory if it does not exist."""
    path.parent.mkdir(parents=True, exist_ok=True)


def _token_estimate(obj: Any) -> int:
    return len(str(obj)) // 4


def _error(
    msg: str,
    hint: str,
    progress: list[dict[str, Any]],
) -> dict[str, Any]:
    progress.append(fail(msg))
    return {
        "success": False,
        "error": msg,
        "hint": hint,
        "progress": progress,
        "token_estimate": _token_estimate(progress),
    }


def _set_placeholder_text(slide: Any, index: int, text: str) -> None:
    """
    Set placeholder text safely by index.

    Uses direct .text assignment — acceptable when creating new content
    because there is no existing run-level formatting to preserve.
    """
    try:
        slide.placeholders[index].text = text
    except (IndexError, KeyError, AttributeError):
        pass


# ─── Public engine functions ─────────────────────────────────────────────────


def create_presentation(
    output_path: str,
    title: str = "",
    subtitle: str = "",
    open_after: bool = True,
) -> dict[str, Any]:
    """Create a blank presentation with a single title slide."""
    progress: list[dict[str, Any]] = []
    path = Path(output_path)

    try:
        _ensure_parent(path)
        progress.append(info("Creating new presentation", path.name))

        prs = Presentation()
        # Layout index 0 is the Title Slide layout in default templates
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        if title:
            _set_placeholder_text(slide, 0, title)
            progress.append(ok("Set title", title[:60]))

        if subtitle:
            _set_placeholder_text(slide, 1, subtitle)
            progress.append(ok("Set subtitle", subtitle[:60]))

        prs.save(str(path))
        progress.append(ok(f"Saved {path.name}", "1 slide"))

        if open_after:
            open_file(path)
            progress.append(ok("Opened in default application"))

        return {
            "success": True,
            "op": "create_presentation",
            "output": str(path),
            "output_name": path.name,
            "slide_count": 1,
            "progress": progress,
            "token_estimate": _token_estimate(progress),
        }

    except Exception as exc:
        logger.exception("create_presentation failed")
        return _error(str(exc), "Check output_path is writable.", progress)


def create_from_outline(
    output_path: str,
    slides: list[dict[str, Any]],
    open_after: bool = True,
) -> dict[str, Any]:
    """Create a presentation from a list of slide descriptor dicts."""
    progress: list[dict[str, Any]] = []
    path = Path(output_path)

    if not slides:
        return _error(
            "slides list is empty",
            "Provide at least one slide dict with a 'title' key.",
            progress,
        )

    try:
        _ensure_parent(path)
        progress.append(info("Creating presentation from outline", path.name))

        prs = Presentation()
        slide_count = 0

        for i, slide_def in enumerate(slides):
            raw_title = slide_def.get("title", "")
            content = slide_def.get("content", "")
            layout_hint = slide_def.get("layout", "content").lower()

            if layout_hint == "title":
                # Layout index 0 — Title Slide (title + subtitle)
                layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)
                _set_placeholder_text(slide, 0, raw_title)
                if content:
                    _set_placeholder_text(slide, 1, content)
            else:
                # Layout index 1 — Title and Content (default)
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                _set_placeholder_text(slide, 0, raw_title)
                if content:
                    # Acceptable here: creating new content with no prior formatting
                    try:
                        tf = slide.placeholders[1].text_frame
                        tf.clear()
                        tf.text = content
                    except (IndexError, KeyError, AttributeError):
                        pass

            slide_count += 1
            progress.append(
                ok(
                    f"Added slide {i + 1}: {raw_title[:40] or '(no title)'}",
                    layout_hint,
                )
            )

        prs.save(str(path))
        progress.append(ok(f"Saved {path.name}", f"{slide_count} slides"))

        if open_after:
            open_file(path)
            progress.append(ok("Opened in default application"))

        return {
            "success": True,
            "op": "create_from_outline",
            "output": str(path),
            "output_name": path.name,
            "slide_count": slide_count,
            "progress": progress,
            "token_estimate": _token_estimate(progress),
        }

    except Exception as exc:
        logger.exception("create_from_outline failed")
        return _error(str(exc), "Check output_path is writable and slides list is valid.", progress)


def create_deck_from_data(
    output_path: str,
    title: str,
    data_slides: list[dict[str, Any]],
    open_after: bool = True,
) -> dict[str, Any]:
    """Create a deck with a title slide followed by one content slide per data item."""
    progress: list[dict[str, Any]] = []
    path = Path(output_path)

    if not data_slides:
        return _error(
            "data_slides list is empty",
            "Provide at least one dict with 'heading' and 'bullets' keys.",
            progress,
        )

    try:
        _ensure_parent(path)
        progress.append(info("Creating deck from data", path.name))

        prs = Presentation()

        # Slide 1: title slide
        title_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_layout)
        _set_placeholder_text(title_slide, 0, title)
        progress.append(ok("Added title slide", title[:60]))

        # Subsequent slides: title + content
        content_layout = prs.slide_layouts[1]
        for i, item in enumerate(data_slides):
            heading = item.get("heading", "")
            bullets = item.get("bullets", [])
            body_text = "\n".join(str(b) for b in bullets)

            slide = prs.slides.add_slide(content_layout)
            _set_placeholder_text(slide, 0, heading)

            if body_text:
                # Acceptable here: creating new content, no prior formatting to preserve
                try:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    tf.text = body_text
                except (IndexError, KeyError, AttributeError):
                    pass

            progress.append(
                ok(
                    f"Added slide {i + 2}: {heading[:40] or '(no heading)'}",
                    f"{len(bullets)} bullet(s)",
                )
            )

        slide_count = 1 + len(data_slides)
        prs.save(str(path))
        progress.append(ok(f"Saved {path.name}", f"{slide_count} slides"))

        if open_after:
            open_file(path)
            progress.append(ok("Opened in default application"))

        return {
            "success": True,
            "op": "create_deck_from_data",
            "output": str(path),
            "output_name": path.name,
            "slide_count": slide_count,
            "progress": progress,
            "token_estimate": _token_estimate(progress),
        }

    except Exception as exc:
        logger.exception("create_deck_from_data failed")
        return _error(str(exc), "Check output_path is writable and data_slides is valid.", progress)


def create_from_template(
    template_path: str,
    output_path: str,
    open_after: bool = True,
) -> dict[str, Any]:
    """Copy an existing .pptx as a new presentation starting point."""
    progress: list[dict[str, Any]] = []
    tmpl = Path(template_path)
    path = Path(output_path)

    if not tmpl.exists():
        return _error(
            f"Template not found: {tmpl.name}",
            "Check that template_path is absolute and the file exists.",
            progress,
        )

    if tmpl.suffix.lower() != ".pptx":
        return _error(
            f"Expected .pptx template, got {tmpl.suffix}",
            "Provide a .pptx file as the template.",
            progress,
        )

    try:
        _ensure_parent(path)
        progress.append(info(f"Copying template {tmpl.name}", path.name))

        shutil.copy2(str(tmpl), str(path))
        progress.append(ok(f"Copied template to {path.name}"))

        # Open the copy to report slide count
        prs = Presentation(str(path))
        slide_count = len(prs.slides)
        progress.append(ok(f"Template has {slide_count} slide(s)"))

        if open_after:
            open_file(path)
            progress.append(ok("Opened in default application"))

        return {
            "success": True,
            "op": "create_from_template",
            "output": str(path),
            "output_name": path.name,
            "template": str(tmpl),
            "template_name": tmpl.name,
            "slide_count": slide_count,
            "progress": progress,
            "token_estimate": _token_estimate(progress),
        }

    except Exception as exc:
        logger.exception("create_from_template failed")
        return _error(str(exc), "Check template_path and output_path are valid.", progress)
