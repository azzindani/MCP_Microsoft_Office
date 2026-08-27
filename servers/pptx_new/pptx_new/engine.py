"""PPTX New engine — create PowerPoint presentations from scratch."""

from __future__ import annotations

import logging
import shutil
import sys
from pathlib import Path
from typing import Any

# Allow imports from the repo root (shared/)
_ROOT = Path(__file__).resolve().parents[2]
if str(_ROOT) not in sys.path:
    sys.path.insert(0, str(_ROOT))

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402,F401

from shared.file_utils import embed_content, hint_for_message  # noqa: E402
from shared.platform_utils import open_file, resolve_output_path  # noqa: E402
from shared.progress import describe_error, fail, info, ok, warn  # noqa: E402
from shared.slide_text import strip_list_markers  # noqa: E402
from shared.template_fill import ordered_pairs, resolve_targets, substitute_once  # noqa: E402

logging.basicConfig(stream=sys.stderr, level=logging.WARNING)
logger = logging.getLogger(__name__)


# ─── Internal helpers ────────────────────────────────────────────────────────


def _ensure_parent(path: Path) -> None:
    """Create parent directory if it does not exist."""
    path.parent.mkdir(parents=True, exist_ok=True)


def _token_estimate(obj: Any) -> int:
    return len(str(obj)) // 4


def _error(
    msg: str,
    hint: str,
    progress: list[dict[str, Any]],
) -> dict[str, Any]:
    progress.append(fail(msg))
    return {
        "success": False,
        "error": describe_error(msg),
        "hint": hint_for_message(msg, hint),
        "progress": progress,
        "token_estimate": _token_estimate(progress),
    }


def _set_placeholder_text(slide: Any, index: int, text: str) -> None:
    """
    Set placeholder text safely by index.

    Uses direct .text assignment — acceptable when creating new content
    because there is no existing run-level formatting to preserve.
    """
    try:
        slide.placeholders[index].text = text
    except (IndexError, KeyError, AttributeError):
        pass


# ─── Public engine functions ─────────────────────────────────────────────────


def create_presentation(
    output_path: str,
    title: str = "",
    subtitle: str = "",
    open_after: bool = True,
    return_content: bool = False,
) -> dict[str, Any]:
    """Create a blank presentation with a single title slide."""
    progress: list[dict[str, Any]] = []
    path = resolve_output_path(output_path, "presentation.pptx")

    try:
        _ensure_parent(path)
        progress.append(info("Creating new presentation", path.name))

        prs = Presentation()
        # Layout index 0 is the Title Slide layout in default templates
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        if title:
            _set_placeholder_text(slide, 0, title)
            progress.append(ok("Set title", title[:60]))

        if subtitle:
            _set_placeholder_text(slide, 1, subtitle)
            progress.append(ok("Set subtitle", subtitle[:60]))

        prs.save(str(path))
        progress.append(ok(f"Saved {path.name}", "1 slide"))

        if open_after:
            open_file(path)
            progress.append(ok("Opened in default application"))

        result: dict[str, Any] = {
            "success": True,
            "op": "create_presentation",
            "output": str(path),
            "output_name": path.name,
            "slide_count": 1,
            "progress": progress,
        }
        embed_content(result, path, return_content)
        result["token_estimate"] = _token_estimate(result)
        return result

    except Exception as exc:
        logger.exception("create_presentation failed")
        return _error(str(exc), "Check output_path is writable.", progress)


# Both deck builders take "a list of slide dicts" and each invented its own key
# names: create_from_outline read title/content, create_deck_from_data read
# heading/bullets. data_slides is typed list[dict[str, Any]], so tools/list shows
# an opaque array and the caller cannot discover which spelling a tool wants --
# it guesses, and the wrong guess returns success:true with every slide title
# blank. Both spellings are now accepted by both tools, and a slide that yields
# no heading at all says so in progress instead of going quietly.
_HEADING_KEYS = ("title", "heading", "header", "name")

# The layouts create_from_outline can build, and the spellings that mean each.
# The keys are matched after lowercasing and turning spaces and hyphens into
# underscores, so "Title Only" and "title-only" both arrive as title_only.
# PowerPoint's own layout names are on the left because they are what a caller
# reads off a deck and types back.
_OUTLINE_LAYOUTS: dict[str, str] = {
    "title": "title",
    "title_slide": "title",
    "cover": "title",
    "title_only": "title_only",
    "titleonly": "title_only",
    "heading_only": "title_only",
    "content": "content",
    "title_and_content": "content",
    "bullets": "content",
    "body": "content",
}
# `items` and `points` are what a caller writes for a bullet list about as often
# as `bullets`; without them create_deck_from_data built the slide, set its
# heading, and left the body blank under success.
_BODY_KEYS = ("content", "bullets", "body", "text", "items", "points", "lines")


def _slide_heading(item: dict[str, Any]) -> str:
    for key in _HEADING_KEYS:
        value = item.get(key)
        if value:
            return str(value)
    return ""


def _slide_body(item: dict[str, Any]) -> str:
    """Return the slide body, joining a list of bullets into lines."""
    for key in _BODY_KEYS:
        value = item.get(key)
        if not value:
            continue
        if isinstance(value, (list, tuple)):
            return "\n".join(str(v) for v in value)
        return str(value)
    return ""


def _note_unnamed(progress: list[dict[str, Any]], index: int, item: dict[str, Any]) -> None:
    """A slide with no recognised heading key is almost always a typo, not intent."""
    if _slide_heading(item) or not isinstance(item, dict):
        return
    progress.append(
        warn(
            f"Slide {index} has no heading",
            f"keys seen: {', '.join(sorted(map(str, item)))} — use 'title' or 'heading'",
        )
    )


def create_from_outline(
    output_path: str,
    slides: list[dict[str, Any]],
    open_after: bool = True,
    return_content: bool = False,
) -> dict[str, Any]:
    """Create a presentation from a list of slide descriptor dicts."""
    progress: list[dict[str, Any]] = []
    path = resolve_output_path(output_path, "presentation.pptx")

    if not slides:
        return _error(
            "slides list is empty",
            "Provide at least one slide dict with a 'title' (or 'heading') key.",
            progress,
        )

    try:
        _ensure_parent(path)
        progress.append(info("Creating presentation from outline", path.name))

        prs = Presentation()
        slide_count = 0

        for i, slide_def in enumerate(slides):
            _note_unnamed(progress, i + 1, slide_def)
            raw_title = _slide_heading(slide_def)
            content = _slide_body(slide_def)
            # Only the exact token "title" was ever recognised; everything else
            # fell through to the content layout and was then echoed back
            # verbatim in the progress line, so layout="Title Only" reported
            # "title only" while the file used Title and Content, and
            # layout="zzz" reported "zzz". The response asserted a layout the
            # artifact contradicted, which is worse than ignoring the argument
            # quietly. Normalised to what will actually be applied, and an
            # unrecognised value is refused rather than silently downgraded.
            layout_hint = _OUTLINE_LAYOUTS.get(
                str(slide_def.get("layout", "content")).strip().lower().replace("-", "_").replace(" ", "_"),
                "",
            )
            if not layout_hint:
                given = slide_def.get("layout")
                return _error(
                    f"Slide {i + 1}: unknown layout {given!r}.",
                    f"Use one of: {', '.join(sorted(set(_OUTLINE_LAYOUTS.values())))}. "
                    f"Accepted spellings: {', '.join(sorted(_OUTLINE_LAYOUTS))}.",
                    progress,
                )

            if layout_hint == "title_only":
                # Layout index 5 — Title Only, a real layout this used to
                # silently turn into Title and Content.
                layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(layout)
                _set_placeholder_text(slide, 0, raw_title)
            elif layout_hint == "title":
                # Layout index 0 — Title Slide (title + subtitle)
                layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)
                _set_placeholder_text(slide, 0, raw_title)
                if content:
                    _set_placeholder_text(slide, 1, content)
            else:
                # Layout index 1 — Title and Content (default)
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                _set_placeholder_text(slide, 0, raw_title)
                if content:
                    # Acceptable here: creating new content with no prior formatting
                    try:
                        tf = slide.placeholders[1].text_frame  # type: ignore[reportAttributeAccessIssue]
                        tf.clear()
                        tf.text = strip_list_markers(content)
                    except (IndexError, KeyError, AttributeError):
                        pass

            slide_count += 1
            progress.append(
                ok(
                    f"Added slide {i + 1}: {raw_title[:40] or '(no title)'}",
                    layout_hint,
                )
            )

        prs.save(str(path))
        progress.append(ok(f"Saved {path.name}", f"{slide_count} slides"))

        if open_after:
            open_file(path)
            progress.append(ok("Opened in default application"))

        result: dict[str, Any] = {
            "success": True,
            "op": "create_from_outline",
            "output": str(path),
            "output_name": path.name,
            "slide_count": slide_count,
            "progress": progress,
        }
        embed_content(result, path, return_content)
        result["token_estimate"] = _token_estimate(result)
        return result

    except Exception as exc:
        logger.exception("create_from_outline failed")
        return _error(str(exc), "Check output_path is writable and slides list is valid.", progress)


def create_deck_from_data(
    output_path: str,
    title: str,
    data_slides: list[dict[str, Any]],
    open_after: bool = True,
    return_content: bool = False,
) -> dict[str, Any]:
    """Create a deck with a title slide followed by one content slide per data item."""
    progress: list[dict[str, Any]] = []
    path = resolve_output_path(output_path, "presentation.pptx")

    if not data_slides:
        return _error(
            "data_slides list is empty",
            "Provide at least one dict with a 'title' (or 'heading') and 'bullets' keys.",
            progress,
        )

    try:
        _ensure_parent(path)
        progress.append(info("Creating deck from data", path.name))

        prs = Presentation()

        # Slide 1: title slide
        title_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_layout)
        _set_placeholder_text(title_slide, 0, title)
        progress.append(ok("Added title slide", title[:60]))

        # Subsequent slides: title + content
        content_layout = prs.slide_layouts[1]
        for i, item in enumerate(data_slides):
            _note_unnamed(progress, i + 2, item)
            heading = _slide_heading(item)
            body_text = _slide_body(item)
            # Was item.get("bullets") alone, so a list given under any of the
            # other accepted body keys reached the slide as no bullets at all.
            raw_bullets = next((item.get(k) for k in _BODY_KEYS if isinstance(item.get(k), list | tuple)), None)
            bullets = list(raw_bullets) if raw_bullets else ([] if not body_text else body_text.split("\n"))

            slide = prs.slides.add_slide(content_layout)
            _set_placeholder_text(slide, 0, heading)

            if body_text:
                # Acceptable here: creating new content, no prior formatting to preserve
                try:
                    tf = slide.placeholders[1].text_frame  # type: ignore[reportAttributeAccessIssue]
                    tf.clear()
                    tf.text = strip_list_markers(body_text)
                except (IndexError, KeyError, AttributeError):
                    pass

            progress.append(
                ok(
                    f"Added slide {i + 2}: {heading[:40] or '(no heading)'}",
                    f"{len(bullets)} bullet(s)",
                )
            )

        slide_count = 1 + len(data_slides)
        prs.save(str(path))
        progress.append(ok(f"Saved {path.name}", f"{slide_count} slides"))

        if open_after:
            open_file(path)
            progress.append(ok("Opened in default application"))

        result: dict[str, Any] = {
            "success": True,
            "op": "create_deck_from_data",
            "output": str(path),
            "output_name": path.name,
            "slide_count": slide_count,
            "progress": progress,
        }
        embed_content(result, path, return_content)
        result["token_estimate"] = _token_estimate(result)
        return result

    except Exception as exc:
        logger.exception("create_deck_from_data failed")
        return _error(str(exc), "Check output_path is writable and data_slides is valid.", progress)


def _deck_text(prs: Any) -> str:
    """Every run of text a substitution can reach, as one string.

    Only used to decide what a key should match, so joining is enough.
    """
    parts: list[str] = []

    def _walk(shape: Any) -> None:
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                _walk(child)
            return
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.extend(r.text for p in cell.text_frame.paragraphs for r in p.runs)
            return
        if getattr(shape, "has_text_frame", False):
            parts.extend(r.text for p in shape.text_frame.paragraphs for r in p.runs)

    for slide in prs.slides:
        for shape in slide.shapes:
            _walk(shape)
    return "\n".join(parts)


def _substitute_in_text_frame(text_frame: Any, pairs: list[tuple[str, str]]) -> int:
    """Replace every target in one pass per run.

    A sequential pass over the keys let one key consume another -- `platform`
    eating the middle of `platform_spend`, leaving a slide reading
    "Google Ads_spend" -- and let a value produced by one key be re-replaced by
    the next. substitute_once() matches all targets in a single regex pass, so
    neither can happen.
    """
    applied = 0
    for para in text_frame.paragraphs:
        for run in para.runs:
            updated, made = substitute_once(run.text, pairs)
            if updated != run.text:
                run.text = updated
                # `made`, not 1. substitute_once already counts the
                # replacements it performed and that count was being thrown
                # away in favour of counting *runs touched*: a slide reading
                # "Channel {channel} delivered {impressions} impressions with
                # CTR {ctr}" is a single run, so all three placeholders were
                # filled correctly and the response said
                # substitutions_applied: 1.
                applied += made
    return applied


def _substitute_in_shape(shape: Any, pairs: list[tuple[str, str]]) -> int:
    """Replace placeholder text inside a shape, run by run.

    Assigning to text_frame.text would flatten the run structure and lose every
    font, colour and size the template author set -- the rule pptx_basic's
    _set_shape_text already follows. Replacing within each run keeps that
    formatting, so a filled template still looks like the template.

    Recurses into grouped shapes and table cells, which is where template
    placeholders usually live.
    """
    applied = 0
    if hasattr(shape, "shapes"):  # grouped shape
        for child in shape.shapes:
            applied += _substitute_in_shape(child, pairs)
        return applied

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                applied += _substitute_in_text_frame(cell.text_frame, pairs)
        return applied

    if getattr(shape, "has_text_frame", False):
        applied += _substitute_in_text_frame(shape.text_frame, pairs)
    return applied


def create_from_template(
    template_path: str,
    output_path: str,
    substitutions: dict | None = None,
    open_after: bool = True,
    return_content: bool = False,
) -> dict[str, Any]:
    """Copy an existing .pptx and apply {key: value} text substitutions."""
    progress: list[dict[str, Any]] = []
    tmpl = Path(template_path)
    path = resolve_output_path(output_path, "presentation.pptx")

    if not tmpl.exists():
        return _error(
            f"Template not found: {tmpl.name}",
            "Check that template_path is absolute and the file exists.",
            progress,
        )

    if tmpl.suffix.lower() != ".pptx":
        return _error(
            f"Expected .pptx template, got {tmpl.suffix}",
            "Provide a .pptx file as the template.",
            progress,
        )

    try:
        _ensure_parent(path)
        progress.append(info(f"Copying template {tmpl.name}", path.name))

        shutil.copy2(str(tmpl), str(path))
        progress.append(ok(f"Copied template to {path.name}"))

        # Open the copy to report slide count
        prs = Presentation(str(path))
        slide_count = len(prs.slides)
        progress.append(ok(f"Template has {slide_count} slide(s)"))

        subs = substitutions or {}
        if not isinstance(subs, dict):
            return _error(
                "substitutions must be a dict",
                'Pass a dict like {"PLACEHOLDER": "replacement value"}.',
                progress,
            )

        # Resolve what each key should match against the whole deck before
        # replacing anything -- a key that only appears delimited is matched in
        # that form, so a caller passing `platform` at a `{platform}` template
        # no longer leaves the braces behind.
        targets, notes = resolve_targets(_deck_text(prs), subs)
        for note in notes:
            progress.append(info(note))
        pairs = ordered_pairs(targets, subs)
        for key in subs:
            if str(key) not in targets:
                progress.append(warn(f"Placeholder '{key}' not found in template"))

        applied = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                applied += _substitute_in_shape(shape, pairs)
        if subs:
            prs.save(str(path))
            progress.append(ok(f"Applied {applied} substitution(s)", f"{len(subs)} key(s) searched"))

        if open_after:
            open_file(path)
            progress.append(ok("Opened in default application"))

        result: dict[str, Any] = {
            "success": True,
            "op": "create_from_template",
            "output": str(path),
            "output_name": path.name,
            "template": str(tmpl),
            "template_name": tmpl.name,
            "slide_count": slide_count,
            "substitutions_applied": applied,
            "progress": progress,
        }
        embed_content(result, path, return_content)
        result["token_estimate"] = _token_estimate(result)
        return result

    except Exception as exc:
        logger.exception("create_from_template failed")
        return _error(str(exc), "Check template_path and output_path are valid.", progress)


def create_agenda(
    output_path: str,
    meeting_title: str,
    date: str,
    items: list,
    presenter: str = "",
    open_after: bool = True,
    return_content: bool = False,
) -> dict[str, Any]:
    """Create a meeting agenda presentation with title and agenda slides."""
    progress: list[dict[str, Any]] = []
    path = resolve_output_path(output_path, "presentation.pptx")

    try:
        _ensure_parent(path)
        progress.append(info("Creating agenda presentation", path.name))

        prs = Presentation()

        # Slide 1 — Title slide (layout index 0)
        title_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_layout)
        _set_placeholder_text(title_slide, 0, meeting_title)
        subtitle = date
        if presenter:
            subtitle = f"{date}\nPresented by: {presenter}"
        _set_placeholder_text(title_slide, 1, subtitle)
        progress.append(ok("Added title slide", meeting_title[:60]))

        # Slide 2 — Agenda (layout index 1)
        content_layout = prs.slide_layouts[1]
        agenda_slide = prs.slides.add_slide(content_layout)
        _set_placeholder_text(agenda_slide, 0, "Agenda")

        agenda_lines = []
        for item in items:
            topic = item.get("topic", "")
            duration = item.get("duration", "")
            owner = item.get("owner", "")
            agenda_lines.append(f"• {topic} ({duration}) — {owner}")
        agenda_text = "\n".join(agenda_lines)

        try:
            tf = agenda_slide.placeholders[1].text_frame  # type: ignore[reportAttributeAccessIssue]
            tf.clear()
            tf.text = strip_list_markers(agenda_text)
        except (IndexError, KeyError, AttributeError):
            pass

        progress.append(ok("Added agenda slide", f"{len(items)} item(s)"))

        prs.save(str(path))
        progress.append(ok(f"Saved {path.name}", "2 slides"))

        if open_after:
            open_file(path)
            progress.append(ok("Opened in default application"))

        result: dict[str, Any] = {
            "success": True,
            "op": "create_agenda",
            "output": str(path),
            "output_name": path.name,
            "slide_count": 2,
            "item_count": len(items),
            "progress": progress,
        }
        embed_content(result, path, return_content)
        result["token_estimate"] = _token_estimate(result)
        return result

    except Exception as exc:
        logger.exception("create_agenda failed")
        return _error(str(exc), "Check output_path is writable and items is a valid list.", progress)


def _heading_level(para: Any) -> int:
    """The heading depth of a paragraph, or 0 if it is body text."""
    name = getattr(getattr(para, "style", None), "name", "") or ""
    if not name.startswith("Heading"):
        return 0
    tail = name[len("Heading") :].strip()
    return int(tail) if tail.isdigit() else 0


def create_from_docx(
    docx_path: str,
    output_path: str,
    max_slides: int = 20,
    open_after: bool = True,
    return_content: bool = False,
) -> dict[str, Any]:
    """Convert a Word document outline into a PowerPoint presentation."""
    progress: list[dict[str, Any]] = []
    path = resolve_output_path(output_path, "presentation.pptx")
    docx_file = Path(docx_path)

    if not docx_file.exists():
        return _error(
            f"File not found: {docx_file.name}",
            "Check that docx_path is absolute and the file exists.",
            progress,
        )
    if docx_file.suffix.lower() != ".docx":
        return _error(
            f"Expected .docx file, got {docx_file.suffix}",
            "Provide a .docx file as the source document.",
            progress,
        )

    try:
        from docx import Document as DocxDocument  # type: ignore[import-untyped]

        _ensure_parent(path)
        progress.append(info(f"Reading {docx_file.name}", str(docx_file)))

        doc = DocxDocument(str(docx_file))
        paragraphs = doc.paragraphs
        source_para_count = len(paragraphs)
        progress.append(ok(f"Opened {docx_file.name}", f"{source_para_count} paragraphs"))

        # Extract slide structure from headings
        slides_data: list[dict[str, Any]] = []
        current_slide: dict[str, Any] | None = None
        has_headings = any(p.style.name.startswith("Heading") for p in paragraphs)  # type: ignore[reportOptionalMemberAccess]

        # A new slide used to start only on Heading 1, with Heading 2 demoted to
        # a bullet. create_from_sections -- this fleet's own .docx builder --
        # writes the document title as Heading 1 and every section heading as
        # Heading 2, so a three-section document round-tripped through the two
        # tools came back as ONE slide with all three sections crammed into its
        # body, under success and a `slide_count: 1` that was accurate about a
        # deck nobody asked for. max_slides=20 sat in the schema promising
        # otherwise.
        #
        # Split on the shallowest heading level that occurs more than once
        # instead: several Heading 1s still split on Heading 1 exactly as
        # before, and a lone title over repeated Heading 2s splits on the
        # Heading 2s, which is the structure the document is actually showing.
        levels = [_heading_level(p) for p in paragraphs if p.text.strip()]
        present = sorted({lv for lv in levels if lv})
        repeated = [lv for lv in present if levels.count(lv) > 1]
        split_level = repeated[0] if repeated else (present[0] if present else 1)
        if split_level != 1:
            progress.append(
                info(
                    f"Splitting slides on Heading {split_level}",
                    f"Heading 1 appears {levels.count(1)} time(s), so it titles the deck rather than each slide",
                )
            )

        if has_headings:
            for para in paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                level = _heading_level(para)
                if level and level <= split_level:
                    current_slide = {"title": text, "content_lines": []}
                    slides_data.append(current_slide)
                elif level:
                    # A heading below the split level is a sub-point of its slide.
                    if current_slide is not None:
                        current_slide["content_lines"].append(f"  • {text}")
                    else:
                        current_slide = {"title": text, "content_lines": []}
                        slides_data.append(current_slide)
                else:
                    # Normal / Body Text — add to current slide content
                    if current_slide is not None:
                        current_slide["content_lines"].append(text)
        else:
            # No headings: group every non-empty paragraph into slides of 5
            bullets: list[str] = [p.text.strip() for p in paragraphs if p.text.strip()]
            group_size = 5
            for i in range(0, len(bullets), group_size):
                group = bullets[i : i + group_size]
                title = group[0] if group else f"Slide {len(slides_data) + 1}"
                slides_data.append(
                    {
                        "title": title,
                        "content_lines": group[1:],
                    }
                )

        # Cap at max_slides
        slides_data = slides_data[:max_slides]

        prs = Presentation()
        slide_count = 0

        for i, slide_def in enumerate(slides_data):
            title_text = slide_def["title"]
            content_text = "\n".join(slide_def.get("content_lines", []))

            if i == 0:
                layout = prs.slide_layouts[0]  # Title Slide
                slide = prs.slides.add_slide(layout)
                _set_placeholder_text(slide, 0, title_text)
                if content_text:
                    _set_placeholder_text(slide, 1, content_text)
            else:
                layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(layout)
                _set_placeholder_text(slide, 0, title_text)
                if content_text:
                    try:
                        tf = slide.placeholders[1].text_frame  # type: ignore[reportAttributeAccessIssue]
                        tf.clear()
                        tf.text = strip_list_markers(content_text)
                    except (IndexError, KeyError, AttributeError):
                        pass

            slide_count += 1
            progress.append(ok(f"Added slide {i + 1}: {title_text[:40] or '(no title)'}"))

        prs.save(str(path))
        progress.append(ok(f"Saved {path.name}", f"{slide_count} slides"))

        if open_after:
            open_file(path)
            progress.append(ok("Opened in default application"))

        result: dict[str, Any] = {
            "success": True,
            "op": "create_from_docx",
            "output": str(path),
            "output_name": path.name,
            "slide_count": slide_count,
            "source_paragraph_count": source_para_count,
            "progress": progress,
        }
        embed_content(result, path, return_content)
        result["token_estimate"] = _token_estimate(result)
        return result

    except Exception as exc:
        logger.exception("create_from_docx failed")
        return _error(str(exc), "Check docx_path and output_path are valid.", progress)
