"""PPTX Basic engine — pure python-pptx logic, no MCP imports."""

from pathlib import Path
from typing import Any

from shared.address_resolver import build_pptx_index
from shared.doc_diff import diff_pptx
from shared.file_utils import hint_for_error, resolve_path
from shared.live_edit import notify_reload
from shared.platform_utils import open_file
from shared.progress import describe_error, fail, index_range, info, ok
from shared.receipt import append_receipt
from shared.slide_text import strip_list_markers
from shared.version_control import get_history, snapshot

# ─── Helpers ─────────────────────────────────────────────────────────────────


def _not_found(path: Path, progress: list[dict[str, Any]]) -> dict[str, Any]:
    progress.append(fail("File not found", str(path)))
    return {
        "success": False,
        "progress": progress,
        "error": f"File not found: {path}",
        "hint": "Check that file_path is absolute and the file exists.",
        "token_estimate": len(str(progress)) // 4,
    }


def _wrong_type(path: Path, expected: str, progress: list[dict[str, Any]]) -> dict[str, Any]:
    progress.append(fail(f"Wrong file type: {path.suffix}", f"Expected {expected}"))
    return {
        "success": False,
        "progress": progress,
        "error": f"Expected {expected} file, got {path.suffix}",
        "hint": "Use the correct server for this file type.",
        "token_estimate": len(str(progress)) // 4,
    }


def _error(msg: str, hint: str, progress: list[dict[str, Any]], backup: str | None = None) -> dict[str, Any]:
    result: dict[str, Any] = {
        "success": False,
        "progress": progress,
        "error": describe_error(msg),
        "hint": hint,
        "token_estimate": len(str(progress)) // 4,
    }
    if backup is not None:
        result["backup"] = backup
    return result


def _shape_type_str(shape: Any) -> str:
    """Return a human-readable shape type string."""
    if shape.has_text_frame:
        return "text"
    if shape.has_table:
        return "table"
    if shape.has_chart:
        return "chart"
    # MSO_SHAPE_TYPE.PICTURE == 13
    if shape.shape_type == 13:
        return "image"
    return "other"


def _set_shape_text(shape: Any, new_text: str) -> None:
    """
    Set text in a shape using run-level editing.

    Never assigns to text_frame.text directly — that destroys formatting.
    Instead clears all paragraphs/runs and sets a single run on the first
    paragraph, preserving the paragraph and run XML structure.
    """
    tf = shape.text_frame
    # Remove all paragraphs except the first
    while len(tf.paragraphs) > 1:
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)
    # Remove all runs from first paragraph
    first_para = tf.paragraphs[0]
    for run in first_para.runs:
        run._r.getparent().remove(run._r)
    # Add a single run with the new text
    run = first_para.add_run()
    run.text = new_text


# ─── Tool functions ───────────────────────────────────────────────────────────


def read_presentation(file_path: str) -> dict[str, Any]:
    """Return full index: slide count, titles, shape info, and layouts."""
    progress: list[dict[str, Any]] = []
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        path = resolve_path(file_path)
        if not path.exists():
            return _not_found(path, progress)
        if path.suffix.lower() != ".pptx":
            return _wrong_type(path, ".pptx", progress)

        prs = Presentation(str(path))
        index = build_pptx_index(prs)
        slide_count = index["slide_count"]
        progress.append(ok(f"Read {path.name}", f"{slide_count} slides"))

        return {
            "success": True,
            "file": str(path),
            "slide_count": slide_count,
            "available_layouts": index["available_layouts"],
            "slides": index["slides"],
            "progress": progress,
            "token_estimate": len(str(index)) // 4,
        }
    except Exception as e:
        progress.append(fail(str(e)))
        return _error(str(e), "Check that file_path points to a valid .pptx file.", progress)


def read_slide(file_path: str, slide_index: int) -> dict[str, Any]:
    """Return all shapes on one slide with name, type, index, text."""
    progress: list[dict[str, Any]] = []
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        path = resolve_path(file_path)
        if not path.exists():
            return _not_found(path, progress)
        if path.suffix.lower() != ".pptx":
            return _wrong_type(path, ".pptx", progress)

        prs = Presentation(str(path))
        total = len(prs.slides)
        progress.append(ok(f"Opened {path.name}", f"{total} slides"))

        if slide_index < 0 or slide_index >= total:
            progress.append(fail(f"Slide index {slide_index} out of range"))
            return _error(
                f"slide_index {slide_index} out of range {index_range(total, 'slides')}",
                "Use read_presentation to get current slide count.",
                progress,
            )

        slide = prs.slides[slide_index]
        shapes = []
        for idx, shape in enumerate(slide.shapes):
            shape_dict: dict[str, Any] = {
                "name": shape.name,
                "index": idx,
                "type": _shape_type_str(shape),
            }
            if shape.has_text_frame:
                shape_dict["text"] = shape.text_frame.text  # type: ignore[reportAttributeAccessIssue]
                shape_dict["paragraphs"] = [
                    {
                        "index": p_idx,
                        "text": para.text,
                        "runs": [
                            {
                                "text": run.text,
                                "bold": run.font.bold,
                                "italic": run.font.italic,
                            }
                            for run in para.runs
                        ],
                    }
                    for p_idx, para in enumerate(shape.text_frame.paragraphs)  # type: ignore[reportAttributeAccessIssue]
                ]
            shapes.append(shape_dict)

        progress.append(ok(f"Read slide {slide_index}", f"{len(shapes)} shapes"))

        return {
            "success": True,
            "file": str(path),
            "slide_index": slide_index,
            "shape_count": len(shapes),
            "shapes": shapes,
            "progress": progress,
            "token_estimate": len(str(shapes)) // 4,
        }
    except Exception as e:
        progress.append(fail(str(e)))
        return _error(str(e), "Check slide_index is valid.", progress)


def search_slides(file_path: str, query: str) -> dict[str, Any]:
    """Scan all slides for query text, return slide index and shape name."""
    progress: list[dict[str, Any]] = []
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        if not query:
            return _error(
                "query cannot be empty",
                "Provide a non-empty search string.",
                progress,
            )

        path = resolve_path(file_path)
        if not path.exists():
            return _not_found(path, progress)
        if path.suffix.lower() != ".pptx":
            return _wrong_type(path, ".pptx", progress)

        prs = Presentation(str(path))
        total = len(prs.slides)
        progress.append(ok(f"Opened {path.name}", f"{total} slides"))

        q_lower = query.lower()
        matches: list[dict[str, Any]] = []

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text  # type: ignore[reportAttributeAccessIssue]
                    if q_lower in text.lower():
                        matches.append(
                            {
                                "slide_index": slide_idx,
                                "shape_name": shape.name,
                                "text": text,
                            }
                        )

        progress.append(
            ok(
                f"Found {len(matches)} match{'es' if len(matches) != 1 else ''}",
                f"query: {query!r}",
            )
        )

        return {
            "success": True,
            "query": query,
            "matches": matches,
            "total_slides_scanned": total,
            "progress": progress,
            "token_estimate": len(str(matches)) // 4,
        }
    except Exception as e:
        progress.append(fail(str(e)))
        return _error(str(e), "Check that file_path points to a valid .pptx file.", progress)


def read_slide_text(file_path: str, slide_index: int) -> dict[str, Any]:
    """Return text content of all shapes on one slide (no formatting)."""
    progress: list[dict[str, Any]] = []
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        path = resolve_path(file_path)
        if not path.exists():
            return _not_found(path, progress)
        if path.suffix.lower() != ".pptx":
            return _wrong_type(path, ".pptx", progress)

        prs = Presentation(str(path))
        total = len(prs.slides)
        progress.append(ok(f"Opened {path.name}", f"{total} slides"))

        if slide_index < 0 or slide_index >= total:
            progress.append(fail(f"Slide index {slide_index} out of range"))
            return _error(
                f"slide_index {slide_index} out of range {index_range(total, 'slides')}",
                "Use read_presentation to get current slide count.",
                progress,
            )

        slide = prs.slides[slide_index]
        shapes = [{"name": shape.name, "text": shape.text_frame.text} for shape in slide.shapes if shape.has_text_frame]  # type: ignore[reportAttributeAccessIssue]

        progress.append(ok(f"Read slide {slide_index} text", f"{len(shapes)} shapes"))

        return {
            "success": True,
            "slide_index": slide_index,
            "shapes": shapes,
            "progress": progress,
            "token_estimate": len(str(shapes)) // 4,
        }
    except Exception as e:
        progress.append(fail(str(e)))
        return _error(str(e), "Check slide_index is valid.", progress)


def set_text(
    file_path: str, slide_index: int, shape_name: str, new_text: str, open_after: bool = False
) -> dict[str, Any]:
    """Replace text in a shape using run-level editing."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    path: Path | None = None
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        path = resolve_path(file_path)
        if not path.exists():
            return _not_found(path, progress)
        if path.suffix.lower() != ".pptx":
            return _wrong_type(path, ".pptx", progress)

        prs = Presentation(str(path))
        total = len(prs.slides)
        progress.append(ok(f"Opened {path.name}", f"{total} slides"))

        if slide_index < 0 or slide_index >= total:
            progress.append(fail(f"Slide index {slide_index} out of range"))
            return _error(
                f"slide_index {slide_index} out of range {index_range(total, 'slides')}",
                "Use read_presentation to get current slide count.",
                progress,
            )

        slide = prs.slides[slide_index]
        target_shape = None
        for shape in slide.shapes:
            if shape.name == shape_name:
                target_shape = shape
                break

        if target_shape is None:
            available = [s.name for s in slide.shapes]
            progress.append(fail(f"Shape '{shape_name}' not found", f"Available: {available}"))
            return _error(
                f"Shape '{shape_name}' not found on slide {slide_index}",
                f"Available shapes: {available}",
                progress,
            )

        if not target_shape.has_text_frame:
            progress.append(fail(f"Shape '{shape_name}' has no text frame"))
            return _error(
                f"Shape '{shape_name}' has no text frame",
                "Choose a shape that contains text.",
                progress,
            )

        old_text = target_shape.text_frame.text  # type: ignore[reportAttributeAccessIssue]
        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        _set_shape_text(target_shape, new_text)
        prs.save(str(path))
        if open_after:
            open_file(path)
        reload_step = notify_reload(str(path), "pptx")
        progress.append(reload_step)

        progress.append(
            ok(
                f"Updated '{shape_name}' on slide {slide_index}",
                new_text[:40],
            )
        )

        append_receipt(
            str(path),
            "set_text",
            "pptx_basic",
            {"slide_index": slide_index, "shape_name": shape_name, "new_text": new_text},
            f"✔ Updated '{shape_name}' on slide {slide_index}",
            backup,
            True,
        )

        return {
            "success": True,
            "op": "set_text",
            "slide_index": slide_index,
            "shape_name": shape_name,
            "old_text": old_text,
            "new_text": new_text,
            "backup": backup,
            "progress": progress,
            "token_estimate": len(str(progress)) // 4,
        }
    except Exception as e:
        progress.append(fail(str(e)))
        append_receipt(
            file_path,
            "set_text",
            "pptx_basic",
            {"slide_index": slide_index, "shape_name": shape_name, "new_text": new_text},
            f"✘ {e}",
            backup,
            False,
        )
        return _error(
            str(e),
            hint_for_error(e, path),
            progress,
            backup,
        )


def add_slide(
    file_path: str, layout_name: str, title: str = "", body: str = "", open_after: bool = False
) -> dict[str, Any]:
    """Append a new slide with layout, title, and body text."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    path: Path | None = None
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        path = resolve_path(file_path)
        if not path.exists():
            return _not_found(path, progress)
        if path.suffix.lower() != ".pptx":
            return _wrong_type(path, ".pptx", progress)

        prs = Presentation(str(path))
        total = len(prs.slides)
        progress.append(ok(f"Opened {path.name}", f"{total} slides"))

        # Validate layout name
        available_layouts = {layout.name: layout for layout in prs.slide_layouts}
        if layout_name not in available_layouts:
            available = list(available_layouts.keys())
            progress.append(fail(f"Layout '{layout_name}' not found", str(available)))
            return _error(
                f"Layout '{layout_name}' not found",
                f"Available layouts: {available}",
                progress,
            )

        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        layout = available_layouts[layout_name]
        new_slide = prs.slides.add_slide(layout)
        new_index = len(prs.slides) - 1

        # Set title if provided
        if title:
            for shape in new_slide.shapes:
                if shape.has_text_frame and shape.name.startswith("Title"):
                    _set_shape_text(shape, title)
                    break

        # Set body if provided
        if body:
            for shape in new_slide.shapes:
                if shape.has_text_frame and not shape.name.startswith("Title"):
                    _set_shape_text(shape, body)
                    break

        prs.save(str(path))
        if open_after:
            open_file(path)
        reload_step = notify_reload(str(path), "pptx")
        progress.append(reload_step)
        progress.append(ok(f"Added slide {new_index} '{title}'", layout_name))

        append_receipt(
            str(path),
            "add_slide",
            "pptx_basic",
            {"layout_name": layout_name, "title": title, "body": body},
            f"✔ Added slide {new_index} '{title}'",
            backup,
            True,
        )

        return {
            "success": True,
            "op": "add_slide",
            "slide_index": new_index,
            "layout_name": layout_name,
            "title": title,
            "backup": backup,
            "progress": progress,
            "token_estimate": len(str(progress)) // 4,
        }
    except Exception as e:
        progress.append(fail(str(e)))
        append_receipt(
            file_path,
            "add_slide",
            "pptx_basic",
            {"layout_name": layout_name, "title": title},
            f"✘ {e}",
            backup,
            False,
        )
        return _error(
            str(e),
            hint_for_error(e, path),
            progress,
            backup,
        )


def delete_slide(file_path: str, slide_index: int, open_after: bool = False) -> dict[str, Any]:
    """Remove a slide by index."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    path: Path | None = None
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        path = resolve_path(file_path)
        if not path.exists():
            return _not_found(path, progress)
        if path.suffix.lower() != ".pptx":
            return _wrong_type(path, ".pptx", progress)

        prs = Presentation(str(path))
        total = len(prs.slides)
        progress.append(ok(f"Opened {path.name}", f"{total} slides"))

        if slide_index < 0 or slide_index >= total:
            progress.append(fail(f"Slide index {slide_index} out of range"))
            return _error(
                f"slide_index {slide_index} out of range {index_range(total, 'slides')}",
                "Use read_presentation to get current slide count.",
                progress,
            )

        # Capture slide title for the receipt
        slide = prs.slides[slide_index]
        deleted_title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name.startswith("Title"):
                deleted_title = shape.text_frame.text  # type: ignore[reportAttributeAccessIssue]
                break

        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        # Remove the slide's entry from the XML slide ID list, and drop the
        # relationship that points at its part.
        #
        # Removing only the sldIdLst entry hides the slide from every reader --
        # PowerPoint and python-pptx both report the right count -- while the
        # slide part stays inside the .pptx with a live relationship naming it.
        # A sweep unzipped a deck after deleting a slide and found
        # ppt/slides/slide4.xml still shipping with all its text, and
        # _rels/presentation.xml.rels still carrying rId9 -> slides/slide4.xml.
        # Anyone handed the file could read what had been deleted from it,
        # which for a deck circulated after removing a slide is the entire
        # point of removing it. Dropping the relationship drops the part too:
        # the package serialiser writes only what is still reachable.
        sldIdLst = prs.slides._sldIdLst
        sld_id = sldIdLst[slide_index]
        prs.part.drop_rel(sld_id.rId)
        sldIdLst.remove(sld_id)

        prs.save(str(path))
        if open_after:
            open_file(path)
        reload_step = notify_reload(str(path), "pptx")
        progress.append(reload_step)
        progress.append(ok(f"Deleted slide {slide_index}", deleted_title[:40]))

        append_receipt(
            str(path),
            "delete_slide",
            "pptx_basic",
            {"slide_index": slide_index},
            f"✔ Deleted slide {slide_index} '{deleted_title}'",
            backup,
            True,
        )

        return {
            "success": True,
            "op": "delete_slide",
            "deleted_index": slide_index,
            "deleted_title": deleted_title,
            "remaining_slides": total - 1,
            "backup": backup,
            "progress": progress,
            "token_estimate": len(str(progress)) // 4,
        }
    except Exception as e:
        progress.append(fail(str(e)))
        append_receipt(
            file_path,
            "delete_slide",
            "pptx_basic",
            {"slide_index": slide_index},
            f"✘ {e}",
            backup,
            False,
        )
        return _error(
            str(e),
            hint_for_error(e, path),
            progress,
            backup,
        )


def reorder_slide(file_path: str, from_index: int, to_index: int, open_after: bool = False) -> dict[str, Any]:
    """Move a slide from one position to another."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    path: Path | None = None
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        path = resolve_path(file_path)
        if not path.exists():
            return _not_found(path, progress)
        if path.suffix.lower() != ".pptx":
            return _wrong_type(path, ".pptx", progress)

        prs = Presentation(str(path))
        total = len(prs.slides)
        progress.append(ok(f"Opened {path.name}", f"{total} slides"))

        for label, idx in (("from_index", from_index), ("to_index", to_index)):
            if idx < 0 or idx >= total:
                progress.append(fail(f"{label} {idx} out of range"))
                return _error(
                    f"{label} {idx} out of range {index_range(total, 'slides')}",
                    "Use read_presentation to get current slide count.",
                    progress,
                )

        if from_index == to_index:
            progress.append(info("No move needed — from_index equals to_index"))
            return {
                "success": True,
                "op": "reorder_slide",
                "from_index": from_index,
                "to_index": to_index,
                "message": "No change — indices are the same.",
                "progress": progress,
                "token_estimate": len(str(progress)) // 4,
            }

        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        sldIdLst = prs.slides._sldIdLst
        slide_id = sldIdLst[from_index]
        sldIdLst.remove(slide_id)
        sldIdLst.insert(to_index, slide_id)

        prs.save(str(path))
        if open_after:
            open_file(path)
        reload_step = notify_reload(str(path), "pptx")
        progress.append(reload_step)
        progress.append(ok(f"Moved slide {from_index} → {to_index}"))

        append_receipt(
            str(path),
            "reorder_slide",
            "pptx_basic",
            {"from_index": from_index, "to_index": to_index},
            f"✔ Moved slide {from_index} → {to_index}",
            backup,
            True,
        )

        return {
            "success": True,
            "op": "reorder_slide",
            "from_index": from_index,
            "to_index": to_index,
            "backup": backup,
            "progress": progress,
            "token_estimate": len(str(progress)) // 4,
        }
    except Exception as e:
        progress.append(fail(str(e)))
        append_receipt(
            file_path,
            "reorder_slide",
            "pptx_basic",
            {"from_index": from_index, "to_index": to_index},
            f"✘ {e}",
            backup,
            False,
        )
        return _error(
            str(e),
            hint_for_error(e, path),
            progress,
            backup,
        )


def add_text_box(
    file_path: str,
    slide_index: int,
    text: str,
    left: float = 1.0,
    top: float = 1.0,
    width: float = 5.0,
    height: float = 1.0,
    open_after: bool = False,
) -> dict[str, Any]:
    """Add a text box to a slide at the specified position (inches)."""
    progress: list[dict[str, Any]] = []
    backup: str | None = None
    path: Path | None = None
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.util import Inches  # type: ignore[import-untyped]

        path = resolve_path(file_path)
        if not path.exists():
            return _not_found(path, progress)
        if path.suffix.lower() != ".pptx":
            return _wrong_type(path, ".pptx", progress)

        prs = Presentation(str(path))
        total = len(prs.slides)
        progress.append(ok(f"Opened {path.name}", f"{total} slides"))

        if slide_index < 0 or slide_index >= total:
            progress.append(fail(f"Slide index {slide_index} out of range"))
            return _error(
                f"slide_index {slide_index} out of range {index_range(total, 'slides')}",
                "Use read_presentation to get current slide count.",
                progress,
            )

        backup = snapshot(str(path))
        progress.append(ok("Snapshot saved", Path(backup).name))

        slide = prs.slides[slide_index]
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.text = strip_list_markers(text)

        prs.save(str(path))
        if open_after:
            open_file(path)
        reload_step = notify_reload(str(path), "pptx")
        progress.append(reload_step)
        progress.append(
            ok(
                f"Added text box on slide {slide_index}",
                f"{text[:40]}",
            )
        )

        append_receipt(
            str(path),
            "add_text_box",
            "pptx_basic",
            {"slide_index": slide_index, "text": text, "left": left, "top": top},
            f"✔ Added text box on slide {slide_index}",
            backup,
            True,
        )

        return {
            "success": True,
            "op": "add_text_box",
            "slide_index": slide_index,
            "text": text,
            "position": {"left": left, "top": top, "width": width, "height": height},
            "backup": backup,
            "progress": progress,
            "token_estimate": len(str(progress)) // 4,
        }
    except Exception as e:
        progress.append(fail(str(e)))
        append_receipt(
            file_path,
            "add_text_box",
            "pptx_basic",
            {"slide_index": slide_index, "text": text},
            f"✘ {e}",
            backup,
            False,
        )
        return _error(
            str(e),
            hint_for_error(e, path),
            progress,
            backup,
        )


def diff_versions(file_path: str, timestamp_a: str, timestamp_b: str = "current") -> dict[str, Any]:
    """Compare two presentation versions by snapshot timestamp."""
    progress: list[dict[str, Any]] = []
    try:
        path = resolve_path(file_path)
        if not path.exists():
            return _not_found(path, progress)
        if path.suffix.lower() != ".pptx":
            return _wrong_type(path, ".pptx", progress)

        progress.append(ok(f"Opened {path.name}"))

        history = get_history(str(path))
        ts_map = {entry["timestamp"]: entry["backup_path"] for entry in history}

        if timestamp_a not in ts_map:
            available = [e["timestamp"] for e in history]
            progress.append(fail(f"Snapshot '{timestamp_a}' not found"))
            return _error(
                f"Snapshot '{timestamp_a}' not found",
                f"Available timestamps: {available}",
                progress,
            )

        path_a = ts_map[timestamp_a]

        if timestamp_b == "current":
            path_b = str(path)
        else:
            if timestamp_b not in ts_map:
                available = [e["timestamp"] for e in history]
                progress.append(fail(f"Snapshot '{timestamp_b}' not found"))
                return _error(
                    f"Snapshot '{timestamp_b}' not found",
                    f"Available timestamps: {available}",
                    progress,
                )
            path_b = ts_map[timestamp_b]

        result = diff_pptx(path_a, path_b)
        if not result.get("success"):
            progress.append(fail("Diff failed", result.get("error", "")))
            return _error(
                result.get("error", "Diff failed"),
                result.get("hint", "Check that both versions are valid .pptx files."),
                progress,
            )

        progress.append(
            ok(
                "Compared versions",
                f"{result['change_count']} change{'s' if result['change_count'] != 1 else ''}",
            )
        )

        return {
            "success": True,
            "op": "diff_versions",
            "timestamp_a": timestamp_a,
            "timestamp_b": timestamp_b,
            "summary": result.get("summary", ""),
            "slide_count_a": result["slide_count_a"],
            "slide_count_b": result["slide_count_b"],
            "slide_count_changed": result["slide_count_changed"],
            "text_changes": result["text_changes"],
            "change_count": result["change_count"],
            "progress": progress,
            "token_estimate": len(str(result)) // 4,
        }
    except Exception as e:
        progress.append(fail(str(e)))
        return _error(str(e), "Check that file_path points to a valid .pptx file.", progress)
