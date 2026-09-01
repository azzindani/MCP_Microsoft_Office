"""diff_pptx said `change_count: 0` about a deck that had just lost a slide.

The response contradicted itself in one object:

    "summary": "Slide count changed: 2 -> 1 (removed 1).",
    "slide_count_changed": true,
    "change_count": 0,
    progress: "Compared versions", detail "0 changes"

`change_count` was `len(text_changes)`, and text_changes only ever held
shape-text edits between *paired* slides. A caller checking the count first --
which is the whole point of a count -- concluded nothing had happened.

Two things were wrong, not one. The zip that pairs the slides stops at the
shorter deck, so the shapes on a dropped slide were never compared either: its
text vanished from the diff entirely rather than being reported as removed.

diff_docx beside it has always summed every kind of change it found
(`len(changes) + len(table_changes)`). This is that rule, applied to the sibling
that was missing it.
"""

from __future__ import annotations

from pptx import Presentation

from shared.doc_diff import diff_pptx


def _deck(path, titles):
    prs = Presentation()
    for title in titles:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
    prs.save(str(path))
    return str(path)


def test_a_removed_slide_is_counted(tmp_path):
    a = _deck(tmp_path / "a.pptx", ["one", "two"])
    b = _deck(tmp_path / "b.pptx", ["one"])

    result = diff_pptx(a, b)

    assert result["slide_count_changed"] is True
    assert result["change_count"] > 0, "a slide was deleted and the count said nothing changed"


def test_an_added_slide_is_counted(tmp_path):
    a = _deck(tmp_path / "a.pptx", ["one"])
    b = _deck(tmp_path / "b.pptx", ["one", "two"])

    assert diff_pptx(a, b)["change_count"] > 0


def test_the_count_never_disagrees_with_the_summary(tmp_path):
    # The exact contradiction from the sweep: these two fields are read
    # together and must not tell opposite stories.
    a = _deck(tmp_path / "a.pptx", ["one", "two"])
    b = _deck(tmp_path / "b.pptx", ["one"])

    result = diff_pptx(a, b)

    assert (result["summary"] == "No changes detected.") == (result["change_count"] == 0)


def test_the_text_on_a_dropped_slide_is_reported(tmp_path):
    # Beyond the shorter deck the zip never looked, so this text was lost.
    a = _deck(tmp_path / "a.pptx", ["kept", "GONE"])
    b = _deck(tmp_path / "b.pptx", ["kept"])

    result = diff_pptx(a, b)

    removed = [c for c in result["slide_changes"] if c["change"] == "removed"]
    assert removed, "the dropped slide is not in slide_changes"
    assert any("GONE" in text for c in removed for text in c["shape_texts"].values())


def test_the_text_on_an_added_slide_is_reported(tmp_path):
    a = _deck(tmp_path / "a.pptx", ["kept"])
    b = _deck(tmp_path / "b.pptx", ["kept", "NEW"])

    added = [c for c in diff_pptx(a, b)["slide_changes"] if c["change"] == "added"]
    assert any("NEW" in text for c in added for text in c["shape_texts"].values())


def test_equal_length_decks_still_report_no_slide_changes(tmp_path):
    a = _deck(tmp_path / "a.pptx", ["one", "two"])
    b = _deck(tmp_path / "b.pptx", ["one", "CHANGED"])

    result = diff_pptx(a, b)

    assert result["slide_changes"] == []
    assert result["change_count"] == len(result["text_changes"])


def test_two_identical_decks_still_count_zero(tmp_path):
    a = _deck(tmp_path / "a.pptx", ["same"])
    b = _deck(tmp_path / "b.pptx", ["same"])

    result = diff_pptx(a, b)

    assert result["change_count"] == 0
    assert result["summary"] == "No changes detected."
