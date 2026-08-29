"""diff_versions is the tool for confirming what changed. It could not see tables.

python-docx's `doc.paragraphs` excludes paragraphs inside table cells, and a
pptx table shape has no `text_frame`, so both diffs filtered tables out
completely. A document that gained a 2x3 table -- six cells of text -- reported
`change_count: 0` and "No changes detected."; so did a slide that gained a 2x2
table. Every table edit (set_cell, add_row, delete_row) was equally invisible.
"""

from __future__ import annotations

import docx
import pytest
from pptx import Presentation
from pptx.util import Inches

from shared.doc_diff import diff_docx, diff_pptx


def _docx(path, paragraphs=("one", "two"), table=None):
    d = docx.Document()
    for text in paragraphs:
        d.add_paragraph(text)
    if table is not None:
        t = d.add_table(rows=len(table), cols=len(table[0]))
        for r, row in enumerate(table):
            for c, val in enumerate(row):
                t.cell(r, c).text = val
    d.save(str(path))
    return str(path)


def _pptx(path, table=None):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if table is not None:
        shape = slide.shapes.add_table(len(table), len(table[0]), Inches(1), Inches(2), Inches(6), Inches(2))
        for r, row in enumerate(table):
            for c, val in enumerate(row):
                shape.table.cell(r, c).text = val
    prs.save(str(path))
    return str(path)


# --- docx -------------------------------------------------------------------


def test_a_document_that_gained_a_table_is_not_reported_unchanged(tmp_path):
    a = _docx(tmp_path / "a.docx")
    b = _docx(tmp_path / "b.docx", table=[["h1", "h2", "h3"], ["v1", "v2", "v3"]])

    result = diff_docx(a, b)

    assert result["success"] is True
    assert result["change_count"] > 0, "a whole table was added and the diff saw nothing"
    assert result["summary"] != "No changes detected."
    assert result["table_count_a"] == 0
    assert result["table_count_b"] == 1


def test_the_added_cells_are_reported(tmp_path):
    a = _docx(tmp_path / "a.docx")
    b = _docx(tmp_path / "b.docx", table=[["h1", "h2"], ["v1", "v2"]])

    cells = [c for change in diff_docx(a, b)["table_changes"] for row in change["b_cells"] for c in row]
    assert "h1" in cells and "v2" in cells


def test_a_changed_cell_is_reported(tmp_path):
    a = _docx(tmp_path / "a.docx", table=[["keep", "before"]])
    b = _docx(tmp_path / "b.docx", table=[["keep", "after"]])

    result = diff_docx(a, b)

    assert result["change_count"] > 0
    assert "table row" in result["summary"]
    flat = [c for change in result["table_changes"] for row in change["b_cells"] for c in row]
    assert "after" in flat


def test_a_deleted_row_is_reported(tmp_path):
    a = _docx(tmp_path / "a.docx", table=[["r1"], ["r2"]])
    b = _docx(tmp_path / "b.docx", table=[["r1"]])

    result = diff_docx(a, b)

    assert result["change_count"] > 0
    assert "deleted" in result["summary"]


def test_two_identical_documents_still_report_no_changes(tmp_path):
    a = _docx(tmp_path / "a.docx", table=[["same", "same"]])
    b = _docx(tmp_path / "b.docx", table=[["same", "same"]])

    result = diff_docx(a, b)

    assert result["change_count"] == 0
    assert result["summary"] == "No changes detected."


def test_paragraph_changes_are_still_reported_as_before(tmp_path):
    a = _docx(tmp_path / "a.docx", paragraphs=("one", "two"))
    b = _docx(tmp_path / "b.docx", paragraphs=("one", "CHANGED"))

    result = diff_docx(a, b)

    assert "paragraph" in result["summary"]
    assert any("CHANGED" in t for c in result["changes"] for t in c["b_text"])


# --- pptx -------------------------------------------------------------------


def test_a_slide_that_gained_a_table_is_not_reported_unchanged(tmp_path):
    a = _pptx(tmp_path / "a.pptx")
    b = _pptx(tmp_path / "b.pptx", table=[["k1", "k2"], ["w1", "w2"]])

    result = diff_pptx(a, b)

    assert result["success"] is True
    assert result["change_count"] > 0, "a table was added to the slide and the diff saw nothing"
    assert result["summary"] != "No changes detected."


def test_a_changed_slide_table_cell_is_reported(tmp_path):
    a = _pptx(tmp_path / "a.pptx", table=[["keep", "before"]])
    b = _pptx(tmp_path / "b.pptx", table=[["keep", "after"]])

    result = diff_pptx(a, b)

    assert result["change_count"] > 0
    assert any("after" in (c["new_text"] or "") for c in result["text_changes"])


def test_two_identical_decks_still_report_no_changes(tmp_path):
    a = _pptx(tmp_path / "a.pptx", table=[["same"]])
    b = _pptx(tmp_path / "b.pptx", table=[["same"]])

    assert diff_pptx(a, b)["summary"] == "No changes detected."


@pytest.mark.parametrize("fn", [diff_docx, diff_pptx])
def test_a_bad_path_still_fails_cleanly(fn, tmp_path):
    result = fn(str(tmp_path / "nope.bin"), str(tmp_path / "also_nope.bin"))
    assert result["success"] is False
    assert "hint" in result


def test_the_text_rendering_shows_table_rows_too(tmp_path):
    # format_diff_as_text rendered only `changes`, so a document whose only
    # change was a table produced a diff with a summary and no body.
    from shared.doc_diff import format_diff_as_text

    a = _docx(tmp_path / "a.docx")
    b = _docx(tmp_path / "b.docx", table=[["h1", "h2"]])

    text = format_diff_as_text(diff_docx(a, b))

    assert "h1" in text and "h2" in text


def test_the_summary_counts_rows_not_opcodes(tmp_path):
    # difflib returns one opcode per contiguous run, so counting opcodes
    # described a 2-row table insert as "1 table row added" -- a summary that
    # disagreed with the b_cells list printed directly beneath it.
    a = _docx(tmp_path / "a.docx")
    b = _docx(tmp_path / "b.docx", table=[["h1", "h2"], ["v1", "v2"], ["x1", "x2"]])

    result = diff_docx(a, b)

    assert "3 table rows added" in result["summary"]
    assert sum(len(c["b_cells"]) for c in result["table_changes"]) == 3


def test_the_paragraph_summary_counts_paragraphs_not_opcodes(tmp_path):
    a = _docx(tmp_path / "a.docx", paragraphs=("one", "two", "three", "four"))
    b = _docx(tmp_path / "b.docx", paragraphs=("one",))

    result = diff_docx(a, b)

    assert "3 paragraphs deleted" in result["summary"]
