"""Two pptx_new tools that reported one thing and built another.

**create_from_outline(layout=...)** recognised the exact token "title" and
nothing else. Everything else fell through to the Title-and-Content branch --
and was then echoed back verbatim in the progress line, so:

    layout="Title Only"  -> progress "title only", file uses Title and Content
    layout="zzz"         -> progress "zzz",        file uses Title and Content

The response asserted a layout the artifact contradicted, which is worse than
ignoring the argument quietly: it gave the caller a reason to believe it. Title
Only is a real layout in the default template (index 5), so the fix is to build
it, and to refuse a value that names no layout at all.

**create_from_docx** started a new slide only on Heading 1. create_from_sections
-- this fleet's own .docx builder -- writes the document title as Heading 1 and
every section heading as Heading 2, so a three-section document round-tripped
through the two tools came back as ONE slide with all three sections crammed
into its body. success: true, slide_count: 1, and max_slides=20 sitting in the
schema promising otherwise. Neither tool is wrong on its own; they simply never
had to agree, which is the shape technique 5 exists to catch.
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
for _p in (str(ROOT), str(ROOT / "servers" / "docx_new"), str(ROOT / "servers" / "pptx_new")):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from pptx import Presentation  # noqa: E402

from docx_new.engine import create_from_sections  # noqa: E402
from pptx_new.engine import _OUTLINE_LAYOUTS, create_from_docx, create_from_outline  # noqa: E402


def layouts_of(path: Path) -> list[str]:
    return [s.slide_layout.name for s in Presentation(str(path)).slides]


def slides_of(path: Path) -> list[list[str]]:
    return [
        [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame] for slide in Presentation(str(path)).slides
    ]


# --- the layout the caller asked for ----------------------------------------


@pytest.mark.parametrize(
    ("given", "expected"),
    [
        ("title", "Title Slide"),
        ("Title Slide", "Title Slide"),
        ("cover", "Title Slide"),
        ("title_only", "Title Only"),
        ("Title Only", "Title Only"),
        ("title-only", "Title Only"),
        ("content", "Title and Content"),
        ("Title and Content", "Title and Content"),
        ("bullets", "Title and Content"),
    ],
)
def test_the_file_uses_the_layout_that_was_asked_for(tmp_path, given, expected):
    out = tmp_path / "o.pptx"
    r = create_from_outline(str(out), [{"title": "T", "content": "body", "layout": given}], open_after=False)
    assert r["success"] is True, r.get("error")
    assert layouts_of(out) == [expected], f"{given!r} built the wrong layout"


def test_a_layout_that_names_nothing_is_refused(tmp_path):
    """It used to be accepted, downgraded, and reported back as itself."""
    out = tmp_path / "o.pptx"
    r = create_from_outline(str(out), [{"title": "T", "content": "b", "layout": "zzz"}], open_after=False)
    assert r["success"] is False
    assert "zzz" in r["error"]
    assert "title_only" in r["hint"]
    assert "Slide 1" in r["error"], "the refusal should say which slide"


def test_the_default_layout_is_unchanged(tmp_path):
    out = tmp_path / "o.pptx"
    r = create_from_outline(str(out), [{"title": "T", "content": "b"}], open_after=False)
    assert r["success"] is True
    assert layouts_of(out) == ["Title and Content"]


def test_title_only_really_drops_the_body(tmp_path):
    """The point of the layout: one placeholder, not two."""
    out = tmp_path / "o.pptx"
    create_from_outline(
        str(out), [{"title": "OnlyTitle", "content": "ignored", "layout": "title_only"}], open_after=False
    )
    text = " ".join(slides_of(out)[0])
    assert "OnlyTitle" in text
    assert "ignored" not in text


def test_every_accepted_spelling_maps_to_a_layout_that_exists():
    assert set(_OUTLINE_LAYOUTS.values()) == {"title", "title_only", "content"}


# --- the two tools compose --------------------------------------------------


def test_a_sectioned_document_becomes_one_slide_per_section(tmp_path):
    """The round trip that produced a single slide."""
    docx = tmp_path / "src.docx"
    create_from_sections(
        str(docx),
        "Deck Title",
        [{"heading": "Alpha", "body": "a"}, {"heading": "Beta", "body": "b"}, {"heading": "Gamma", "body": "c"}],
        open_after=False,
    )
    out = tmp_path / "d.pptx"
    r = create_from_docx(str(docx), str(out), open_after=False)
    assert r["success"] is True, r.get("error")

    slides = slides_of(out)
    assert len(slides) == 4, f"expected a title slide plus three sections, got {len(slides)}"
    titles = [s[0] for s in slides]
    assert titles == ["Deck Title", "Alpha", "Beta", "Gamma"]
    # And the bodies went with their own headings, not into one pile.
    assert "a" in " ".join(slides[1])
    assert "c" in " ".join(slides[3])


def test_repeated_heading_ones_still_split_on_heading_one(tmp_path):
    """The previous behaviour, for the documents that already worked."""
    from docx import Document

    docx = tmp_path / "h1.docx"
    doc = Document()
    for name in ("First", "Second", "Third"):
        doc.add_paragraph(name, style="Heading 1")
        doc.add_paragraph(f"body of {name}")
    doc.save(str(docx))

    out = tmp_path / "h1.pptx"
    r = create_from_docx(str(docx), str(out), open_after=False)
    assert r["success"] is True, r.get("error")
    assert [s[0] for s in slides_of(out)] == ["First", "Second", "Third"]


def test_max_slides_still_caps_the_deck(tmp_path):
    from docx import Document

    docx = tmp_path / "many.docx"
    doc = Document()
    for i in range(8):
        doc.add_paragraph(f"H{i}", style="Heading 1")
    doc.save(str(docx))

    out = tmp_path / "many.pptx"
    r = create_from_docx(str(docx), str(out), max_slides=3, open_after=False)
    assert r["success"] is True, r.get("error")
    assert len(slides_of(out)) == 3


def test_a_document_with_no_headings_is_untouched(tmp_path):
    """The no-headings path groups paragraphs in fives and must keep doing so."""
    from docx import Document

    docx = tmp_path / "flat.docx"
    doc = Document()
    for i in range(10):
        doc.add_paragraph(f"line {i}")
    doc.save(str(docx))

    out = tmp_path / "flat.pptx"
    r = create_from_docx(str(docx), str(out), open_after=False)
    assert r["success"] is True, r.get("error")
    assert len(slides_of(out)) == 2
