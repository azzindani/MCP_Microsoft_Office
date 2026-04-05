"""Tests for pptx_design engine functions."""

import shutil
from pathlib import Path

import pytest
from pptx import Presentation

from servers.pptx_design.engine import (
    add_chart,
    add_table,
    duplicate_slide,
    export_pdf,
    set_background,
    set_font_style,
)

FIXTURES = Path(__file__).parent / "fixtures"
DECK_SIMPLE = FIXTURES / "deck_simple.pptx"


@pytest.fixture()
def deck(tmp_path: Path) -> Path:
    """Return a writable copy of deck_simple.pptx."""
    dest = tmp_path / "deck_simple.pptx"
    shutil.copy(DECK_SIMPLE, dest)
    return dest


def _slide_count(path: Path) -> int:
    prs = Presentation(str(path))
    return len(prs.slides)


# ---------------------------------------------------------------------------
# set_background
# ---------------------------------------------------------------------------


def test_set_background_solid_color(deck: Path) -> None:
    result = set_background(str(deck), 0, color_hex="FF0000")
    assert result["success"] is True
    assert result["color_hex"] == "FF0000"
    assert result["slide_index"] == 0
    assert "backup" in result


def test_set_background_solid_color_with_hash(deck: Path) -> None:
    result = set_background(str(deck), 0, color_hex="#00FF00")
    assert result["success"] is True


def test_set_background_no_args_error(deck: Path) -> None:
    result = set_background(str(deck), 0)
    assert result["success"] is False
    assert "hint" in result


def test_set_background_invalid_slide_index(deck: Path) -> None:
    result = set_background(str(deck), 9999, color_hex="FF0000")
    assert result["success"] is False
    assert "hint" in result


def test_set_background_image(deck: Path, tmp_path: Path) -> None:
    """Create a tiny PNG and use it as background."""
    import struct
    import zlib

    def make_minimal_png() -> bytes:
        # Minimal 1x1 red PNG
        sig = b"\x89PNG\r\n\x1a\n"
        ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
        ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
        ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)
        raw = b"\x00\xFF\x00\x00"  # filter byte + RGB
        compressed = zlib.compress(raw)
        idat_crc = zlib.crc32(b"IDAT" + compressed) & 0xFFFFFFFF
        idat = struct.pack(">I", len(compressed)) + b"IDAT" + compressed + struct.pack(">I", idat_crc)
        iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
        iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)
        return sig + ihdr + idat + iend

    img_path = tmp_path / "bg.png"
    img_path.write_bytes(make_minimal_png())

    result = set_background(str(deck), 0, image_path=str(img_path))
    assert result["success"] is True


def test_set_background_image_not_found(deck: Path) -> None:
    result = set_background(str(deck), 0, image_path="/nonexistent/image.png")
    assert result["success"] is False
    assert "hint" in result


def test_set_background_creates_snapshot(deck: Path) -> None:
    result = set_background(str(deck), 0, color_hex="0000FF")
    assert result["success"] is True
    assert Path(result["backup"]).exists()


# ---------------------------------------------------------------------------
# set_font_style
# ---------------------------------------------------------------------------


def test_set_font_style(deck: Path) -> None:
    # Find the first slide's first text shape
    prs = Presentation(str(deck))
    slide = prs.slides[0]
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    if not text_shapes:
        pytest.skip("No text shapes on first slide of fixture")
    shape_name = text_shapes[0].name
    prs.save(str(deck))  # close without changes

    result = set_font_style(
        str(deck), 0, shape_name,
        font_name="Arial", font_size=18, bold=True
    )
    assert result["success"] is True
    assert result["shape_name"] == shape_name
    assert "backup" in result


def test_set_font_style_shape_not_found(deck: Path) -> None:
    result = set_font_style(str(deck), 0, "NonExistentShape123", font_name="Arial")
    assert result["success"] is False
    assert "hint" in result


def test_set_font_style_with_color(deck: Path) -> None:
    prs = Presentation(str(deck))
    slide = prs.slides[0]
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    if not text_shapes:
        pytest.skip("No text shapes on first slide")
    shape_name = text_shapes[0].name

    result = set_font_style(str(deck), 0, shape_name, color_hex="FF0000")
    assert result["success"] is True


def test_set_font_style_creates_snapshot(deck: Path) -> None:
    prs = Presentation(str(deck))
    slide = prs.slides[0]
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    if not text_shapes:
        pytest.skip("No text shapes on first slide")
    shape_name = text_shapes[0].name

    result = set_font_style(str(deck), 0, shape_name, bold=True)
    assert result["success"] is True
    assert Path(result["backup"]).exists()


# ---------------------------------------------------------------------------
# add_table
# ---------------------------------------------------------------------------


def test_add_table_on_slide(deck: Path) -> None:
    data = [
        ["Name", "Value"],
        ["Alpha", "100"],
        ["Beta", "200"],
    ]
    result = add_table(str(deck), 0, 3, 2, data)
    assert result["success"] is True
    assert result["rows"] == 3
    assert result["cols"] == 2
    assert "backup" in result


def test_add_table_empty_data(deck: Path) -> None:
    result = add_table(str(deck), 0, 2, 2, [])
    assert result["success"] is True  # empty data is fine, just blank cells


def test_add_table_invalid_rows(deck: Path) -> None:
    result = add_table(str(deck), 0, 0, 2, [])
    assert result["success"] is False
    assert "hint" in result


def test_add_table_invalid_slide(deck: Path) -> None:
    result = add_table(str(deck), 9999, 2, 2, [["a", "b"]])
    assert result["success"] is False
    assert "hint" in result


def test_add_table_creates_snapshot(deck: Path) -> None:
    result = add_table(str(deck), 0, 2, 2, [["A", "B"], ["C", "D"]])
    assert result["success"] is True
    assert Path(result["backup"]).exists()


# ---------------------------------------------------------------------------
# add_chart
# ---------------------------------------------------------------------------


def test_add_chart_bar(deck: Path) -> None:
    data = {
        "categories": ["Q1", "Q2", "Q3"],
        "series": [{"name": "Revenue", "values": [100, 200, 150]}],
    }
    result = add_chart(str(deck), 0, "bar", data, title="Revenue Chart")
    assert result["success"] is True
    assert result["chart_type"] == "bar"
    assert "backup" in result


def test_add_chart_line(deck: Path) -> None:
    data = {
        "categories": ["Jan", "Feb", "Mar"],
        "series": [{"name": "Sales", "values": [10, 20, 15]}],
    }
    result = add_chart(str(deck), 0, "line", data)
    assert result["success"] is True


def test_add_chart_pie(deck: Path) -> None:
    data = {
        "categories": ["A", "B", "C"],
        "series": [{"name": "Share", "values": [30, 50, 20]}],
    }
    result = add_chart(str(deck), 0, "pie", data)
    assert result["success"] is True


def test_add_chart_unsupported_type_error(deck: Path) -> None:
    data = {
        "categories": ["A", "B"],
        "series": [{"name": "S", "values": [1, 2]}],
    }
    result = add_chart(str(deck), 0, "donut", data)
    assert result["success"] is False
    assert "donut" in result["error"]
    assert "hint" in result


def test_add_chart_creates_snapshot(deck: Path) -> None:
    data = {
        "categories": ["X", "Y"],
        "series": [{"name": "N", "values": [5, 10]}],
    }
    result = add_chart(str(deck), 0, "bar", data)
    assert result["success"] is True
    assert Path(result["backup"]).exists()


# ---------------------------------------------------------------------------
# duplicate_slide
# ---------------------------------------------------------------------------


def test_duplicate_slide_increases_count(deck: Path) -> None:
    original_count = _slide_count(deck)
    result = duplicate_slide(str(deck), 0)
    assert result["success"] is True
    assert result["slide_count"] == original_count + 1
    assert "backup" in result

    new_count = _slide_count(deck)
    assert new_count == original_count + 1


def test_duplicate_slide_append_at_end(deck: Path) -> None:
    original_count = _slide_count(deck)
    result = duplicate_slide(str(deck), 0, insert_at=-1)
    assert result["success"] is True
    assert result["new_index"] == original_count  # appended at end


def test_duplicate_slide_insert_at_position(deck: Path) -> None:
    original_count = _slide_count(deck)
    if original_count < 2:
        pytest.skip("Need at least 2 slides to test insert_at positioning")

    result = duplicate_slide(str(deck), 0, insert_at=1)
    assert result["success"] is True
    new_count = _slide_count(deck)
    assert new_count == original_count + 1


def test_duplicate_slide_invalid_index(deck: Path) -> None:
    result = duplicate_slide(str(deck), 9999)
    assert result["success"] is False
    assert "hint" in result


def test_duplicate_slide_creates_snapshot(deck: Path) -> None:
    result = duplicate_slide(str(deck), 0)
    assert result["success"] is True
    assert Path(result["backup"]).exists()


# ---------------------------------------------------------------------------
# export_pdf
# ---------------------------------------------------------------------------


def test_export_pdf_no_converter(deck: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    """When no PDF converter is available, return a clear error."""
    from shared import platform_utils
    monkeypatch.setattr(platform_utils, "get_pdf_converter", lambda: None)
    # Re-import engine to pick up monkeypatched value
    from servers.pptx_design import engine as pptx_engine
    monkeypatch.setattr(pptx_engine, "get_pdf_converter", lambda: None)

    result = pptx_engine.export_pdf(str(deck))
    assert result["success"] is False
    assert "hint" in result
    assert "progress" in result


def test_export_pdf_file_not_found() -> None:
    result = export_pdf("/nonexistent/file.pptx")
    assert result["success"] is False
    assert "hint" in result


def test_export_pdf_wrong_type(deck: Path) -> None:
    # Rename to .docx to trigger wrong-type error
    wrong = deck.parent / "wrong.docx"
    shutil.copy(deck, wrong)
    result = export_pdf(str(wrong))
    assert result["success"] is False
    assert "hint" in result


# ---------------------------------------------------------------------------
# All write tools create snapshots and have progress
# ---------------------------------------------------------------------------


def test_all_write_tools_create_snapshot(deck: Path) -> None:
    prs = Presentation(str(deck))
    slide = prs.slides[0]
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    shape_name = text_shapes[0].name if text_shapes else None

    results = []

    # set_background
    r = set_background(str(deck), 0, color_hex="123456")
    results.append(("set_background", r))

    if shape_name:
        r = set_font_style(str(deck), 0, shape_name, bold=True)
        results.append(("set_font_style", r))

    r = add_table(str(deck), 0, 2, 2, [["A", "B"], ["C", "D"]])
    results.append(("add_table", r))

    data = {"categories": ["X"], "series": [{"name": "N", "values": [1]}]}
    r = add_chart(str(deck), 0, "bar", data)
    results.append(("add_chart", r))

    r = duplicate_slide(str(deck), 0)
    results.append(("duplicate_slide", r))

    for name, result in results:
        assert result["success"] is True, f"{name} failed: {result}"
        assert "backup" in result, f"{name} missing 'backup'"
        assert Path(result["backup"]).exists(), f"{name} backup file missing"


def test_all_responses_have_progress(deck: Path) -> None:
    data = {"categories": ["X"], "series": [{"name": "N", "values": [1]}]}
    results = [
        set_background(str(deck), 0, color_hex="FF0000"),
        set_background(str(deck), 9999, color_hex="FF0000"),
        set_background(str(deck), 0),
        add_table(str(deck), 0, 2, 2, [["A", "B"]]),
        add_chart(str(deck), 0, "bar", data),
        add_chart(str(deck), 0, "badtype", data),
        duplicate_slide(str(deck), 0),
        duplicate_slide(str(deck), 9999),
    ]
    for r in results:
        assert "progress" in r, f"Missing 'progress' key in: {r}"
        assert isinstance(r["progress"], list)
