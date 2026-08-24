"""Content given under a reasonable-but-undocumented key reaches the document.

`paragraphs`, `sections` and `data_slides` are bare list[dict] in their
schemas, so pydantic never sees the keys and nothing refuses a wrong one. The
entry was built from a .get() default instead, and the file came out short
while the response said it had been written:

    create_from_text(paragraphs=[{"content": "hello"}])
    -> success: true, "1 paragraphs written", a .docx with no text in it

    create_from_sections(sections=[{"header": "H", "body": "B"}])
    -> success: true, and a document whose headings are simply gone

    create_deck_from_data(data_slides=[{"title": "H", "items": ["a"]}])
    -> success: true, and a slide with a heading and an empty body

Only visible by opening the file -- every response was correct about what it
had been asked to do and silent about what it did.

pptx_new already aliased its heading and body keys and warned about an entry it
could not name; docx_new, written the same week, got neither. The vocabulary
now lives in shared/arg_alias.py so the next module cannot be the one that
misses out.
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
for _p in (str(ROOT), str(ROOT / "servers" / "docx_new"), str(ROOT / "servers" / "pptx_new")):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from docx import Document  # noqa: E402
from pptx import Presentation  # noqa: E402

from docx_new.engine import create_from_sections, create_from_text  # noqa: E402
from pptx_new.engine import create_deck_from_data, create_from_outline  # noqa: E402


def docx_text(path: Path) -> list[str]:
    return [p.text for p in Document(str(path)).paragraphs if p.text.strip()]


def pptx_text(path: Path) -> list[str]:
    out = []
    for slide in Presentation(str(path)).slides:
        out.append(" ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame))
    return out


# --- docx: the aliases now land in the file ---------------------------------


@pytest.mark.parametrize("key", ["text", "content", "body", "paragraph", "value"])
def test_a_paragraph_arrives_under_any_reasonable_key(tmp_path, key):
    out = tmp_path / "d.docx"
    r = create_from_text(str(out), [{key: "hello"}], open_after=False)
    assert r["success"] is True
    assert "hello" in docx_text(out), f"{key} was dropped"


@pytest.mark.parametrize("heading_key", ["heading", "title", "header", "name"])
def test_a_section_heading_survives_its_spelling(tmp_path, heading_key):
    out = tmp_path / "s.docx"
    r = create_from_sections(str(out), "T", [{heading_key: "H", "body": "B"}], open_after=False)
    assert r["success"] is True
    text = docx_text(out)
    assert "H" in text, f"{heading_key} was dropped: {text}"
    assert "B" in text


def test_paragraphs_carrying_nothing_are_refused_not_written_blank(tmp_path):
    out = tmp_path / "junk.docx"
    r = create_from_text(str(out), [{"zzz": "hello"}], open_after=False)
    assert r["success"] is False
    assert "zzz" in r["error"]
    assert "text" in r["error"]
    assert not out.exists()


def test_a_section_carrying_nothing_says_so(tmp_path):
    """Softer than paragraphs: a section may legitimately be heading-only or
    body-only, so this warns rather than refusing the whole document."""
    out = tmp_path / "s2.docx"
    r = create_from_sections(str(out), "T", [{"zzz": "H"}], open_after=False)
    assert r["success"] is True
    warnings = [p for p in r["progress"] if p.get("status") == "warn"]
    assert any("Section 0" in p["message"] for p in warnings), r["progress"]


# --- pptx: the bullet keys --------------------------------------------------


@pytest.mark.parametrize("key", ["bullets", "items", "points", "lines"])
def test_deck_bullets_arrive_under_any_reasonable_key(tmp_path, key):
    out = tmp_path / "deck.pptx"
    r = create_deck_from_data(str(out), "T", [{"heading": "H", key: ["alpha", "beta"]}], open_after=False)
    assert r["success"] is True
    body = " ".join(pptx_text(out))
    assert "alpha" in body and "beta" in body, f"{key} was dropped: {body!r}"


@pytest.mark.parametrize("key", ["title", "heading", "header", "name"])
def test_a_slide_heading_survives_its_spelling(tmp_path, key):
    out = tmp_path / "o.pptx"
    r = create_from_outline(str(out), [{key: "MyHeading", "content": "B"}], open_after=False)
    assert r["success"] is True
    assert "MyHeading" in " ".join(pptx_text(out))


# --- the documented spellings are untouched ---------------------------------


def test_the_documented_shapes_still_produce_what_they_always_did(tmp_path):
    d1 = tmp_path / "a.docx"
    create_from_text(str(d1), [{"text": "hello", "style": "Normal"}], open_after=False)
    assert docx_text(d1) == ["hello"]

    d2 = tmp_path / "b.docx"
    create_from_sections(str(d2), "T", [{"heading": "H", "body": "B"}], open_after=False)
    assert docx_text(d2) == ["T", "H", "B"]

    p1 = tmp_path / "c.pptx"
    create_deck_from_data(str(p1), "T", [{"heading": "H", "bullets": ["a"]}], open_after=False)
    joined = " ".join(pptx_text(p1))
    assert "H" in joined and "a" in joined
