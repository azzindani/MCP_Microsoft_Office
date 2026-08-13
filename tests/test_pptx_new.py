"""Tests for pptx_new engine functions."""

from pathlib import Path

from pptx import Presentation

from pptx_new.engine import (
    create_deck_from_data,
    create_from_docx,
    create_from_outline,
    create_from_template,
    create_presentation,
)

FIXTURES = Path(__file__).parent / "fixtures"
DECK_SIMPLE = FIXTURES / "deck_simple.pptx"
CONTRACT_SIMPLE = FIXTURES / "contract_simple.docx"


# ---------------------------------------------------------------------------
# create_presentation
# ---------------------------------------------------------------------------


def test_create_presentation_with_title_and_subtitle(tmp_path: Path) -> None:
    out = tmp_path / "deck.pptx"
    result = create_presentation(str(out), title="Q4 Review", subtitle="Finance Team", open_after=False)
    assert result["success"] is True
    assert result["slide_count"] == 1

    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_create_presentation_blank(tmp_path: Path) -> None:
    out = tmp_path / "deck.pptx"
    result = create_presentation(str(out), open_after=False)
    assert result["success"] is True
    assert out.exists()


# ---------------------------------------------------------------------------
# create_from_outline
# ---------------------------------------------------------------------------


def test_create_from_outline_builds_all_slides(tmp_path: Path) -> None:
    out = tmp_path / "deck.pptx"
    slides = [
        {"title": "Welcome", "layout": "title", "content": "Intro subtitle"},
        {"title": "Agenda", "content": "Item 1\nItem 2"},
    ]
    result = create_from_outline(str(out), slides, open_after=False)
    assert result["success"] is True
    assert result["slide_count"] == 2

    prs = Presentation(str(out))
    assert len(prs.slides) == 2


def test_create_from_outline_rejects_empty_list(tmp_path: Path) -> None:
    out = tmp_path / "deck.pptx"
    result = create_from_outline(str(out), [], open_after=False)
    assert result["success"] is False
    assert "empty" in result["error"].lower()


# ---------------------------------------------------------------------------
# create_deck_from_data
# ---------------------------------------------------------------------------


def test_create_deck_from_data_adds_title_plus_one_slide_per_item(tmp_path: Path) -> None:
    out = tmp_path / "deck.pptx"
    data_slides = [
        {"heading": "Revenue", "bullets": ["Up 18% YoY", "All regions growing"]},
        {"heading": "Costs", "bullets": ["Flat QoQ"]},
    ]
    result = create_deck_from_data(str(out), "Q3 Review", data_slides, open_after=False)
    assert result["success"] is True
    assert result["slide_count"] == 3  # 1 title + 2 data slides

    prs = Presentation(str(out))
    assert len(prs.slides) == 3


def test_create_deck_from_data_rejects_empty_list(tmp_path: Path) -> None:
    out = tmp_path / "deck.pptx"
    result = create_deck_from_data(str(out), "Title", [], open_after=False)
    assert result["success"] is False
    assert "empty" in result["error"].lower()


# ---------------------------------------------------------------------------
# create_from_template
# ---------------------------------------------------------------------------


def test_create_from_template_copies_deck(tmp_path: Path) -> None:
    out = tmp_path / "copy.pptx"
    result = create_from_template(str(DECK_SIMPLE), str(out), open_after=False)
    assert result["success"] is True

    src = Presentation(str(DECK_SIMPLE))
    dst = Presentation(str(out))
    assert len(dst.slides) == len(src.slides)


def test_create_from_template_missing_template(tmp_path: Path) -> None:
    out = tmp_path / "copy.pptx"
    result = create_from_template(str(tmp_path / "ghost.pptx"), str(out), open_after=False)
    assert result["success"] is False
    assert "not found" in result["error"].lower()


def test_create_from_template_wrong_file_type(tmp_path: Path) -> None:
    bad = tmp_path / "notes.txt"
    bad.write_text("hello")
    out = tmp_path / "copy.pptx"
    result = create_from_template(str(bad), str(out), open_after=False)
    assert result["success"] is False


# ---------------------------------------------------------------------------
# create_from_docx
# ---------------------------------------------------------------------------


def test_create_from_docx_builds_slides_from_headings(tmp_path: Path) -> None:
    out = tmp_path / "deck.pptx"
    result = create_from_docx(str(CONTRACT_SIMPLE), str(out), open_after=False)
    assert result["success"] is True
    assert result["slide_count"] >= 1

    prs = Presentation(str(out))
    assert len(prs.slides) == result["slide_count"]


def test_create_from_docx_missing_file(tmp_path: Path) -> None:
    out = tmp_path / "deck.pptx"
    result = create_from_docx(str(tmp_path / "ghost.docx"), str(out), open_after=False)
    assert result["success"] is False
    assert "not found" in result["error"].lower()


def test_create_from_docx_wrong_file_type(tmp_path: Path) -> None:
    bad = tmp_path / "notes.txt"
    bad.write_text("hello")
    out = tmp_path / "deck.pptx"
    result = create_from_docx(str(bad), str(out), open_after=False)
    assert result["success"] is False


def test_create_from_docx_respects_max_slides(tmp_path: Path) -> None:
    out = tmp_path / "deck.pptx"
    result = create_from_docx(str(CONTRACT_SIMPLE), str(out), max_slides=0, open_after=False)
    assert result["success"] is True
    assert result["slide_count"] == 0


# ---------------------------------------------------------------------------
# Global contract: every tool response has a progress field
# ---------------------------------------------------------------------------


def test_all_responses_have_progress_field(tmp_path: Path) -> None:
    results = [
        create_presentation(str(tmp_path / "a.pptx"), title="T", open_after=False),
        create_from_outline(str(tmp_path / "b.pptx"), [{"title": "S1"}], open_after=False),
        create_deck_from_data(str(tmp_path / "c.pptx"), "T", [{"heading": "H", "bullets": ["x"]}], open_after=False),
        create_from_template(str(DECK_SIMPLE), str(tmp_path / "d.pptx"), open_after=False),
        create_from_docx(str(CONTRACT_SIMPLE), str(tmp_path / "e.pptx"), open_after=False),
    ]
    for r in results:
        assert "progress" in r, f"Missing 'progress' in: {r}"
        assert isinstance(r["progress"], list)
