"""duplicate_slide copied the shapes and left the background behind.

A coverage sweep set every slide of a deck to navy, duplicated one, then applied
white text to all slides -- and the copy came out white on white. The model
noticed and re-applied the background by hand; nothing in the response had said
the background was gone, because nothing knew it had been dropped.

A slide's background lives in <p:bg> on its own <p:cSld>, not in the shape tree.
duplicate_slide builds the new slide from the source's layout and then moves the
shapes across, so anything outside the shape tree is simply not carried: the
same is true of speaker notes, which live in a separate part reached by a
relationship, and which the duplicate also lost.

Both now travel with the copy. Notes are copied as text -- cloning the whole
notesSlide part to preserve run formatting would mean rebuilding the package
relationship by hand for a gain nobody has asked for.
"""

from __future__ import annotations

import struct
import zlib
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from pptx_design.engine import (  # type: ignore[reportMissingImports]
    add_chart,
    add_image_to_all_slides,
    duplicate_slide,
    set_background,
)


def _png_chunk(tag: bytes, payload: bytes) -> bytes:
    return struct.pack(">I", len(payload)) + tag + payload + struct.pack(">I", zlib.crc32(tag + payload) & 0xFFFFFFFF)


NAVY = "1F3864"
TITLE = "Where the money went"
NOTES = "Say that Google Ads carried 77% of the spend."


def background_colors(path: str, index: int) -> list[str]:
    slide = Presentation(path).slides[index]
    cSld = slide._element.find(qn("p:cSld"))
    assert cSld is not None
    bg = cSld.find(qn("p:bg"))
    if bg is None:
        return []
    return [c.get("val") for c in bg.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")]


def notes_of(path: str, index: int) -> str | None:
    slide = Presentation(path).slides[index]
    if not slide.has_notes_slide:
        return None
    return slide.notes_slide.notes_text_frame.text  # type: ignore[reportOptionalMemberAccess]


@pytest.fixture()
def deck(tmp_path: Path) -> str:
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = TITLE  # type: ignore[reportOptionalMemberAccess]
    slide.placeholders[1].text = "Google Ads 1,939,000"  # type: ignore[reportAttributeAccessIssue]
    slide.notes_slide.notes_text_frame.text = NOTES  # type: ignore[reportOptionalMemberAccess]
    prs.save(str(path))
    return str(path)


@pytest.fixture()
def navy_deck(deck: str) -> str:
    assert set_background(deck, -1, color_hex=NAVY)["success"] is True
    return deck


class TestTheBackgroundTravelsWithTheCopy:
    def test_the_source_still_has_it(self, navy_deck: str):
        duplicate_slide(navy_deck, 0)
        assert background_colors(navy_deck, 0) == [NAVY]

    def test_the_copy_has_it_too(self, navy_deck: str):
        r = duplicate_slide(navy_deck, 0)
        assert r["success"] is True, r.get("error")
        assert background_colors(navy_deck, 1) == [NAVY], background_colors(navy_deck, 1)

    def test_it_survives_being_inserted_earlier(self, navy_deck: str):
        prs = Presentation(navy_deck)
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(navy_deck)
        set_background(navy_deck, -1, color_hex=NAVY)
        duplicate_slide(navy_deck, 0, insert_at=1)
        assert background_colors(navy_deck, 1) == [NAVY]


class TestTheNotesTravelWithTheCopy:
    def test_the_copy_has_notes(self, deck: str):
        duplicate_slide(deck, 0)
        assert notes_of(deck, 1) == NOTES, notes_of(deck, 1)

    def test_the_source_keeps_its_own(self, deck: str):
        duplicate_slide(deck, 0)
        assert notes_of(deck, 0) == NOTES

    def test_a_slide_with_no_notes_does_not_gain_any(self, tmp_path: Path):
        path = tmp_path / "plain.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(path))
        assert duplicate_slide(str(path), 0)["success"] is True
        assert notes_of(str(path), 1) in (None, "")


class TestAChartTravelsWithTheCopy:
    """The copy claimed to hold a chart and could not draw one.

    A chart is not self-contained XML: the <p:graphicFrame> holds an r:id into a
    relationship on the slide part. Deep-copying the element copied the id and
    left the relationship behind, so read_slide reported a chart -- has_chart
    was True -- while reading it raised "no relationship with key 'rId2'" and
    LibreOffice rendered an empty space. A sweep duplicated a slide holding a
    chart and a table and got back the table alone.
    """

    @pytest.fixture()
    def charted(self, tmp_path: Path) -> str:
        path = tmp_path / "chart.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])
        prs.save(str(path))
        r = add_chart(
            str(path),
            0,
            "bar",
            {"categories": ["Google Ads", "Facebook Ads"], "series": {"Spend": [1939000, 564100]}},
            title="Spend by platform",
        )
        assert r["success"] is True, r.get("error")
        return str(path)

    def test_the_copy_still_has_a_chart_shape(self, charted: str):
        assert duplicate_slide(charted, 0)["success"] is True
        copy_shapes = Presentation(charted).slides[1].shapes
        assert any(s.has_chart for s in copy_shapes)

    def test_the_chart_can_actually_be_read(self, charted: str):
        duplicate_slide(charted, 0)
        chart = next(s.chart for s in Presentation(charted).slides[1].shapes if s.has_chart)  # type: ignore[reportAttributeAccessIssue]
        assert [s.name for s in chart.plots[0].series] == ["Spend"]

    def test_it_carries_the_same_numbers(self, charted: str):
        duplicate_slide(charted, 0)
        prs = Presentation(charted)
        got = []
        for index in (0, 1):
            chart = next(s.chart for s in prs.slides[index].shapes if s.has_chart)  # type: ignore[reportAttributeAccessIssue]
            got.append(list(chart.plots[0].series[0].values))
        assert got[0] == got[1] == [1939000, 564100], got


class TestAnImageTravelsWithTheCopy:
    @pytest.fixture()
    def with_image(self, deck: str, tmp_path: Path) -> str:
        png = tmp_path / "logo.png"
        png.write_bytes(
            b"\x89PNG\r\n\x1a\n"
            + _png_chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
            + _png_chunk(b"IDAT", zlib.compress(b"\x00\xff\xff\xff"))
            + _png_chunk(b"IEND", b"")
        )
        assert add_image_to_all_slides(deck, str(png))["success"] is True
        return deck

    def test_the_copy_still_shows_it(self, with_image: str):
        duplicate_slide(with_image, 0)
        pictures = [s for s in Presentation(with_image).slides[1].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert pictures, [s.name for s in Presentation(with_image).slides[1].shapes]

    def test_its_bytes_are_reachable(self, with_image: str):
        duplicate_slide(with_image, 0)
        picture = next(s for s in Presentation(with_image).slides[1].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
        assert len(picture.image.blob) > 0  # type: ignore[reportAttributeAccessIssue]


class TestWhatWasAlreadyRightStaysRight:
    def test_the_shapes_are_still_copied_in_order(self, navy_deck: str):
        duplicate_slide(navy_deck, 0)
        prs = Presentation(navy_deck)
        assert [s.name for s in prs.slides[0].shapes] == [s.name for s in prs.slides[1].shapes]

    def test_no_empty_placeholder_is_left_behind(self, navy_deck: str):
        duplicate_slide(navy_deck, 0)
        copy_shapes = Presentation(navy_deck).slides[1].shapes
        empty = [
            s.name
            for s in copy_shapes
            if s.has_text_frame and not s.text_frame.text.strip()  # type: ignore[reportAttributeAccessIssue]
        ]
        assert not empty, empty

    def test_the_title_is_still_there(self, navy_deck: str):
        duplicate_slide(navy_deck, 0)
        title = Presentation(navy_deck).slides[1].shapes.title
        assert title is not None and title.text == TITLE

    def test_a_deck_with_no_background_set_is_unaffected(self, deck: str):
        duplicate_slide(deck, 0)
        assert background_colors(deck, 1) == []
