"""A template filled with two related keys came out corrupted, and said success.

The substitution loop replaced each key in turn with a plain substring replace,
in whatever order the caller's dict happened to be in. Two keys where one
contains the other destroyed the document:

    {"platform": "Google Ads", "platform_spend": "1,939,003.26", "total": "..."}

    Platform: {platform}        ->  Platform: {Google Ads}
    Spend: {platform_spend}     ->  Spend: {Google Ads_spend}
    Total: {total}              ->  Total: {2,503,118.77}

    success: True, substitutions_applied: 3
    warn: "Placeholder 'platform_spend' not found in template"

Three failures in one call. `platform` matched inside `platform_spend` and ate
it; the warning then blamed the template for a placeholder the tool had just
destroyed; and the braces survived every replacement, because the caller passed
bare names and nothing anywhere says whether a key includes its delimiters --
the docstring is 80 characters and `substitutions` is an opaque dict in the
schema.

docx and pptx had the same sequential loop. Both now resolve what each key
should match before touching anything, then replace all targets at once, so a
key cannot consume another key nor a value an earlier key produced.

Found by reading the notes column of a round-7 sweep report that said PASS.
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).parent.parent / "servers" / "docx_new"))
sys.path.insert(0, str(Path(__file__).parent.parent / "servers" / "pptx_new"))

from docx_new.engine import create_from_template as docx_from_template  # noqa: E402
from pptx_new.engine import create_from_template as pptx_from_template  # noqa: E402
from shared.template_fill import ordered_pairs, resolve_targets, substitute_once  # noqa: E402

# The sweep's own keys: one is a prefix of the other.
COLLIDING = {
    "platform": "Google Ads",
    "platform_spend": "1,939,003.26",
    "total": "2,503,118.77",
}
BRACED = {f"{{{k}}}": v for k, v in COLLIDING.items()}
FILLED = ["Platform: Google Ads", "Spend: 1,939,003.26", "Total: 2,503,118.77"]


@pytest.fixture()
def docx_template(tmp_path: Path) -> str:
    d = Document()
    d.add_paragraph("Platform: {platform}")
    d.add_paragraph("Spend: {platform_spend}")
    d.add_paragraph("Total: {total}")
    d.save(str(tmp_path / "tpl.docx"))
    return str(tmp_path / "tpl.docx")


@pytest.fixture()
def pptx_template(tmp_path: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    assert title is not None
    title.text = "Platform: {platform}"
    frame = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2)).text_frame
    frame.text = "Spend: {platform_spend}"
    frame.add_paragraph().text = "Total: {total}"
    prs.save(str(tmp_path / "tpl.pptx"))
    return str(tmp_path / "tpl.pptx")


def docx_lines(path: str) -> list[str]:
    return [p.text for p in Document(path).paragraphs if p.text.strip()]


def pptx_lines(path: str) -> list[str]:
    out = []
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            frame = getattr(shape, "text_frame", None)
            if frame is None:
                continue
            out.extend(p.text for p in frame.paragraphs if p.text.strip())
    return out


class TestTheSweepsCallOnDocx:
    def test_the_longer_key_is_not_eaten_by_the_shorter(self, docx_template: str, tmp_path: Path):
        out = str(tmp_path / "f.docx")
        docx_from_template(docx_template, out, COLLIDING, open_after=False)
        assert "Google Ads_spend" not in "\n".join(docx_lines(out))

    def test_every_placeholder_is_filled(self, docx_template: str, tmp_path: Path):
        out = str(tmp_path / "f.docx")
        docx_from_template(docx_template, out, COLLIDING, open_after=False)
        assert docx_lines(out) == FILLED

    def test_no_braces_are_left_behind(self, docx_template: str, tmp_path: Path):
        out = str(tmp_path / "f.docx")
        docx_from_template(docx_template, out, COLLIDING, open_after=False)
        body = "\n".join(docx_lines(out))
        assert "{" not in body and "}" not in body, body

    def test_nothing_is_reported_as_missing(self, docx_template: str, tmp_path: Path):
        r = docx_from_template(docx_template, str(tmp_path / "f.docx"), COLLIDING, open_after=False)
        missing = [p["msg"] for p in r["progress"] if "not found" in p.get("msg", "")]
        assert not missing, missing

    def test_the_delimited_match_is_explained(self, docx_template: str, tmp_path: Path):
        r = docx_from_template(docx_template, str(tmp_path / "f.docx"), COLLIDING, open_after=False)
        notes = " ".join(p.get("msg", "") for p in r["progress"])
        assert "{platform}" in notes, notes

    def test_keys_written_with_their_braces_still_work(self, docx_template: str, tmp_path: Path):
        out = str(tmp_path / "f.docx")
        docx_from_template(docx_template, out, BRACED, open_after=False)
        assert docx_lines(out) == FILLED


class TestTheSameCallOnPptx:
    def test_the_longer_key_is_not_eaten_by_the_shorter(self, pptx_template: str, tmp_path: Path):
        out = str(tmp_path / "f.pptx")
        pptx_from_template(pptx_template, out, COLLIDING, open_after=False)
        assert "Google Ads_spend" not in "\n".join(pptx_lines(out))

    def test_every_placeholder_is_filled(self, pptx_template: str, tmp_path: Path):
        out = str(tmp_path / "f.pptx")
        pptx_from_template(pptx_template, out, COLLIDING, open_after=False)
        assert sorted(pptx_lines(out)) == sorted(FILLED)

    def test_keys_written_with_their_braces_still_work(self, pptx_template: str, tmp_path: Path):
        out = str(tmp_path / "f.pptx")
        pptx_from_template(pptx_template, out, BRACED, open_after=False)
        assert sorted(pptx_lines(out)) == sorted(FILLED)


class TestAValueIsNeverReplacedAgain:
    """The second half of the ordering bug: a value can contain a later key."""

    def test_docx_leaves_an_inserted_value_alone(self, docx_template: str, tmp_path: Path):
        out = str(tmp_path / "f.docx")
        docx_from_template(
            docx_template,
            out,
            {"total": "platform total", "platform": "Google Ads"},
            open_after=False,
        )
        assert "Total: platform total" in docx_lines(out)

    def test_substitute_once_does_not_rewrite_its_own_output(self):
        pairs = ordered_pairs({"a": "a", "b": "b"}, {"a": "b", "b": "z"})
        assert substitute_once("a b", pairs)[0] == "b z"


class TestAMissingKeyIsStillReported:
    def test_docx_warns_and_changes_nothing(self, docx_template: str, tmp_path: Path):
        out = str(tmp_path / "f.docx")
        r = docx_from_template(docx_template, out, {"nope": "x"}, open_after=False)
        assert r["success"] is True
        assert r["substitutions_applied"] == 0
        assert any("nope" in p.get("msg", "") for p in r["progress"])
        assert docx_lines(out) == ["Platform: {platform}", "Spend: {platform_spend}", "Total: {total}"]

    def test_pptx_warns_too(self, pptx_template: str, tmp_path: Path):
        r = pptx_from_template(pptx_template, str(tmp_path / "f.pptx"), {"nope": "x"}, open_after=False)
        assert any("nope" in p.get("msg", "") for p in r["progress"])


class TestTheResolver:
    def test_a_delimited_form_beats_the_bare_literal(self):
        targets, _ = resolve_targets("Platform: {platform}", {"platform": "x"})
        assert targets["platform"] == "{platform}"

    def test_an_undelimited_template_still_matches_the_bare_key(self):
        targets, _ = resolve_targets("Platform: platform", {"platform": "x"})
        assert targets["platform"] == "platform"

    def test_a_key_present_in_no_form_is_absent(self):
        targets, _ = resolve_targets("nothing here", {"platform": "x"})
        assert "platform" not in targets

    def test_double_braces_are_preferred_over_single(self):
        targets, _ = resolve_targets("{{name}} and {name}", {"name": "x"})
        assert targets["name"] == "{{name}}"

    def test_targets_come_back_longest_first(self):
        targets, _ = resolve_targets("a ab abc", {"a": "1", "ab": "2", "abc": "3"})
        lengths = [len(t) for t, _ in ordered_pairs(targets, {"a": "1", "ab": "2", "abc": "3"})]
        assert lengths == sorted(lengths, reverse=True)

    def test_an_empty_key_is_ignored(self):
        targets, _ = resolve_targets("anything", {"": "x"})
        assert targets == {}
