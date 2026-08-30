"""bold and italic could be switched ON and never OFF, across three servers.

Found by round 22's sweep, whose axis was to cross-check a tool's answer rather
than trust its success flag. `set_font_all_slides(bold=false)` returned:

    success: true, shapes_modified: 6, bold: false

and the title was still bold on read-back. The success flag, the count and the
echoed argument all agreed with each other and disagreed with the file, which is
the hardest defect shape to see and the reason the sweep recomputes.

The cause is one line, repeated six times in three servers:

    if bold:
        run.font.bold = True

`False` is a real value, not an absence, so truthiness cannot tell "make this
not bold" from "bold was not mentioned". The neighbouring arguments escape it
because their types carry a value outside the useful range -- `font_name=""` and
`font_size=0` are genuinely unset. A bool has none.

Fixed by spelling the three states: "" leaves, "true" turns on, "false" turns
off (shared/tristate.py). That is a deliberate schema change -- a caller passing
a JSON boolean now gets a validation error instead of an argument that was
accepted, ignored, and reported as applied.

`shapes_modified` was a second, quieter fault at the same site: it counted every
shape with a text frame, so a call that changed nothing still reported six.
"""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest

FIXTURES = Path(__file__).parent / "fixtures"


def _copy(name: str, tmp_path: Path) -> Path:
    dest = tmp_path / name
    shutil.copy(FIXTURES / name, dest)
    return dest


# ---------------------------------------------------------------------------
# pptx_design — where the sweep found it
# ---------------------------------------------------------------------------


@pytest.fixture()
def deck(tmp_path: Path) -> Path:
    return _copy("deck_simple.pptx", tmp_path)


def _first_run_bold(path: Path, slide_index: int = 0):
    from pptx import Presentation

    shape = Presentation(str(path)).slides[slide_index].shapes[0]
    return shape.text_frame.paragraphs[0].runs[0].font.bold


def _shape_name(path: Path, slide_index: int = 0) -> str:
    from pptx import Presentation

    return Presentation(str(path)).slides[slide_index].shapes[0].name


class TestSetFontStyle:
    def test_bold_true_turns_it_on(self, deck):
        from pptx_design.engine import set_font_style

        assert set_font_style(str(deck), 0, _shape_name(deck), bold="true")["success"]
        assert _first_run_bold(deck) is True

    def test_bold_false_turns_it_off(self, deck):
        """The operation that was impossible through any argument."""
        from pptx_design.engine import set_font_style

        name = _shape_name(deck)
        set_font_style(str(deck), 0, name, bold="true")
        assert _first_run_bold(deck) is True

        result = set_font_style(str(deck), 0, name, bold="false")
        assert result["success"], result
        assert _first_run_bold(deck) is False, "bold=false reported success and changed nothing"

    def test_omitting_bold_leaves_it_alone(self, deck):
        """The reason this is three states and not two.

        Changing only the font name must not silently un-bold a title -- which
        is what simply applying the boolean would have done.
        """
        from pptx_design.engine import set_font_style

        name = _shape_name(deck)
        set_font_style(str(deck), 0, name, bold="true")
        set_font_style(str(deck), 0, name, font_name="Georgia")
        assert _first_run_bold(deck) is True

    def test_the_response_reports_the_effect_not_the_argument(self, deck):
        from pptx_design.engine import set_font_style

        name = _shape_name(deck)
        assert set_font_style(str(deck), 0, name, bold="false")["bold"] == "false"
        assert set_font_style(str(deck), 0, name, font_name="Georgia")["bold"] == "unchanged"

    def test_a_value_that_is_neither_is_refused_with_the_three_words(self, deck):
        from pptx_design.engine import set_font_style

        result = set_font_style(str(deck), 0, _shape_name(deck), bold="maybe")
        assert result["success"] is False
        assert "'true'" in result["hint"] and "'false'" in result["hint"]


class TestSetFontAllSlides:
    def test_bold_false_turns_it_off_everywhere(self, deck):
        from pptx_design.engine import set_font_all_slides

        set_font_all_slides(str(deck), bold="true")
        assert _first_run_bold(deck) is True

        result = set_font_all_slides(str(deck), bold="false")
        assert result["success"], result
        assert result["bold"] == "false"
        assert _first_run_bold(deck) is False

    def test_shapes_modified_counts_shapes_that_changed(self, deck):
        """It counted every shape with a text frame, changed or not.

        A call naming no font argument at all changes nothing, and must not
        report a count that agrees with the success flag and disagrees with the
        file.
        """
        from pptx_design.engine import set_font_all_slides

        result = set_font_all_slides(str(deck))
        assert result["success"], result
        assert result["shapes_modified"] == 0
        assert result["slides_modified"] == 0

    def test_a_real_change_still_counts(self, deck):
        from pptx_design.engine import set_font_all_slides

        result = set_font_all_slides(str(deck), font_name="Georgia")
        assert result["shapes_modified"] > 0


# ---------------------------------------------------------------------------
# docx_layout — bold AND italic, same line, same fault
# ---------------------------------------------------------------------------


class TestDocxSetFont:
    def _bold_italic(self, path: Path):
        from docx import Document

        run = Document(str(path)).paragraphs[0].runs[0]
        return run.bold, run.italic

    def test_bold_and_italic_can_both_be_turned_off(self, tmp_path):
        from docx_layout.engine import set_font

        path = _copy("contract_simple.docx", tmp_path)
        set_font(str(path), 0, bold="true", italic="true")
        assert self._bold_italic(path) == (True, True)

        result = set_font(str(path), 0, bold="false", italic="false")
        assert result["success"], result
        assert self._bold_italic(path) == (False, False)

    def test_omitting_them_leaves_them_alone(self, tmp_path):
        from docx_layout.engine import set_font

        path = _copy("contract_simple.docx", tmp_path)
        set_font(str(path), 0, bold="true", italic="true")
        set_font(str(path), 0, font_size=14)
        assert self._bold_italic(path) == (True, True)

    def test_changes_names_the_state_it_set(self, tmp_path):
        """`changes` said "bold=True" or nothing. It can now say False."""
        from docx_layout.engine import set_font

        path = _copy("contract_simple.docx", tmp_path)
        result = set_font(str(path), 0, bold="false")
        assert "bold=False" in result["changes"]


# ---------------------------------------------------------------------------
# xlsx_charts — the same fault twice in one expression
# ---------------------------------------------------------------------------


class TestSetCellStyle:
    def test_bold_false_reaches_the_font(self, tmp_path):
        """Two guards had to be passed, and False failed both.

            "bold": bold if bold else (existing.bold if existing else None)
            Font(**{k: v for k, v in kwargs.items() if v is not None})

        The first fell through to the existing value; even past it, False was
        a legitimate value the second would have kept, so only the first
        needed fixing -- but both had to be read to know that.
        """
        import openpyxl

        from xlsx_charts.engine import set_cell_style

        path = _copy("dashboard.xlsx", tmp_path)
        sheet = openpyxl.load_workbook(str(path)).sheetnames[0]

        set_cell_style(str(path), sheet, "A1", bold="true")
        assert openpyxl.load_workbook(str(path))[sheet]["A1"].font.bold is True

        result = set_cell_style(str(path), sheet, "A1", bold="false")
        assert result["success"], result
        assert result["bold"] == "false"
        assert openpyxl.load_workbook(str(path))[sheet]["A1"].font.bold is False

    def test_omitting_bold_keeps_what_the_cell_had(self, tmp_path):
        import openpyxl

        from xlsx_charts.engine import set_cell_style

        path = _copy("dashboard.xlsx", tmp_path)
        sheet = openpyxl.load_workbook(str(path)).sheetnames[0]

        set_cell_style(str(path), sheet, "A1", bold="true")
        set_cell_style(str(path), sheet, "A1", number_format="0.00")
        assert openpyxl.load_workbook(str(path))[sheet]["A1"].font.bold is True


# ---------------------------------------------------------------------------
# The shared parser
# ---------------------------------------------------------------------------


class TestTriState:
    def test_none_means_leave_and_is_not_falsy_by_accident(self):
        """`if parse(...)` would reintroduce the bug. Callers must use `is not None`."""
        from shared import tristate

        assert tristate.parse("", "bold") is None
        assert tristate.parse("false", "bold") is False
        # Both are falsy. That is exactly why the contract is `is not None`.
        assert not tristate.parse("", "bold")
        assert not tristate.parse("false", "bold")

    @pytest.mark.parametrize("value", ["true", "TRUE", " True ", "yes", "on", "1"])
    def test_the_spellings_of_true(self, value):
        from shared import tristate

        assert tristate.parse(value, "bold") is True

    @pytest.mark.parametrize("value", ["false", "FALSE", " False ", "no", "off", "0"])
    def test_the_spellings_of_false(self, value):
        from shared import tristate

        assert tristate.parse(value, "bold") is False

    def test_a_python_bool_still_works_for_engine_level_callers(self):
        """The engine is called directly by tests and by sibling modules.

        The MCP schema is `str`, so a JSON boolean is refused at the boundary --
        loudly, which is the point. Below that boundary a real bool is
        unambiguous and is accepted.
        """
        from shared import tristate

        assert tristate.parse(True, "bold") is True  # type: ignore[arg-type]
        assert tristate.parse(False, "bold") is False  # type: ignore[arg-type]

    def test_an_unknown_value_names_all_three_states(self):
        from shared import tristate

        with pytest.raises(tristate.TriStateError) as caught:
            tristate.parse("sometimes", "italic")
        assert "italic" in caught.value.hint
        assert "'true'" in caught.value.hint and "'false'" in caught.value.hint
