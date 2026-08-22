"""pptx create_from_template now fills placeholders, like docx and xlsx.

It used to only copy the file, so of the three template tools one could not do
the thing the other two exist for. A caller who filled a Word template and an
Excel template and then tried the same call on a deck got an "unexpected keyword
argument" -- the same trap that `output_path` and the required-`substitutions`
mapping had already set twice.

Formatting is the whole point of starting from a template, so replacement is
done run by run. Assigning to text_frame.text would flatten the runs and drop
every font, size and colour the template author chose -- the rule pptx_basic's
_set_shape_text already follows.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx_new.engine import create_from_template


@pytest.fixture()
def template(tmp_path: Path) -> Path:
    p = tmp_path / "tpl.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Report for {{CLIENT}}"  # type: ignore[reportOptionalMemberAccess]
    body = slide.placeholders[1].text_frame  # type: ignore[reportAttributeAccessIssue]
    body.text = "Prepared by {{AUTHOR}}"
    body.paragraphs[0].runs[0].font.size = Pt(28)
    body.paragraphs[0].runs[0].font.bold = True

    table_slide = prs.slides.add_slide(prs.slide_layouts[5])
    shape = table_slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(1.5))
    shape.table.cell(0, 0).text = "Client"
    shape.table.cell(0, 1).text = "{{CLIENT}}"
    prs.save(str(p))
    return p


SUBS = {"{{CLIENT}}": "Acme Corporation", "{{AUTHOR}}": "Alex"}


class TestItFillsPlaceholders:
    def test_the_call_is_accepted(self, template: Path, tmp_path: Path):
        out = tmp_path / "filled.pptx"
        r = create_from_template(str(template), str(out), SUBS, open_after=False)
        assert r["success"] is True, r.get("error")

    def test_the_title_is_filled(self, template: Path, tmp_path: Path):
        out = tmp_path / "filled.pptx"
        create_from_template(str(template), str(out), SUBS, open_after=False)
        title = Presentation(str(out)).slides[0].shapes.title.text  # type: ignore[reportOptionalMemberAccess]
        assert title == "Report for Acme Corporation"

    def test_the_body_is_filled(self, template: Path, tmp_path: Path):
        out = tmp_path / "filled.pptx"
        create_from_template(str(template), str(out), SUBS, open_after=False)
        body = Presentation(str(out)).slides[0].placeholders[1].text_frame.text  # type: ignore[reportAttributeAccessIssue]
        assert "Alex" in body
        assert "{{AUTHOR}}" not in body

    def test_table_cells_are_filled(self, template: Path, tmp_path: Path):
        """Template placeholders usually live in tables."""
        out = tmp_path / "filled.pptx"
        create_from_template(str(template), str(out), SUBS, open_after=False)
        table = next(s for s in Presentation(str(out)).slides[1].shapes if s.has_table).table  # type: ignore[reportAttributeAccessIssue]
        assert table.cell(0, 1).text == "Acme Corporation"

    def test_it_reports_how_many_it_applied(self, template: Path, tmp_path: Path):
        out = tmp_path / "filled.pptx"
        r = create_from_template(str(template), str(out), SUBS, open_after=False)
        assert r["substitutions_applied"] >= 3


class TestFormattingSurvives:
    def test_font_size_is_kept(self, template: Path, tmp_path: Path):
        out = tmp_path / "filled.pptx"
        create_from_template(str(template), str(out), SUBS, open_after=False)
        run = Presentation(str(out)).slides[0].placeholders[1].text_frame.paragraphs[0].runs[0]  # type: ignore[reportAttributeAccessIssue]
        assert run.font.size == Pt(28), "assigning text_frame.text would have dropped this"

    def test_bold_is_kept(self, template: Path, tmp_path: Path):
        out = tmp_path / "filled.pptx"
        create_from_template(str(template), str(out), SUBS, open_after=False)
        run = Presentation(str(out)).slides[0].placeholders[1].text_frame.paragraphs[0].runs[0]  # type: ignore[reportAttributeAccessIssue]
        assert run.font.bold is True


class TestTheOldBehaviourStillWorks:
    def test_no_substitutions_is_a_plain_copy(self, template: Path, tmp_path: Path):
        out = tmp_path / "copy.pptx"
        r = create_from_template(str(template), str(out), open_after=False)
        assert r["success"] is True, r.get("error")
        assert r["substitutions_applied"] == 0
        assert Presentation(str(out)).slides[0].shapes.title.text == "Report for {{CLIENT}}"  # type: ignore[reportOptionalMemberAccess]

    def test_the_source_template_is_untouched(self, template: Path, tmp_path: Path):
        before = template.read_bytes()
        create_from_template(str(template), str(tmp_path / "filled.pptx"), SUBS, open_after=False)
        assert template.read_bytes() == before

    def test_slide_count_is_still_reported(self, template: Path, tmp_path: Path):
        r = create_from_template(str(template), str(tmp_path / "c.pptx"), open_after=False)
        assert r["slide_count"] == 2


class TestBadInput:
    def test_a_non_dict_is_rejected_with_a_usable_hint(self, template: Path, tmp_path: Path):
        r = create_from_template(str(template), str(tmp_path / "x.pptx"), "nope", open_after=False)  # type: ignore[reportArgumentType]
        assert r["success"] is False
        assert "dict" in r["error"]
        assert "PLACEHOLDER" in r["hint"]

    def test_a_missing_placeholder_is_not_an_error(self, template: Path, tmp_path: Path):
        """Same as docx: a key that is not in the template is a no-op."""
        out = tmp_path / "x.pptx"
        r = create_from_template(str(template), str(out), {"{{NOPE}}": "x"}, open_after=False)
        assert r["success"] is True
        assert r["substitutions_applied"] == 0
