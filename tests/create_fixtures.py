"""Script to generate test fixture files for office-mcp tests.

Run once: uv run python tests/create_fixtures.py
"""

from pathlib import Path

FIXTURES = Path(__file__).parent / "fixtures"
FIXTURES.mkdir(exist_ok=True)


def create_docx_fixtures() -> None:
    from docx import Document

    # contract_simple.docx — plain paragraphs, no tables
    doc = Document()
    doc.add_heading("Service Agreement", level=1)
    p = doc.add_paragraph()
    p.add_run("This Service Agreement (the ").bold = False
    run = p.add_run('"Agreement"')
    run.bold = True
    p.add_run(") is entered into as of EFFECTIVE_DATE.")

    doc.add_heading("1. Parties", level=2)
    doc.add_paragraph('PARTY_A_NAME ("Client") and PARTY_B_NAME ("Service Provider") agree to the following terms.')

    doc.add_heading("2. Scope of Work", level=2)
    doc.add_paragraph(
        "Service Provider shall deliver software development services as described in Exhibit A attached hereto."
    )
    doc.add_paragraph("All deliverables shall conform to the specifications agreed upon in writing.")

    doc.add_heading("3. Payment Terms", level=2)
    p2 = doc.add_paragraph()
    p2.add_run("Payment shall be made within ")
    run2 = p2.add_run("30 days")
    run2.bold = True
    p2.add_run(" of invoice receipt.")
    doc.add_paragraph("Contract value is CONTRACT_VALUE USD. Late payments incur a 1.5% monthly fee.")

    doc.add_heading("4. Termination", level=2)
    doc.add_paragraph("Either party may terminate this Agreement with 30 days written notice.")
    doc.add_paragraph("Upon termination, all outstanding invoices become immediately due.")

    doc.save(str(FIXTURES / "contract_simple.docx"))
    print("Created contract_simple.docx")

    # contract_complex.docx — multi-run bold/italic, tables, header
    doc2 = Document()
    section = doc2.sections[0]
    header = section.header
    header.paragraphs[0].text = "CONFIDENTIAL — PARTY_A_NAME"

    doc2.add_heading("Master Services Agreement", level=1)

    p = doc2.add_paragraph()
    p.add_run("This Master Services Agreement is between ")
    run_a = p.add_run("PARTY_A_NAME")
    run_a.bold = True
    p.add_run(' ("Client") and ')
    run_b = p.add_run("PARTY_B_NAME")
    run_b.bold = True
    p.add_run(' ("Vendor"), effective ')
    run_d = p.add_run("EFFECTIVE_DATE")
    run_d.italic = True
    p.add_run(".")

    doc2.add_heading("1. Services", level=2)
    for i in range(1, 6):
        doc2.add_paragraph(f"Section 1.{i}: Lorem ipsum dolor sit amet, clause {i}.", style="Body Text")

    doc2.add_heading("2. Fees and Payment", level=2)
    p3 = doc2.add_paragraph()
    p3.add_run("Total contract value: ")
    run_v = p3.add_run("CONTRACT_VALUE")
    run_v.bold = True
    run_v.italic = True
    p3.add_run(" USD.")

    # Add a table
    table = doc2.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    headers = ["Item", "Quantity", "Price"]
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    data = [
        ("Development", "1", "$10,000"),
        ("Testing", "1", "$3,000"),
        ("Deployment", "1", "$2,000"),
    ]
    for r, (item, qty, price) in enumerate(data, start=1):
        table.rows[r].cells[0].text = item
        table.rows[r].cells[1].text = qty
        table.rows[r].cells[2].text = price

    doc2.add_heading("3. Termination", level=2)
    doc2.add_paragraph("Either party may terminate with 60 days notice.")

    doc2.save(str(FIXTURES / "contract_complex.docx"))
    print("Created contract_complex.docx")

    # report_tables.docx — multiple tables
    doc3 = Document()
    doc3.add_heading("Quarterly Report", level=1)
    doc3.add_paragraph("This report covers Q1 financial performance.")

    for t_num in range(1, 3):
        doc3.add_heading(f"Table {t_num}: Revenue Data", level=2)
        table = doc3.add_table(rows=5, cols=4)
        table.style = "Table Grid"
        hdr = table.rows[0]
        for i, h in enumerate(["Region", "Q1", "Q2", "Total"]):
            hdr.cells[i].text = h
        regions = ["North", "South", "East", "West"]
        for r, region in enumerate(regions, start=1):
            table.rows[r].cells[0].text = region
            table.rows[r].cells[1].text = str(10000 * r)
            table.rows[r].cells[2].text = str(12000 * r)
            table.rows[r].cells[3].text = str(22000 * r)

    doc3.save(str(FIXTURES / "report_tables.docx"))
    print("Created report_tables.docx")


def create_xlsx_fixtures() -> None:
    from openpyxl import Workbook

    # budget_simple.xlsx — flat data, no formulas
    wb = Workbook()
    ws = wb.active
    ws.title = "Q3 Revenue"
    headers = ["Region", "Jan", "Feb", "Mar", "Q3 Total", "YoY%"]
    ws.append(headers)
    data = [
        ("North", 120000, 135000, 118000, 373000, 0.12),
        ("South", 98000, 102000, 115000, 315000, 0.08),
        ("East", 87000, 91000, 95000, 273000, 0.15),
        ("West", 143000, 138000, 151000, 432000, 0.21),
    ]
    for row in data:
        ws.append(row)

    ws2 = wb.create_sheet("Dashboard")
    ws2["A1"] = "Summary Dashboard"
    ws2["A2"] = "Total Revenue"
    ws2["B2"] = 1393000

    wb.save(str(FIXTURES / "budget_simple.xlsx"))
    print("Created budget_simple.xlsx")

    # budget_formulas.xlsx — SUM, IF, VLOOKUP formulas
    wb2 = Workbook()
    ws3 = wb2.active
    ws3.title = "Q3 Revenue"
    ws3.append(["Region", "Jan", "Feb", "Mar", "Q3 Total", "YoY%"])
    raw_data = [
        ("North", 120000, 135000, 118000),
        ("South", 98000, 102000, 115000),
        ("East", 87000, 91000, 95000),
        ("West", 143000, 138000, 151000),
    ]
    for i, (region, jan, feb, mar) in enumerate(raw_data, start=2):
        ws3[f"A{i}"] = region
        ws3[f"B{i}"] = jan
        ws3[f"C{i}"] = feb
        ws3[f"D{i}"] = mar
        ws3[f"E{i}"] = f"=SUM(B{i}:D{i})"
        ws3[f"F{i}"] = f'=IF(E{i}>300000,"High","Normal")'

    ws3["A7"] = "Total"
    ws3["B7"] = "=SUM(B2:B5)"
    ws3["C7"] = "=SUM(C2:C5)"
    ws3["D7"] = "=SUM(D2:D5)"
    ws3["E7"] = "=SUM(E2:E5)"

    wb2.save(str(FIXTURES / "budget_formulas.xlsx"))
    print("Created budget_formulas.xlsx")

    # dashboard.xlsx — with some cell styles (no charts — need openpyxl chart objects)
    wb3 = Workbook()
    ws4 = wb3.active
    ws4.title = "Dashboard"
    ws4["A1"] = "KPI Dashboard"
    ws4["A3"] = "Metric"
    ws4["B3"] = "Value"
    ws4["C3"] = "Target"
    ws4["D3"] = "Status"
    kpis = [
        ("Revenue", 1393000, 1200000, "Above"),
        ("Costs", 850000, 900000, "Good"),
        ("NPS Score", 72, 70, "Good"),
        ("Churn Rate", 0.03, 0.05, "Good"),
    ]
    for i, (metric, val, target, status) in enumerate(kpis, start=4):
        ws4[f"A{i}"] = metric
        ws4[f"B{i}"] = val
        ws4[f"C{i}"] = target
        ws4[f"D{i}"] = status

    wb3.save(str(FIXTURES / "dashboard.xlsx"))
    print("Created dashboard.xlsx")


def create_pptx_fixtures() -> None:
    from pptx import Presentation

    # deck_simple.pptx — title + content slides
    prs = Presentation()
    title_layout = prs.slide_layouts[0]  # Title Slide
    content_layout = prs.slide_layouts[1]  # Title and Content

    slide0 = prs.slides.add_slide(title_layout)
    slide0.shapes.title.text = "Q3 2026 Business Review"
    slide0.placeholders[1].text = "Prepared by Strategy Team"

    titles_bodies = [
        ("Revenue Performance", "Q3 total: $1.4M\nUp 18% YoY\nAll regions growing"),
        ("Key Metrics", "CAC: $420\nLTV: $2,100\nNPS: 72"),
        ("Roadmap", "Q4 priorities:\n• Launch v2.0\n• Expand to EU\n• Hire 5 engineers"),
        ("Summary", "Strong quarter.\nMomentum continues into Q4."),
    ]
    for title, body in titles_bodies:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = body

    prs.save(str(FIXTURES / "deck_simple.pptx"))
    print("Created deck_simple.pptx")

    # deck_images.pptx — slides with text (images need actual image files)
    prs2 = Presentation()
    content_layout2 = prs2.slide_layouts[1]

    slide_t = prs2.slides.add_slide(prs2.slide_layouts[0])
    slide_t.shapes.title.text = "Visual Deck"
    slide_t.placeholders[1].text = "Image-heavy presentation"

    for i in range(1, 4):
        slide = prs2.slides.add_slide(content_layout2)
        slide.shapes.title.text = f"Slide {i} — Content"
        slide.placeholders[1].text = f"Content for slide {i}\nBullet point A\nBullet point B"

    prs2.save(str(FIXTURES / "deck_images.pptx"))
    print("Created deck_images.pptx")


if __name__ == "__main__":
    print("Creating test fixtures...")
    try:
        create_docx_fixtures()
    except ImportError as e:
        print(f"Skipping DOCX fixtures (missing dep): {e}")
    try:
        create_xlsx_fixtures()
    except ImportError as e:
        print(f"Skipping XLSX fixtures (missing dep): {e}")
    try:
        create_pptx_fixtures()
    except ImportError as e:
        print(f"Skipping PPTX fixtures (missing dep): {e}")
    print("Done.")
