"""A content placeholder draws its own bullet, so text carrying its own renders
"• • Point 1". Found by converting a generated deck to PDF and looking at it —
the file was valid OOXML and every tool involved reported success."""

from __future__ import annotations

import pytest

from shared.slide_text import strip_list_markers


class TestStripListMarkers:
    @pytest.mark.parametrize("marker", ["•", "-", "*", "–", "—", "‣", "◦", "∙"])
    def test_each_marker_style_is_removed(self, marker: str):
        assert strip_list_markers(f"{marker} Point 1") == "Point 1"

    def test_every_line_is_cleaned(self):
        assert strip_list_markers("• A\n• B\n• C") == "A\nB\nC"

    def test_mixed_markers_across_lines(self):
        assert strip_list_markers("- A\n* B\n• C") == "A\nB\nC"

    def test_untouched_when_there_is_no_marker(self):
        assert strip_list_markers("Point 1\nPoint 2") == "Point 1\nPoint 2"

    def test_negative_numbers_survive(self):
        """'-5 degrees' is content, not a list item — the marker needs whitespace."""
        assert strip_list_markers("-5 degrees overnight") == "-5 degrees overnight"

    def test_emphasis_markup_survives(self):
        assert strip_list_markers("*emphasis* matters") == "*emphasis* matters"

    def test_a_marker_alone_on_a_line_survives(self):
        assert strip_list_markers("-") == "-"

    def test_nested_marker_keeps_one_level(self):
        assert strip_list_markers("- - sub item") == "- sub item"

    def test_empty_and_blank_input(self):
        assert strip_list_markers("") == ""
        assert strip_list_markers("\n\n") == "\n\n"

    def test_trailing_whitespace_is_trimmed(self):
        assert strip_list_markers("• Point   ") == "Point"


class TestDeckBodyTextIsClean:
    def test_outline_bullets_are_not_doubled(self, tmp_path):
        from pptx import Presentation

        from servers.pptx_new.pptx_new.engine import create_from_outline

        out = tmp_path / "deck.pptx"
        result = create_from_outline(
            str(out),
            [{"title": "Main Points", "content": "• Point 1\n• Point 2\n• Point 3"}],
            open_after=False,
        )
        assert result["success"] is True

        texts: list[str] = []
        for slide in Presentation(str(out)).slides:
            for shape in slide.shapes:
                frame = getattr(shape, "text_frame", None)
                if frame is None:
                    continue
                texts.extend(run.text for para in frame.paragraphs for run in para.runs)
        assert "Point 1" in texts
        assert not any(t.startswith("•") for t in texts)
