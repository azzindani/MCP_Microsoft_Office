"""Slide titles vanished, and success said nothing was wrong.

Two deck builders in the same module both take "a list of slide dicts" and each
invented its own key names:

    create_from_outline    slide_def.get("title")    .get("content")
    create_deck_from_data  item.get("heading")       .get("bullets")

Both parameters are typed list[dict[str, Any]], so tools/list shows an opaque
array with no item schema. A caller cannot discover which spelling a tool wants;
it guesses. Passing "title" to create_deck_from_data -- the same word that tool
uses for its own deck title, one parameter above -- returned success:true with
every content slide's title blank. Rendered to PDF the deck was bullet lists
floating under empty headers, and nothing in the response said so.

Both spellings now work in both tools, and a slide dict with no recognised
heading key gets a warning naming the keys it did see, so a genuine typo is
visible instead of silent.

Also here: a chart's value axis printed raw "500000 / 1000000 / 2500000". Too
wide for the plot, the labels rotated 45 degrees, and on a chart squeezed under
a bullet list one of the two category names was dropped entirely. Formatted to
"0.5M / 1M / 2.5M" they stay level and both categories fit.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from pptx_design.engine import _axis_number_format, add_chart  # type: ignore[reportMissingImports]
from pptx_new.engine import create_deck_from_data, create_from_outline  # type: ignore[reportMissingImports]

BULLETS = ["Google Ads took 1.9M of 2.5M total spend", "Facebook Ads delivered 564K spend"]


def _titles(path: Path) -> list[str]:
    return [s.shapes.title.text if s.shapes.title else "" for s in Presentation(str(path)).slides]


def _warnings(result: dict) -> list[dict]:
    return [s for s in result.get("progress", []) if s.get("icon") == "⚠"]


class TestEitherSpellingWorks:
    @pytest.mark.parametrize("key", ["title", "heading"])
    def test_create_deck_from_data_keeps_the_heading(self, key: str, tmp_path: Path):
        out = tmp_path / "deck.pptx"
        r = create_deck_from_data(
            str(out), "Deck", [{key: "Where the money went", "bullets": BULLETS}], open_after=False
        )
        assert r["success"] is True, r.get("error")
        assert _titles(out)[1] == "Where the money went"

    @pytest.mark.parametrize("key", ["title", "heading"])
    def test_create_from_outline_keeps_the_heading(self, key: str, tmp_path: Path):
        out = tmp_path / "deck.pptx"
        r = create_from_outline(str(out), [{key: "Quarterly review", "content": "Body text"}], open_after=False)
        assert r["success"] is True, r.get("error")
        assert _titles(out)[0] == "Quarterly review"

    @pytest.mark.parametrize("body_key", ["bullets", "content"])
    def test_either_body_spelling_reaches_the_slide(self, body_key: str, tmp_path: Path):
        out = tmp_path / "deck.pptx"
        value = BULLETS if body_key == "bullets" else "\n".join(BULLETS)
        r = create_deck_from_data(str(out), "Deck", [{"title": "T", body_key: value}], open_after=False)
        assert r["success"] is True, r.get("error")
        shapes = Presentation(str(out)).slides[1].shapes
        text = "\n".join(sh.text_frame.text for sh in shapes if sh.has_text_frame)  # type: ignore[reportAttributeAccessIssue]
        assert "Google Ads" in text


class TestATypoIsNotSilent:
    """The bug was not that a key was missing, it was that nothing said so."""

    def test_an_unrecognised_heading_key_warns(self, tmp_path: Path):
        out = tmp_path / "deck.pptx"
        r = create_deck_from_data(str(out), "Deck", [{"headline": "typo", "bullets": BULLETS}], open_after=False)
        assert r["success"] is True
        assert _warnings(r), "a slide with no heading must say so"

    def test_the_warning_names_the_keys_it_saw(self, tmp_path: Path):
        out = tmp_path / "deck.pptx"
        r = create_deck_from_data(str(out), "Deck", [{"headline": "typo", "bullets": BULLETS}], open_after=False)
        detail = _warnings(r)[0]["detail"]
        assert "headline" in detail
        assert "title" in detail and "heading" in detail

    def test_a_correct_key_produces_no_warning(self, tmp_path: Path):
        out = tmp_path / "deck.pptx"
        r = create_deck_from_data(str(out), "Deck", [{"title": "fine", "bullets": BULLETS}], open_after=False)
        assert not _warnings(r)

    def test_outline_warns_too(self, tmp_path: Path):
        out = tmp_path / "deck.pptx"
        r = create_from_outline(str(out), [{"headline": "typo", "content": "x"}], open_after=False)
        assert _warnings(r)


class TestAxisLabelsAreReadable:
    @pytest.mark.parametrize(
        ("values", "expected"),
        [
            ([1_939_000, 564_100], '#,##0.#,,"M"'),
            ([45_000, 12_000], '#,##0.#,"K"'),
            ([12, 7.5], "#,##0.##"),
            ([], "#,##0.##"),
        ],
    )
    def test_the_format_follows_the_magnitude(self, values: list, expected: str):
        assert _axis_number_format(values) == expected

    def test_non_numeric_values_do_not_raise(self):
        assert _axis_number_format([None, "n/a", 5]) == "#,##0.##"

    def test_a_millions_axis_is_formatted_on_a_real_chart(self, tmp_path: Path):
        out = tmp_path / "deck.pptx"
        create_deck_from_data(str(out), "Deck", [{"title": "T", "bullets": BULLETS}], open_after=False)
        r = add_chart(
            str(out),
            1,
            "bar",
            {"categories": ["Google Ads", "Facebook Ads"], "series": {"Spend": [1_939_000, 564_100]}},
            title="Total spend by platform",
        )
        assert r["success"] is True, r.get("error")
        chart = next(sh.chart for sh in Presentation(str(out)).slides[1].shapes if sh.has_chart)  # type: ignore[reportAttributeAccessIssue]
        assert chart.value_axis.tick_labels.number_format == '#,##0.#,,"M"'

    def test_one_series_gets_no_legend(self, tmp_path: Path):
        """It would only repeat the chart title."""
        out = tmp_path / "deck.pptx"
        create_deck_from_data(str(out), "Deck", [{"title": "T", "bullets": BULLETS}], open_after=False)
        add_chart(str(out), 1, "bar", {"categories": ["A", "B"], "series": {"Spend": [1, 2]}}, title="t")
        chart = next(sh.chart for sh in Presentation(str(out)).slides[1].shapes if sh.has_chart)  # type: ignore[reportAttributeAccessIssue]
        assert chart.has_legend is False

    def test_more_than_one_series_gets_a_legend(self, tmp_path: Path):
        out = tmp_path / "deck.pptx"
        create_deck_from_data(str(out), "Deck", [{"title": "T", "bullets": BULLETS}], open_after=False)
        add_chart(
            str(out),
            1,
            "bar",
            {"categories": ["A", "B"], "series": {"Spend": [1, 2], "Clicks": [3, 4]}},
            title="t",
        )
        chart = next(sh.chart for sh in Presentation(str(out)).slides[1].shapes if sh.has_chart)  # type: ignore[reportAttributeAccessIssue]
        assert chart.has_legend is True

    def test_a_pie_gets_a_legend_and_no_axis_call(self, tmp_path: Path):
        """A pie has no value axis; asking for one used to be the crash risk."""
        out = tmp_path / "deck.pptx"
        create_deck_from_data(str(out), "Deck", [{"title": "T", "bullets": BULLETS}], open_after=False)
        r = add_chart(str(out), 1, "pie", {"categories": ["A", "B"], "series": {"Share": [30, 70]}}, title="t")
        assert r["success"] is True, r.get("error")
        chart = next(sh.chart for sh in Presentation(str(out)).slides[1].shapes if sh.has_chart)  # type: ignore[reportAttributeAccessIssue]
        assert chart.has_legend is True
