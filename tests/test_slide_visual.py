"""Two ways a deck goes wrong that no return value revealed.

Both were found by rendering a swept deck to PDF and looking at it. Every tool
involved returned success:true and produced structurally valid pptx.

- `add_chart` handed its inch box straight to python-pptx, which placed the
  chart past the bottom edge of the slide. Two of five category labels were cut
  off and nothing said so.
- `set_background` sets one slide; `set_font_all_slides` sets every slide. Dark
  background on slide 1 plus white text everywhere left slide 2 with white text
  on white -- present, "successful", invisible.
"""

from __future__ import annotations

from typing import Any

import pytest
from pptx import Presentation
from pptx.util import Inches

from shared.slide_visual import (
    MIN_CONTRAST,
    contrast_ratio,
    contrast_warning,
    fit_to_slide,
    slide_background_hex,
    unreadable_slides,
)


@pytest.fixture()
def prs() -> Any:
    # Presentation is a factory function, not a class, so it cannot annotate.
    p = Presentation()
    for _ in range(2):
        p.slides.add_slide(p.slide_layouts[6])  # blank
    return p


def _paint(slide, hex_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    from pptx.dml.color import RGBColor

    fill.fore_color.rgb = RGBColor.from_string(hex_color)


class TestContrastMaths:
    def test_black_on_white_is_the_maximum(self):
        assert round(contrast_ratio("000000", "FFFFFF"), 1) == 21.0

    def test_a_colour_against_itself_is_the_minimum(self):
        assert contrast_ratio("FFFFFF", "FFFFFF") == 1.0

    def test_white_on_navy_is_readable(self):
        assert contrast_ratio("FFFFFF", "1B2A47") > MIN_CONTRAST


class TestUnreadableSlides:
    def test_white_text_on_a_white_slide_is_reported_invisible(self, prs):
        _paint(prs.slides[0], "FFFFFF")
        offender = next(o for o in unreadable_slides(prs, "FFFFFF") if o["slide"] == 0)
        assert offender["invisible"] is True
        assert offender["assumed"] is False

    def test_white_text_on_a_dark_slide_is_fine(self, prs):
        """Slide 1 is deliberately left inheriting, so it is still flagged; the
        painted dark slide is the one that must not be."""
        _paint(prs.slides[0], "1B2A47")
        assert 0 not in [o["slide"] for o in unreadable_slides(prs, "FFFFFF")]

    def test_the_mixed_deck_that_caused_this(self, prs):
        """Slide 0 dark, slide 1 left at default: white everywhere flags only 1."""
        _paint(prs.slides[0], "1B2A47")
        _paint(prs.slides[1], "FFFFFF")
        assert [o["slide"] for o in unreadable_slides(prs, "FFFFFF")] == [1]

    def test_an_inherited_background_is_still_checked(self, prs):
        """The deck that prompted all this had its unreadable slide left at the
        template default -- white by inheritance, not by an explicit fill. A
        check that only looked at the slide's own fill saw nothing and missed
        the exact bug it was written for."""
        assert slide_background_hex(prs.slides[0]) is None  # nothing set anywhere
        offenders = unreadable_slides(prs, "FFFFFF")
        assert [o["slide"] for o in offenders] == [0, 1]
        assert all(o["assumed"] for o in offenders)

    def test_dark_text_on_an_inherited_background_is_fine(self):
        """The assumption must not fire on the normal case: default template,
        default dark text."""
        from pptx import Presentation as _P

        deck = _P()
        deck.slides.add_slide(deck.slide_layouts[6])
        assert unreadable_slides(deck, "000000") == []

    def test_an_explicit_background_is_not_marked_assumed(self, prs):
        _paint(prs.slides[0], "FFFFFF")
        offender = next(o for o in unreadable_slides(prs, "FFFFFF") if o["slide"] == 0)
        assert offender["assumed"] is False

    def test_the_warning_names_the_slides_and_the_fix(self, prs):
        _paint(prs.slides[1], "FFFFFF")
        message = contrast_warning(unreadable_slides(prs, "FFFFFF"), "FFFFFF")
        assert "1" in message
        assert "set_background" in message

    def test_the_warning_says_when_the_background_was_inherited(self, prs):
        """A caller using a themed master should know the check assumed white."""
        message = contrast_warning(unreadable_slides(prs, "FFFFFF"), "FFFFFF")
        assert "inherit" in message


class TestFitToSlide:
    def test_a_box_inside_the_canvas_is_untouched(self, prs):
        left, top, width, height, note = fit_to_slide(prs, 1.0, 1.0, 4.0, 3.0)
        assert (left, top, width, height) == (1.0, 1.0, 4.0, 3.0)
        assert note == ""

    def test_a_chart_hanging_off_the_bottom_is_pulled_back(self, prs):
        slide_h = prs.slide_height / 914400
        _, top, _, height, note = fit_to_slide(prs, 1.0, slide_h - 1.0, 6.0, 4.5)
        assert top + height <= slide_h
        assert note

    def test_a_box_larger_than_the_slide_is_shrunk(self, prs):
        slide_w = prs.slide_width / 914400
        slide_h = prs.slide_height / 914400
        left, top, width, height, note = fit_to_slide(prs, 0.0, 0.0, slide_w * 2, slide_h * 2)
        assert left + width <= slide_w
        assert top + height <= slide_h
        assert note

    def test_the_note_says_what_moved(self, prs):
        _, _, _, _, note = fit_to_slide(prs, 20.0, 20.0, 6.0, 4.5)
        assert "slide" in note.lower()


class TestAddChartUsesIt:
    def test_an_overflowing_chart_is_fitted_and_reported(self, tmp_path):
        from pptx_design.engine import add_chart

        deck = tmp_path / "d.pptx"
        p = Presentation()
        p.slides.add_slide(p.slide_layouts[6])
        p.save(deck)

        result = add_chart(
            str(deck),
            slide_index=0,
            chart_type="bar",
            data={"categories": ["A", "B"], "series": {"S": [1, 2]}},
            title="Off the edge",
            left=1.0,
            top=6.0,
            width=6.0,
            height=4.5,
        )
        assert result["success"] is True
        assert any("fit" in str(step).lower() for step in result["progress"])

        check = Presentation(str(deck))
        shape = next(s for s in check.slides[0].shapes if s.has_chart)
        assert int(shape.top or 0) + int(shape.height or 0) <= int(check.slide_height or 0)
        assert int(shape.left or 0) + int(shape.width or 0) <= int(check.slide_width or 0)


class TestBackgroundOnEverySlide:
    """set_background took one slide while set_font_all_slides took all of them,
    so "make this deck dark" needed a call per slide -- which is exactly how a
    deck ended up dark on slide 0, untouched white on slide 1, and white text on
    both. slide_index=-1 makes the two operations symmetric."""

    def _deck(self, tmp_path):
        deck = tmp_path / "d.pptx"
        p = Presentation()
        for _ in range(3):
            p.slides.add_slide(p.slide_layouts[6])
        p.save(deck)
        return deck

    def test_minus_one_paints_every_slide(self, tmp_path):
        from pptx_design.engine import set_background

        deck = self._deck(tmp_path)
        result = set_background(str(deck), -1, color_hex="1B2A47")
        assert result["success"] is True

        check = Presentation(str(deck))
        assert [slide_background_hex(s) for s in check.slides] == ["1B2A47"] * 3

    def test_a_single_index_still_paints_only_that_slide(self, tmp_path):
        from pptx_design.engine import set_background

        deck = self._deck(tmp_path)
        set_background(str(deck), 1, color_hex="1B2A47")

        check = Presentation(str(deck))
        assert [slide_background_hex(s) for s in check.slides] == [None, "1B2A47", None]

    def test_out_of_range_is_still_rejected(self, tmp_path):
        from pptx_design.engine import set_background

        result = set_background(str(self._deck(tmp_path)), 99, color_hex="1B2A47")
        assert result["success"] is False

    def test_the_pairing_that_broke_the_deck_now_produces_readable_text(self, tmp_path):
        from pptx_design.engine import set_background, set_font_all_slides

        deck = self._deck(tmp_path)
        set_background(str(deck), -1, color_hex="1B2A47")
        result = set_font_all_slides(str(deck), color_hex="FFFFFF")
        assert result.get("unreadable_slides") is None


class TestChartTextFollowsTheSlide:
    """Colouring only the chart title left the category names and the value
    ticks black on a navy slide -- visible in the render and still wrong."""

    def _deck_with_chart(self, tmp_path, background: str | None):
        from pptx_design.engine import add_chart, set_background

        deck = tmp_path / "c.pptx"
        p = Presentation()
        p.slides.add_slide(p.slide_layouts[6])
        p.save(deck)
        if background:
            set_background(str(deck), 0, color_hex=background)
        add_chart(
            str(deck),
            slide_index=0,
            chart_type="bar",
            data={"categories": ["Instagram", "LinkedIn"], "series": {"Spends": [5, 3]}},
            title="Spends by Platform",
        )
        chart = next(s for s in Presentation(str(deck)).slides[0].shapes if s.has_chart).chart
        return chart

    def test_axis_labels_go_light_on_a_dark_slide(self, tmp_path):
        chart = self._deck_with_chart(tmp_path, "1B2A47")
        assert str(chart.category_axis.tick_labels.font.color.rgb) == "FFFFFF"
        assert str(chart.value_axis.tick_labels.font.color.rgb) == "FFFFFF"

    def test_the_title_goes_light_too(self, tmp_path):
        chart = self._deck_with_chart(tmp_path, "1B2A47")
        colors = {
            str(p.font.color.rgb) for p in chart.chart_title.text_frame.paragraphs if p.font.color.rgb is not None
        }
        assert colors == {"FFFFFF"}

    def test_axis_labels_go_dark_on_a_light_slide(self, tmp_path):
        chart = self._deck_with_chart(tmp_path, "FFFFFF")
        assert str(chart.category_axis.tick_labels.font.color.rgb) == "000000"

    def test_a_pie_chart_has_no_axes_and_must_not_crash(self, tmp_path):
        from pptx_design.engine import add_chart, set_background

        deck = tmp_path / "pie.pptx"
        p = Presentation()
        p.slides.add_slide(p.slide_layouts[6])
        p.save(deck)
        set_background(str(deck), 0, color_hex="1B2A47")
        result = add_chart(
            str(deck),
            slide_index=0,
            chart_type="pie",
            data={"categories": ["A", "B"], "series": {"S": [1, 2]}},
            title="Split",
        )
        assert result["success"] is True


class TestNewShapesDoNotLandOnExistingContent:
    """add_table and add_chart both default to top=2.0in, which is where a
    layout puts its body text. A swept deck ended up with bullets, a table and
    a chart stacked in the same region -- three shapes, none of them readable,
    and all 24 tools in that phase reported PASS."""

    def _slide_with_text(self, tmp_path, text="Total spend across 16834 rows"):
        tmp_path.mkdir(parents=True, exist_ok=True)
        deck = tmp_path / "t.pptx"
        p = Presentation()
        slide = p.slides.add_slide(p.slide_layouts[1])  # title + content
        slide.shapes.title.text = "Spend Overview"
        slide.placeholders[1].text_frame.text = text
        p.save(deck)
        return deck

    def test_a_table_is_pushed_below_the_bullets(self, tmp_path):
        from pptx_design.engine import add_table

        deck = self._slide_with_text(tmp_path)
        result = add_table(str(deck), 0, 2, 2, [["a", "b"], ["c", "d"]], top=2.0)
        assert result["success"] is True

        slide = Presentation(str(deck)).slides[0]
        table = next(s for s in slide.shapes if s.has_table)
        body = slide.placeholders[1]
        # Below where the text is drawn -- not below the placeholder frame,
        # which is the layout's full body area whether or not it is filled.
        assert table.top > body.top, "table starts above the body text"
        assert table.top > Inches(2.0), "table was not moved off the requested position"

    def test_a_chart_is_pushed_below_the_bullets(self, tmp_path):
        from pptx_design.engine import add_chart

        deck = self._slide_with_text(tmp_path)
        result = add_chart(str(deck), 0, "bar", {"categories": ["A"], "series": {"S": [1]}}, top=2.0, height=2.0)
        assert result["success"] is True

        slide = Presentation(str(deck)).slides[0]
        chart = next(s for s in slide.shapes if s.has_chart)
        body = slide.placeholders[1]
        assert chart.top > body.top
        assert chart.top > Inches(2.0)

    def test_more_text_pushes_the_shape_further_down(self, tmp_path):
        """The point of measuring the drawn text rather than the frame: a wall
        of bullets must displace more than a single line does."""
        from pptx_design.engine import add_table

        short_deck = self._slide_with_text(tmp_path / "s", "One line")
        long_deck = self._slide_with_text(tmp_path / "l", "\n".join(f"Bullet number {i}" for i in range(8)))

        add_table(str(short_deck), 0, 2, 2, [["a", "b"], ["c", "d"]], top=2.0)
        add_table(str(long_deck), 0, 2, 2, [["a", "b"], ["c", "d"]], top=2.0)

        short_top = next(s for s in Presentation(str(short_deck)).slides[0].shapes if s.has_table).top
        long_top = next(s for s in Presentation(str(long_deck)).slides[0].shapes if s.has_table).top
        assert long_top > short_top, "displacement ignores how much text there is"

    def test_the_move_is_reported(self, tmp_path):
        from pptx_design.engine import add_table

        deck = self._slide_with_text(tmp_path)
        result = add_table(str(deck), 0, 2, 2, [["a", "b"], ["c", "d"]], top=2.0)
        assert any("reposition" in str(s).lower() for s in result["progress"])

    def test_an_empty_placeholder_does_not_push_anything_down(self, tmp_path):
        """An unfilled layout placeholder renders as nothing. Treating it as
        content would shove every shape to the bottom of an empty slide."""
        from pptx_design.engine import add_table

        deck = tmp_path / "e.pptx"
        p = Presentation()
        p.slides.add_slide(p.slide_layouts[1])  # placeholders present but empty
        p.save(deck)

        result = add_table(str(deck), 0, 2, 2, [["a", "b"], ["c", "d"]], top=2.0)
        assert result["success"] is True
        table = next(s for s in Presentation(str(deck)).slides[0].shapes if s.has_table)
        assert table.top == Inches(2.0), "an empty placeholder displaced the table"

    def test_a_blank_slide_keeps_the_requested_position(self, tmp_path):
        from pptx_design.engine import add_table

        deck = tmp_path / "b.pptx"
        p = Presentation()
        p.slides.add_slide(p.slide_layouts[6])  # blank
        p.save(deck)

        add_table(str(deck), 0, 2, 2, [["a", "b"], ["c", "d"]], top=2.0)
        table = next(s for s in Presentation(str(deck)).slides[0].shapes if s.has_table)
        assert table.top == Inches(2.0)
