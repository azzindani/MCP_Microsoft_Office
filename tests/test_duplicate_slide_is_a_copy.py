"""A duplicated slide carried every placeholder twice, and back to front.

duplicate_slide built the copy with prs.slides.add_slide(src_layout). python-pptx
clones the layout's placeholders into that new slide, and the loop underneath
then deep-copied the source's own shapes -- placeholders included, with their
text -- on top. Both sets stayed:

    original  ['Title 1', 'Content Placeholder 2']
    copy      ['Content Placeholder 2', 'Title 1', 'Title 1', 'Content Placeholder 2']
                                                    ^^^^^^^^^^ empty clones from the layout

Empty placeholders draw nothing, so the rendered deck looked right and the call
reported success. The damage is to anything addressing a slide by shape:
read_slide returned four shapes where the original had two, two of them sharing
a name with the populated ones, and nothing in the response said which "Title 1"
an edit should go to. Duplicating a duplicate compounded it.

The second bug is in the same two lines. Each copied element was inserted at a
fixed index in the shape tree, so every shape landed ahead of the one copied
before it and the copy listed its shapes in reverse -- body before title.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_design.engine import duplicate_slide  # type: ignore[reportMissingImports]


def _names(path: str, index: int) -> list[str]:
    return [sh.name for sh in Presentation(path).slides[index].shapes]


def _texts(path: str, index: int) -> list[str]:
    return [
        sh.text_frame.text  # type: ignore[reportAttributeAccessIssue]
        for sh in Presentation(path).slides[index].shapes
        if sh.has_text_frame
    ]


@pytest.fixture()
def deck(tmp_path: Path) -> str:
    """A title-and-content slide, the shape every real deck is built from."""
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Where the money went"  # type: ignore[reportOptionalMemberAccess]
    slide.placeholders[1].text = "Google Ads took 1.9M of 2.5M total spend"  # type: ignore[reportAttributeAccessIssue]
    prs.save(str(path))
    return str(path)


class TestTheCopyMatchesTheOriginal:
    def test_it_has_the_same_number_of_shapes(self, deck: str):
        duplicate_slide(deck, 0)
        assert len(_names(deck, 1)) == len(_names(deck, 0))

    def test_it_has_the_same_shapes_in_the_same_order(self, deck: str):
        duplicate_slide(deck, 0)
        assert _names(deck, 1) == _names(deck, 0)

    def test_no_shape_name_appears_twice(self, deck: str):
        """Two shapes called "Title 1" leave an edit with no way to choose."""
        duplicate_slide(deck, 0)
        names = _names(deck, 1)
        assert len(names) == len(set(names)), names

    def test_no_empty_placeholder_is_left_behind(self, deck: str):
        duplicate_slide(deck, 0)
        assert all(t.strip() for t in _texts(deck, 1)), _texts(deck, 1)

    def test_the_text_comes_across(self, deck: str):
        duplicate_slide(deck, 0)
        assert _texts(deck, 1) == _texts(deck, 0)

    def test_the_title_is_still_the_title(self, deck: str):
        duplicate_slide(deck, 0)
        title = Presentation(deck).slides[1].shapes.title
        assert title is not None
        assert title.text == "Where the money went"


class TestItStaysCleanUnderRepetition:
    def test_duplicating_a_duplicate_does_not_compound(self, deck: str):
        duplicate_slide(deck, 0)
        duplicate_slide(deck, 1)
        assert _names(deck, 2) == _names(deck, 0)

    def test_three_copies_all_match(self, deck: str):
        for _ in range(3):
            duplicate_slide(deck, 0)
        prs = Presentation(deck)
        assert len(prs.slides) == 4
        for i in range(1, 4):
            assert _names(deck, i) == _names(deck, 0), i


class TestNonPlaceholderShapesSurvive:
    def test_a_text_box_is_copied_too(self, tmp_path: Path):
        path = tmp_path / "box.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = "205 duplicate rows"
        prs.save(str(path))

        duplicate_slide(str(path), 0)
        assert "205 duplicate rows" in _texts(str(path), 1)

    def test_a_blank_layout_copy_gains_nothing_extra(self, tmp_path: Path):
        path = tmp_path / "blank.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(path))

        duplicate_slide(str(path), 0)
        assert _names(str(path), 1) == _names(str(path), 0)


class TestThePositionArgumentStillWorks:
    def test_appending_puts_the_copy_last(self, deck: str):
        r = duplicate_slide(deck, 0)
        assert r["success"] is True, r.get("error")
        assert len(Presentation(deck).slides) == 2

    def test_insert_at_places_it_where_asked(self, deck: str):
        prs = Presentation(deck)
        second = prs.slides.add_slide(prs.slide_layouts[1])
        second.shapes.title.text = "Data quality"  # type: ignore[reportOptionalMemberAccess]
        prs.save(deck)

        duplicate_slide(deck, 1, insert_at=0)
        titles = [s.shapes.title.text if s.shapes.title else "" for s in Presentation(deck).slides]
        assert titles[0] == "Data quality", titles
