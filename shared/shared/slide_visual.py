"""Geometry and contrast checks for slides — the two ways a deck goes wrong silently.

Both defects here were found by rendering a generated deck to PDF and looking at
it. Every tool involved returned success:true and the file was structurally
valid python-pptx output.

1. Placement is unchecked. `add_chart` takes left/top/width/height in inches and
   hands them straight to python-pptx, which will happily place a shape past the
   edge of the slide. The chart was drawn with its lower half beyond the canvas:
   two of the five category labels were cut off, and nothing in the result said
   so.

2. Colour is unchecked. `set_background` sets one slide at a time, while
   `set_font_all_slides` paints every slide. Set a dark background on slide 1 and
   white text everywhere, and slide 2 gets white text on a white background --
   the title is still there, still "successfully" applied, and completely
   invisible.

Neither of these can be caught by asserting on the return value, so the checks
belong next to the operations that cause them.
"""

from __future__ import annotations

# WCAG 2.1 contrast ratios. 4.5 is the AA threshold for body text and 3.0 for
# large text; a deck title is large, so 3.0 is the line that matters here.
# Anything at or below ~1.3 is effectively the same colour.
MIN_CONTRAST = 3.0
INVISIBLE_CONTRAST = 1.3


def _channel(value: float) -> float:
    """Linearise one sRGB channel for the luminance formula."""
    value = value / 255.0
    return value / 12.92 if value <= 0.03928 else ((value + 0.055) / 1.055) ** 2.4


def relative_luminance(hex_color: str) -> float:
    """WCAG relative luminance of a 6-digit hex colour."""
    clean = hex_color.lstrip("#")
    r, g, b = (int(clean[i : i + 2], 16) for i in (0, 2, 4))
    return 0.2126 * _channel(r) + 0.7152 * _channel(g) + 0.0722 * _channel(b)


def contrast_ratio(fg_hex: str, bg_hex: str) -> float:
    """Return the WCAG contrast ratio between two hex colours (1.0 to 21.0)."""
    a, b = relative_luminance(fg_hex), relative_luminance(bg_hex)
    lighter, darker = max(a, b), min(a, b)
    return (lighter + 0.05) / (darker + 0.05)


def slide_background_hex(slide) -> str | None:
    """Return the slide's own solid background colour, or None if it inherits.

    The slide's own fill wins; failing that the layout's, then the master's.
    Returns None only when nothing in that chain sets a solid fill.

    Resolving the chain matters more than it looks. The deck that prompted this
    had its unreadable slide left at the template default -- white by
    inheritance, not by an explicit fill -- so a check that only looked at the
    slide itself saw nothing and missed the very bug it was written for.
    """
    # Imported here rather than at module scope: shared/ is also imported by the
    # docx and xlsx servers, which have no reason to pull in python-pptx.
    from pptx.enum.dml import MSO_FILL

    for source in (slide, getattr(slide, "slide_layout", None), getattr(slide, "slide_master", None)):
        if source is None:
            continue
        try:
            fill = source.background.fill
            # Anything other than SOLID (usually BACKGROUND) means "inherit",
            # so fall through to the next level of the chain.
            if fill.type == MSO_FILL.SOLID:
                return str(fill.fore_color.rgb)
        except AttributeError, TypeError, ValueError:
            continue
    return None


# PowerPoint renders a slide with nothing set anywhere in its inheritance chain
# on white. Assuming that is a guess, but it is the correct guess for every
# default-template deck, and the alternative -- staying silent -- is what let
# white-on-white text ship in the first place. Offenders found this way are
# flagged `assumed` so the caller can tell the two apart.
_DEFAULT_BACKGROUND = "FFFFFF"


def unreadable_slides(prs, color_hex: str) -> list[dict]:
    """Return the slides where `color_hex` text would be hard or impossible to read."""
    clean = color_hex.lstrip("#").upper()
    offenders: list[dict] = []
    for index, slide in enumerate(prs.slides):
        resolved = slide_background_hex(slide)
        background = resolved if resolved is not None else _DEFAULT_BACKGROUND
        ratio = contrast_ratio(clean, background)
        if ratio < MIN_CONTRAST:
            offenders.append(
                {
                    "slide": index,
                    "background": f"#{background}",
                    "contrast": round(ratio, 2),
                    "invisible": ratio <= INVISIBLE_CONTRAST,
                    "assumed": resolved is None,
                }
            )
    return offenders


def contrast_warning(offenders: list[dict], color_hex: str) -> str:
    """One sentence naming the slides and what to do about them."""
    invisible = [o for o in offenders if o["invisible"]]
    slides = ", ".join(str(o["slide"]) for o in offenders)
    colour = color_hex.lstrip("#").upper()
    # Say so when the background was inherited rather than read off the slide,
    # so a caller using a themed master knows why they are being warned.
    basis = (
        " (those slides inherit the template background, taken as white)"
        if all(o.get("assumed") for o in offenders)
        else ""
    )
    if invisible:
        return (
            f"Text set to #{colour} is effectively invisible on slide(s) {slides}{basis} -- "
            f"it nearly matches the background there. Use set_background on those slides, "
            f"or pick a colour that contrasts with them."
        )
    return (
        f"Text set to #{colour} falls below the {MIN_CONTRAST}:1 readable contrast ratio on "
        f"slide(s) {slides}{basis}. Use set_background on those slides, "
        f"or pick a darker or lighter colour."
    )


def fit_to_slide(
    prs,
    left: float,
    top: float,
    width: float,
    height: float,
    margin: float = 0.2,
) -> tuple[float, float, float, float, str]:
    """Clamp a shape's inch-based box so it stays on the canvas.

    Returns (left, top, width, height, note). `note` is empty when the box
    already fitted; otherwise it describes what was moved or shrunk, so the
    caller can surface it instead of silently drawing off the edge.
    """
    slide_w = prs.slide_width / 914400  # EMU per inch
    slide_h = prs.slide_height / 914400

    original = (left, top, width, height)

    # A box larger than the canvas is shrunk first, then pulled back inside it.
    width = max(0.5, min(width, slide_w - 2 * margin))
    height = max(0.5, min(height, slide_h - 2 * margin))
    left = max(margin, min(left, slide_w - width - margin))
    top = max(margin, min(top, slide_h - height - margin))

    if (round(left, 3), round(top, 3), round(width, 3), round(height, 3)) == tuple(round(v, 3) for v in original):
        return left, top, width, height, ""

    note = (
        f"Requested box {original[2]:.2f}x{original[3]:.2f}in at "
        f"({original[0]:.2f}, {original[1]:.2f}) extended past the "
        f"{slide_w:.2f}x{slide_h:.2f}in slide; fitted to {width:.2f}x{height:.2f}in "
        f"at ({left:.2f}, {top:.2f})."
    )
    return left, top, width, height, note


# PowerPoint renders level-1 body text at 28pt in the stock layouts, not the
# 18pt first assumed here. Getting that wrong made a three-bullet placeholder
# measure 2.69in when it drew 4.6in, and the table landed on top of two of the
# bullets -- visible in the render, with the geometry insisting it was clear.
_DEFAULT_BODY_PT = 28.0
# Advance width of a proportional face as a fraction of point size, and the
# usual line-height multiplier, used to estimate how many lines a paragraph
# wraps to. 0.72 is calibrated against a rendered slide: a 37-character bullet
# in a 9in placeholder at 28pt wraps to two lines, i.e. ~33 characters per
# line. A narrower guess counted it as one and put a table over the next
# bullet.
_CHAR_WIDTH_RATIO = 0.72
_LINE_SPACING = 1.25
# Deliberately overestimate. The two failure modes are not symmetric: too small
# and shapes overlap and the slide is unreadable, too large and there is a bit
# of extra whitespace. Bias towards the whitespace.
_HEIGHT_SAFETY = 1.1


def _text_height(shape, width_in: float) -> float:
    """Estimate the height a text frame actually draws, in inches.

    A placeholder's frame is its *maximum* extent, not its content: the default
    layout gives the body 4.95in of a 7.5in slide whether it holds one line or
    twenty. Measuring the frame made a single-line bullet look like it filled
    the slide, so nothing could ever be placed below it.
    """
    frame = shape.text_frame
    lines = 0
    for para in frame.paragraphs:
        text = "".join(run.text for run in para.runs) or ""
        size_pt = _DEFAULT_BODY_PT
        for run in para.runs:
            if run.font.size is not None:
                size_pt = run.font.size.pt
                break
        chars_per_line = max(1, int(width_in * 72 / (size_pt * _CHAR_WIDTH_RATIO)))
        lines += max(1, -(-len(text) // chars_per_line))  # ceil division
    drawn = lines * _DEFAULT_BODY_PT * _LINE_SPACING * _HEIGHT_SAFETY / 72.0
    # Never claim more than the frame itself allows.
    return min(drawn, (shape.height or 0) / 914400)


# Below this, a chart or table is a sliver of unreadable pixels rather than
# content, so the call is refused instead of drawn.
MIN_USABLE_HEIGHT = 1.0


class NoRoomOnSlide(Exception):
    """Raised when a shape cannot be placed without covering existing content."""


def _occupied_bottom(slide, left: float, width: float, margin: float) -> float:
    """Return the lowest bottom edge, in inches, of content already on the slide.

    Only shapes that overlap the new box horizontally count -- something in a
    side column should not push a chart down the page. Text shapes are measured
    by what they draw rather than by their frame, and empty placeholders are
    skipped entirely: an unfilled layout placeholder renders as nothing, and
    treating it as content would shove every shape off the bottom of the slide.
    """
    bottom = 0.0
    for shape in slide.shapes:
        try:
            s_left = (shape.left or 0) / 914400
            s_top = (shape.top or 0) / 914400
            s_width = (shape.width or 0) / 914400
            s_height = (shape.height or 0) / 914400
        except AttributeError, TypeError:
            continue
        if s_width <= 0 or s_height <= 0:
            continue

        if getattr(shape, "has_text_frame", False):
            if not shape.text_frame.text.strip():
                continue  # empty placeholder or textbox draws nothing
            s_height = _text_height(shape, s_width)

        # No horizontal overlap means no collision.
        if s_left + s_width <= left + margin or s_left >= left + width - margin:
            continue
        bottom = max(bottom, s_top + s_height)
    return bottom


def place_below_content(
    prs,
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    margin: float = 0.2,
    gap: float = 0.15,
) -> tuple[float, float, float, float, str]:
    """Fit a box to the slide *and* keep it clear of what is already there.

    fit_to_slide only knows about the slide edges, so a chart and a table added
    to a slide that already had bullet text were all placed at their default
    top and drawn on top of each other -- three overlapping shapes, none of
    them readable, and every call reported success. This drops the new box
    below the existing content when it would collide, shrinking it if the
    remaining space is tight, and says so.
    """
    left, top, width, height, note = fit_to_slide(prs, left, top, width, height, margin)

    occupied = _occupied_bottom(slide, left, width, margin)
    if occupied <= top:
        return left, top, width, height, note

    slide_h = prs.slide_height / 914400
    new_top = occupied + gap
    available = slide_h - margin - new_top

    if available < MIN_USABLE_HEIGHT:
        # Returning the caller's position here was the original mistake: it put
        # the shape straight back on top of the content, which is the exact
        # unreadable slide this function exists to prevent. There is no
        # placement that satisfies the request, so say so rather than produce a
        # file that looks fine to a return-value check and is illegible to read.
        raise NoRoomOnSlide(
            f"Slide already has content down to {occupied:.2f}in, leaving {available:.2f}in "
            f"below it -- not enough for a readable {height:.2f}in shape."
        )

    new_height = min(height, available)
    collision = (
        f"Existing content reaches {occupied:.2f}in, so the box was moved from "
        f"top {top:.2f}in to {new_top:.2f}in"
        + (f" and shrunk to {new_height:.2f}in high" if new_height < height else "")
        + " to avoid overlapping it."
    )
    return left, new_top, width, new_height, f"{note} {collision}".strip()
