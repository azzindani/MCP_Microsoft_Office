"""Document diff engine — paragraph, cell, and shape-level comparison."""

import difflib
from typing import Any

from shared.counts import counted

from .file_utils import scrub_repr

MAX_CHANGED_CELLS = 500


def _docx_table_rows(doc: Any) -> list[tuple[str, ...]]:
    """Every table's cells, one flat row of text per table row.

    python-docx's `doc.paragraphs` walks the body only; a paragraph inside a
    table cell is not in it. Diffing paragraphs alone therefore cannot see a
    table at all -- see diff_docx.

    Tuples, not lists: SequenceMatcher hashes the elements it compares.
    """
    rows: list[tuple[str, ...]] = []
    for t_index, table in enumerate(doc.tables):
        for r_index, row in enumerate(table.rows):
            rows.append((f"t{t_index}r{r_index}", *[cell.text for cell in row.cells]))
    return rows


def diff_docx(path_a: str, path_b: str) -> dict[str, Any]:
    """
    Compare two .docx files at paragraph and table-cell level.

    Returns structured diff with added, removed, changed paragraphs and rows.

    Paragraphs alone were not enough. `doc.paragraphs` excludes anything inside
    a table cell, so a document that gained a whole 2x3 table came back
    `change_count: 0`, "No changes detected." -- the two files differed by six
    cells of text and one table. add_table, set_cell, add_row and delete_row
    were all invisible to the one tool whose job is confirming what changed.
    """
    try:
        from docx import Document  # type: ignore[import-untyped]

        doc_a = Document(path_a)
        doc_b = Document(path_b)

        paras_a = [{"index": i, "text": p.text, "style": p.style.name} for i, p in enumerate(doc_a.paragraphs)]  # type: ignore[reportOptionalMemberAccess]
        paras_b = [{"index": i, "text": p.text, "style": p.style.name} for i, p in enumerate(doc_b.paragraphs)]  # type: ignore[reportOptionalMemberAccess]

        texts_a = [p["text"] for p in paras_a]
        texts_b = [p["text"] for p in paras_b]

        matcher = difflib.SequenceMatcher(None, texts_a, texts_b, autojunk=False)

        changes: list[dict[str, Any]] = []
        for tag, i1, i2, j1, j2 in matcher.get_opcodes():
            if tag == "equal":
                continue
            changes.append(
                {
                    "type": tag,
                    "a_range": [i1, i2],
                    "b_range": [j1, j2],
                    "a_text": texts_a[i1:i2],
                    "b_text": texts_b[j1:j2],
                }
            )

        rows_a = _docx_table_rows(doc_a)
        rows_b = _docx_table_rows(doc_b)
        row_matcher = difflib.SequenceMatcher(None, rows_a, rows_b, autojunk=False)
        table_changes: list[dict[str, Any]] = []
        for tag, i1, i2, j1, j2 in row_matcher.get_opcodes():
            if tag == "equal":
                continue
            table_changes.append(
                {
                    "type": tag,
                    "a_range": [i1, i2],
                    "b_range": [j1, j2],
                    # The row label is carried in the leading element; the cells
                    # are what the reader wants to see.
                    "a_cells": [list(r[1:]) for r in rows_a[i1:i2]],
                    "b_cells": [list(r[1:]) for r in rows_b[j1:j2]],
                }
            )

        summary = _summarise_docx_diff(changes, len(paras_a), len(paras_b), table_changes)

        return {
            "success": True,
            "file_a": path_a,
            "file_b": path_b,
            "paragraph_count_a": len(paras_a),
            "paragraph_count_b": len(paras_b),
            "table_count_a": len(doc_a.tables),
            "table_count_b": len(doc_b.tables),
            "changes": changes,
            "table_changes": table_changes,
            "change_count": len(changes) + len(table_changes),
            "summary": summary,
        }
    except Exception as e:
        return {
            "success": False,
            "error": scrub_repr(e),
            "hint": "Check that both paths point to valid .docx files.",
        }


def diff_xlsx(path_a: str, path_b: str, sheet_name: str | None = None) -> dict[str, Any]:
    """
    Compare two .xlsx files at cell level.

    Returns structured diff with changed cells by sheet.
    If sheet_name provided, only diff that sheet.
    """
    try:
        import openpyxl  # type: ignore[import-untyped]

        wb_a = openpyxl.load_workbook(path_a, data_only=True)
        wb_b = openpyxl.load_workbook(path_b, data_only=True)

        sheets_a = set(wb_a.sheetnames)
        sheets_b = set(wb_b.sheetnames)

        result: dict[str, Any] = {
            "success": True,
            "added_sheets": list(sheets_b - sheets_a),
            "removed_sheets": list(sheets_a - sheets_b),
            "sheet_diffs": {},
        }

        common = sheets_a & sheets_b
        if sheet_name:
            common = {sheet_name} if sheet_name in common else set()

        total_changes = 0
        truncated = False

        for name in sorted(common):
            ws_a = wb_a[name]
            ws_b = wb_b[name]
            changed_cells: list[dict[str, Any]] = []

            # Collect all cell coordinates from both sheets
            all_coords: set[str] = set()
            for row in ws_a.iter_rows():
                for cell in row:
                    all_coords.add(cell.coordinate)
            for row in ws_b.iter_rows():
                for cell in row:
                    all_coords.add(cell.coordinate)

            for coord in sorted(all_coords):
                val_a = ws_a[coord].value if coord in ws_a else None  # type: ignore[reportOperatorIssue]
                val_b = ws_b[coord].value if coord in ws_b else None  # type: ignore[reportOperatorIssue]
                if val_a != val_b:
                    changed_cells.append({"cell": coord, "old": val_a, "new": val_b})
                    total_changes += 1
                    if total_changes >= MAX_CHANGED_CELLS:
                        truncated = True
                        break

            if changed_cells:
                result["sheet_diffs"][name] = {
                    "changed_cells": changed_cells,
                    "change_count": len(changed_cells),
                }
            if truncated:
                break

        # Present on every diff, not only the cut ones: a caller could not tell
        # a clean comparison from one this stopped short of finishing. The walk
        # halts at MAX_CHANGED_CELLS, so when it fills, `total` is a floor and
        # `counted()` marks it as one rather than implying an exact count.
        result.update(counted(total_changes, total_changes, exact=not truncated))
        if truncated:
            result["total_changes_approx"] = f">{MAX_CHANGED_CELLS}"

        result["summary"] = _summarise_xlsx_diff(result)
        return result

    except Exception as e:
        return {
            "success": False,
            "error": scrub_repr(e),
            "hint": "Check that both paths point to valid .xlsx files.",
        }


def _pptx_shape_texts(slide: Any) -> dict[str, str]:
    """Readable text per shape, including tables.

    A table shape has no `text_frame`, so a `has_text_frame` filter drops it
    entirely and a slide that gained a whole table diffed as unchanged. Its rows
    are rendered as text so the existing per-shape comparison covers them.
    """
    texts: dict[str, str] = {}
    for shape in slide.shapes:
        key = shape.name
        # Shape names are not unique within a slide; without this the second
        # shape of a name silently replaces the first.
        if key in texts:
            key = f"{shape.name}#{shape.shape_id}"
        if getattr(shape, "has_text_frame", False):
            texts[key] = shape.text_frame.text
        elif getattr(shape, "has_table", False):
            texts[key] = "\n".join(" | ".join(cell.text for cell in row.cells) for row in shape.table.rows)
    return texts


def diff_pptx(path_a: str, path_b: str) -> dict[str, Any]:
    """
    Compare two .pptx files at shape-text level, tables included.

    Returns structured diff with changed text per slide per shape.

    `change_count` counts slides added and removed as well as shape text edits.
    It used to be `len(text_changes)` alone, so deleting a slide came back
    `change_count: 0`, progress "0 changes", in the same response whose own
    summary read "Slide count changed: 2 -> 1 (removed 1)" and whose
    `slide_count_changed` was true. A caller checking the count first -- which
    is what the count is for -- concluded nothing had happened to a deck that
    had just lost a slide. diff_docx beside it has always summed every kind of
    change it found (`len(changes) + len(table_changes)`); this is that rule,
    applied to the sibling that was missing it.

    The zip below also walks only the slides the two decks have in common, so
    the shapes on a dropped slide were invisible twice over. They are reported
    as `slide_changes` entries carrying the text that went with them.
    """
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        prs_a = Presentation(path_a)
        prs_b = Presentation(path_b)

        count_a = len(prs_a.slides)
        count_b = len(prs_b.slides)

        changes: list[dict[str, Any]] = []
        for i, (slide_a, slide_b) in enumerate(zip(prs_a.slides, prs_b.slides)):
            shapes_a = _pptx_shape_texts(slide_a)
            shapes_b = _pptx_shape_texts(slide_b)

            all_names = set(shapes_a) | set(shapes_b)
            for name in sorted(all_names):
                t_a = shapes_a.get(name)
                t_b = shapes_b.get(name)
                if t_a != t_b:
                    changes.append(
                        {
                            "slide_index": i,
                            "shape_name": name,
                            "old_text": t_a,
                            "new_text": t_b,
                        }
                    )

        # Slides past the end of the shorter deck: never paired above, so their
        # text has to be reported here or it is lost from the diff entirely.
        slide_changes: list[dict[str, Any]] = []
        common = min(count_a, count_b)
        if count_b > count_a:
            for i, slide in enumerate(list(prs_b.slides)[common:], start=common):
                slide_changes.append({"slide_index": i, "change": "added", "shape_texts": _pptx_shape_texts(slide)})
        elif count_a > count_b:
            for i, slide in enumerate(list(prs_a.slides)[common:], start=common):
                slide_changes.append({"slide_index": i, "change": "removed", "shape_texts": _pptx_shape_texts(slide)})

        summary = _summarise_pptx_diff(changes, count_a, count_b)

        return {
            "success": True,
            "slide_count_a": count_a,
            "slide_count_b": count_b,
            "slide_count_changed": count_a != count_b,
            "text_changes": changes,
            "slide_changes": slide_changes,
            "change_count": len(changes) + len(slide_changes),
            "summary": summary,
        }
    except Exception as e:
        return {
            "success": False,
            "error": scrub_repr(e),
            "hint": "Check that both paths point to valid .pptx files.",
        }


def format_diff_as_text(diff: dict[str, Any]) -> str:
    """Format a diff dict as a human-readable unified-diff-style string."""
    lines: list[str] = []

    if not diff.get("success"):
        return f"Diff failed: {diff.get('error', 'unknown error')}"

    summary = diff.get("summary", "")
    if summary:
        lines.append(summary)
        lines.append("")

    # DOCX diff
    if "paragraph_count_a" in diff:
        for change in diff.get("changes", []):
            for text in change.get("a_text", []):
                lines.append(f"- {text[:120]}")
            for text in change.get("b_text", []):
                lines.append(f"+ {text[:120]}")
        # Table rows too, or this renders an empty diff for a document whose
        # only change was a table -- the same blind spot diff_docx had.
        for change in diff.get("table_changes", []):
            for cells in change.get("a_cells", []):
                lines.append(f"- | {' | '.join(cells)}"[:120])
            for cells in change.get("b_cells", []):
                lines.append(f"+ | {' | '.join(cells)}"[:120])

    # XLSX diff
    if "sheet_diffs" in diff:
        for sheet, sdata in diff["sheet_diffs"].items():
            lines.append(f"Sheet: {sheet}")
            for c in sdata.get("changed_cells", []):
                lines.append(f"  {c['cell']}: {c['old']!r} → {c['new']!r}")

    # PPTX diff
    if "text_changes" in diff:
        for change in diff["text_changes"]:
            lines.append(f"Slide {change['slide_index']} / {change['shape_name']}:")
            old = (change["old_text"] or "")[:80]
            new = (change["new_text"] or "")[:80]
            lines.append(f"  - {old}")
            lines.append(f"  + {new}")

    return "\n".join(lines)


def _summarise_docx_diff(
    changes: list[dict[str, Any]],
    count_a: int,
    count_b: int,
    table_changes: list[dict[str, Any]] | None = None,
) -> str:
    table_changes = table_changes or []
    if not changes and not table_changes:
        return "No changes detected."

    # Count the rows and paragraphs, not the opcodes. difflib returns one
    # opcode per contiguous run, so counting opcodes reported a 2-row table
    # insert as "1 table row added" and a replace spanning four paragraphs as
    # "1 paragraph changed" -- a summary that disagreed with the very list of
    # changes printed beneath it.
    def _spans(source: list[dict[str, Any]]) -> tuple[int, int, int]:
        changed = added = deleted = 0
        for c in source:
            a_len = c["a_range"][1] - c["a_range"][0]
            b_len = c["b_range"][1] - c["b_range"][0]
            if c["type"] == "replace":
                changed += max(a_len, b_len)
            elif c["type"] == "insert":
                added += b_len
            elif c["type"] == "delete":
                deleted += a_len
        return changed, added, deleted

    parts = []
    for label, source in (("paragraph", changes), ("table row", table_changes)):
        n_changed, n_added, n_deleted = _spans(source)
        if n_changed:
            parts.append(f"{n_changed} {label}{'s' if n_changed != 1 else ''} changed")
        if n_added:
            parts.append(f"{n_added} {label}{'s' if n_added != 1 else ''} added")
        if n_deleted:
            parts.append(f"{n_deleted} {label}{'s' if n_deleted != 1 else ''} deleted")
    return ". ".join(parts) + "."


def _summarise_xlsx_diff(result: dict[str, Any]) -> str:
    sheet_diffs = result.get("sheet_diffs", {})
    added = result.get("added_sheets", [])
    removed = result.get("removed_sheets", [])
    parts = []
    total = sum(v["change_count"] for v in sheet_diffs.values())
    if total:
        parts.append(f"{total} cell{'s' if total != 1 else ''} changed")
    if added:
        parts.append(f"{len(added)} sheet{'s' if len(added) != 1 else ''} added")
    if removed:
        parts.append(f"{len(removed)} sheet{'s' if len(removed) != 1 else ''} removed")
    if result.get("truncated"):
        parts.append("(truncated at 500 changes)")
    return ". ".join(parts) + "." if parts else "No changes detected."


def _summarise_pptx_diff(changes: list[dict[str, Any]], count_a: int, count_b: int) -> str:
    if not changes and count_a == count_b:
        return "No changes detected."
    parts = []
    if changes:
        parts.append(f"{len(changes)} shape text{'s' if len(changes) != 1 else ''} changed")
    if count_a != count_b:
        diff = count_b - count_a
        parts.append(f"Slide count changed: {count_a} → {count_b} ({'added' if diff > 0 else 'removed'} {abs(diff)})")
    return ". ".join(parts) + "."
